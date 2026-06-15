"""ClonePlan: deep-copy a bespoke slide from the source deck and patch text.

For slides too rich to encode as layouts (timeline ribbons, funnel infographics,
SmartArt) the source designer-authored geometry is the answer. We clone the
shape tree byte-for-byte and rewrite text via XML `<a:t>` runs — the shape
iterator misses grouped / SmartArt text, the XML scan doesn't.
"""

from __future__ import annotations

import copy
from dataclasses import dataclass, field
from io import BytesIO

from lxml import etree
from pptx import Presentation

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS = {"a": _A}


@dataclass
class ClonePlan:
    snippet_id: str
    text_replacements: list[tuple[str, str]] = field(default_factory=list)


def apply_clone_plan(prs, plan: ClonePlan, catalog) -> None:
    snippet = catalog.snippets.get(plan.snippet_id)
    if snippet is None:
        raise KeyError(f"snippet not in catalog: {plan.snippet_id!r}")
    source_deck = catalog.source_deck
    src_prs = Presentation(str(source_deck))
    src_slide = src_prs.slides[snippet.source_idx]
    layout_name = src_slide.slide_layout.name
    dest_layout = _layout_by_name(prs, layout_name)
    new_slide = prs.slides.add_slide(dest_layout)

    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Map src image rels → new rIds in the master before cloning XML, so
    # r:embed / r:link attrs on cloned shapes can be rewritten. When the
    # source deck IS the master, naive relate_to(target_part) collides on
    # the master's /ppt/media/ partname; re-add via get_or_add_image_part.
    rid_map: dict[str, str] = {}
    for src_rid, rel in src_slide.part.rels.items():
        if rel.reltype.endswith("/image"):
            _, new_rid = new_slide.part.get_or_add_image_part(BytesIO(rel.target_part.blob))
            rid_map[src_rid] = new_rid

    for child in list(src_slide.shapes._spTree):
        if etree.QName(child).localname in ("nvGrpSpPr", "grpSpPr"):
            continue
        cloned = copy.deepcopy(child)
        if rid_map:
            for el in cloned.iter():
                embed = el.get(f"{{{_R}}}embed")
                if embed and embed in rid_map:
                    el.set(f"{{{_R}}}embed", rid_map[embed])
                link = el.get(f"{{{_R}}}link")
                if link and link in rid_map:
                    el.set(f"{{{_R}}}link", rid_map[link])
        new_slide.shapes._spTree.append(cloned)

    if plan.text_replacements:
        _patch_text(new_slide, plan.text_replacements)


def _layout_by_name(prs, name: str):
    target = name.strip()
    for layout in prs.slide_layouts:
        if layout.name.strip() == target:
            return layout
    raise KeyError(f"layout not in master: {name!r}")


def _patch_text(slide, replacements: list[tuple[str, str]]) -> None:
    queues: dict[str, list[str]] = {}
    for old, new in replacements:
        queues.setdefault(old, []).append(new)
    runs = slide.shapes._spTree.findall(".//a:t", _NS)
    for run in runs:
        if run.text and run.text in queues and queues[run.text]:
            run.text = queues[run.text].pop(0)
