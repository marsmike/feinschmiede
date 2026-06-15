"""Chart insertion into an OBJECT placeholder.

The chart inherits the master theme — Bosch corporate blue, BSH orange — at
zero styling cost. The placeholder's bbox is read, the placeholder is removed,
and `add_chart` is called with the same coordinates.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Literal

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

ChartKind = Literal["column", "bar", "line", "pie"]

_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


@dataclass(frozen=True)
class ChartSpec:
    kind: ChartKind
    categories: list[str]
    series: list[tuple[str, list[float]]] = field(default_factory=list)

    def to_chart_data(self) -> CategoryChartData:
        data = CategoryChartData()
        data.categories = self.categories
        for name, values in self.series:
            data.add_series(name, tuple(values))
        return data

    @property
    def xl_type(self):
        return _TYPE_MAP[self.kind]


def add_chart_into_placeholder(slide, ph_idx: int, spec: ChartSpec) -> None:
    obj = next(p for p in slide.placeholders if p.placeholder_format.idx == ph_idx)
    left, top, width, height = obj.left, obj.top, obj.width, obj.height
    obj._element.getparent().remove(obj._element)
    slide.shapes.add_chart(spec.xl_type, left, top, width, height, spec.to_chart_data())
