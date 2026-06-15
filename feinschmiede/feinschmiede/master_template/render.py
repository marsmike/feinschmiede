"""Orchestrator: open the master, strip its example slides, apply each plan.

Public entry point:
    render(brand_pack: Path, plans: Iterable[FillPlan | ClonePlan], out: Path) -> Path
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation

from feinschmiede.master_template.catalog import Catalog, load_catalog
from feinschmiede.master_template.clone_plan import ClonePlan, apply_clone_plan
from feinschmiede.master_template.fill_plan import FillPlan, apply_fill_plan
from feinschmiede.master_template.themes import (
    build_swap,
    discover_themes,
    load_theme_colors,
    recolor_element,
)

_PPTX_RELS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def render(
    brand_pack: Path,
    plans: Iterable[FillPlan | ClonePlan],
    out: Path,
    *,
    catalog: Catalog | None = None,
    theme: str | None = None,
) -> Path:
    cat = catalog or load_catalog(Path(brand_pack))
    prs = Presentation(str(cat.master_pptx))
    strip_existing_slides(prs)

    swap = _build_theme_swap(cat, theme) if theme else {}

    for plan in plans:
        before = len(prs.slides)
        if isinstance(plan, FillPlan):
            apply_fill_plan(prs, plan, cat)
        elif isinstance(plan, ClonePlan):
            apply_clone_plan(prs, plan, cat)
        else:
            raise TypeError(f"unknown plan type: {type(plan).__name__}")
        if swap:
            for slide in list(prs.slides)[before:]:
                recolor_element(slide.shapes._spTree, swap)

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def _build_theme_swap(cat: Catalog, theme: str) -> dict[str, str]:
    themes = discover_themes(cat.brand_pack)
    if theme not in themes:
        raise KeyError(
            f"theme {theme!r} not found in {cat.brand_pack}/themes/; "
            f"available: {sorted(themes)}"
        )
    source_theme = cat.master_theme or "default"
    if source_theme not in themes:
        # No source-theme tokens present → cannot derive a swap. Render as-is.
        return {}
    if source_theme == theme:
        return {}
    return build_swap(load_theme_colors(themes[source_theme]),
                      load_theme_colors(themes[theme]))


def strip_existing_slides(prs) -> None:
    sld_id_lst = prs.slides._sldIdLst
    rels = prs.part.rels
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(_PPTX_RELS)
        sld_id_lst.remove(sld_id)
        if r_id in rels:
            rels.pop(r_id)
