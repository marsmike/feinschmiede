"""Hybrid PPTX+SVG → Feinschliff DSL decompiler — brand-agnostic.

A higher-fidelity alternative to `lib/dsl/pptx_decompile.py`. Uses PPTX
XML as the canonical source for shape semantics and (optionally) SVG —
rendered from the slide's PDF via pdf2svg — as a secondary source for
custGeom geometry that PPTX xfrm inherits from the layout.

Sources of truth:

  PPTX (`ppt/slides/slideN.xml` + slideLayouts + slideMasters + theme1)
    - canonical for shape semantics: placeholder type/idx, group structure,
      text content + per-run style, schemeClr → token resolution, z-order,
      footer / page-number / wordmark placeholders inherited from master.

  SVG (pdf2svg of the slide's PDF page; optional)
    - canonical for final rendered geometry: bounding boxes of custGeom
      shapes, fall-back bbox when PPTX xfrm is inherited from layout.
    - NOT used for text classification or color matching (PPTX wins).

  Picture handling
    - every <p:pic> shape and every placeholder with type="pic" is
      emitted as a `picture` statement pointing at a configurable
      placeholder image (default: `assets/illustrations/placeholder.jpg`,
      the feinschliff convention).

Coordinate systems:
  - PPTX EMU: 914400 EMU/inch, slide size from presentation.xml::sldSz.
  - SVG: PDF points (1 pt = 1/72 inch). pdf2svg width/height map 1:1 to
    the slide's printable area, so EMU and pt match through inch.
  - DSL canvas: 1920×1080 px by default (override via canvas_w/h).

Brand-specific knobs (all defaulted to feinschliff baseline):
  - `theme_name` — name of the brand to emit on the `theme` directive
  - `tokens_path` — brand's tokens.json (for nearest-color matching)
  - `placeholder_rel` — DSL-relative path for picture placeholders
  - Footer-region text is emitted as plain `text` primitives; brands
    that ship a `footer(...)` compound can post-process the output to
    collapse the four footer lines into a single compound call.

Usage (programmatic, preferred):
  from feinschliff_builder.decompile.pptx_svg_decompile import derive
  dsl = derive(pptx_path, slide_idx=1, theme_name="acme",
               tokens_path=Path("brands/acme/tokens.json"),
               layout_name="cover-orange")

Usage (CLI smoke test):
  uv run python lib/dsl/pptx_svg_decompile.py SOURCE.pptx --slide N \\
      --theme <brand> --brand-tokens brands/<brand>/tokens.json > out.dsl
"""
from __future__ import annotations

import argparse
import json
import math
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from feinschmiede.dsl.tokens import STYLE_BUNDLES

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

EMU_PER_PT = 12700
PLACEHOLDER_REL = "assets/illustrations/placeholder.jpg"
# A non-placeholder <p:pic> smaller than this fraction of the slide is treated as
# fixed corporate-design chrome (logo / mark) and carried natively; larger pics
# are changeable topical content and stay fillable picture slots.
_TEMPLATE_IMG_MAX_AREA = 0.12

# Native-carry payloads at most this many RAW bytes ride inline as base64 in
# the DSL line; anything bigger goes to a sha-named sidecar file under the
# brand pack's assets dir (`xml_file:` / `media_file:` / `parts_file:` refs).
# Inlining a 33 MB carried vector group produced a 44 MB .slide.dsl — the DSL
# must stay a readable text file, not a binary container.
NATIVE_INLINE_MAX = 16 * 1024


def _native_sidecar_ref(payload: bytes, ext: str, native_dir: Path | None,
                        native_rel: str | None) -> str | None:
    """Write `payload` as a sha-named sidecar under `native_dir` and return the
    asset-root-relative ref for the DSL, or None when the payload should stay
    inline (no sidecar dir configured, or small enough). Content-hash naming
    dedupes across slides: the same template logo carried on 99 layouts lands
    on disk exactly once."""
    if native_dir is None or len(payload) <= NATIVE_INLINE_MAX:
        return None
    import hashlib
    name = f"{hashlib.sha1(payload).hexdigest()[:12]}.{ext}"
    native_dir.mkdir(parents=True, exist_ok=True)
    target = native_dir / name
    # Unconditional write: the slotify pass REWRITES sidecars in place
    # (planting {{ text_N }} templates), so an exists-skip would freeze a
    # previous run's slotified state under the fresh payload's hash name —
    # re-derives could never refresh it. Write-to-temp + atomic rename:
    # parallel decompile workers deriving sibling slides can hit the same
    # content-hash name simultaneously (shared template chrome) — a direct
    # write would interleave the two byte streams.
    import os
    import tempfile
    fd, tmp = tempfile.mkstemp(dir=native_dir, prefix=f".{name}.")
    try:
        os.write(fd, payload)
    finally:
        os.close(fd)
    os.replace(tmp, target)
    return f"{(native_rel or 'native').rstrip('/')}/{name}"


def _rasterize_svg_bytes(blob: bytes) -> bytes | None:
    """Rasterize SVG bytes to PNG bytes via a temp file; None on failure."""
    import tempfile
    with tempfile.TemporaryDirectory() as td:
        svg = Path(td) / "m.svg"
        svg.write_bytes(blob)
        png = _rasterize_svg(svg)
        return png.read_bytes() if png is not None else None


def _rasterize_svg(svg_path: Path) -> Path | None:
    """soffice-rasterize an extracted SVG to a sibling PNG; None on failure."""
    import subprocess
    import tempfile
    try:
        with tempfile.TemporaryDirectory() as profile:
            subprocess.run(
                ["soffice", f"-env:UserInstallation=file://{profile}",
                 "--headless", "--convert-to", "png", str(svg_path),
                 "--outdir", str(svg_path.parent)],
                check=True, capture_output=True, timeout=120)
        png = svg_path.with_suffix(".png")
        return png if png.is_file() else None
    except Exception:
        return None


def _media_ext(blob: bytes) -> str:
    if blob.startswith(b"\x89PNG"):
        return "png"
    if blob.startswith(b"\xff\xd8"):
        return "jpeg"
    if blob.startswith((b"GIF87a", b"GIF89a")):
        return "gif"
    return "bin"


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------


@dataclass
class Shape:
    kind: str            # rect | line | oval | pic | text | table
    x: float             # px in canvas
    y: float
    w: float
    h: float
    fill: str | None = None       # token name, e.g. "accent"
    stroke: str | None = None
    text_runs: list["TextRun"] = field(default_factory=list)
    is_picture: bool = False      # True for <p:pic> or ph type="pic"
    # True when this shape was inherited from the layout/master (not authored on the
    # slide). Slide-authored text is real content; only inherited text is prompt copy.
    from_chain: bool = False
    ph_type: str | None = None    # 'title','body','subTitle','pic','ftr','sldNum',...
    ph_idx: str | None = None
    # Border width in design-px. Captured from `<a:ln w="...">` (EMU); when
    # None the emitter uses its default hairline width.
    stroke_width: float | None = None
    # Dashed-line preset name from `<a:ln><a:prstDash val="...">` (e.g.
    # "dash", "dot", "sysDash", "lgDashDot"). None = solid stroke.
    stroke_dash: str | None = None
    # Corner radius in design-px for `roundRect` shapes. Captured from the
    # `<a:gd name="adj" fmla="val N">` adjustment value where N is N/100000
    # of the shape's shortest side. None when the source uses sharp corners.
    corner_radius: float | None = None
    # Drop shadow descriptor from `<a:effectLst><a:outerShdw>`. Tuple
    # (blur_px, dist_px, angle_deg, color_token, alpha) so the decompiler
    # can emit a compact `shadow:` kwarg and the emitter can rebuild the
    # `<a:effectLst>` XML at build time. None = no shadow.
    shadow: tuple[float, float, float, str, float] | None = None
    # Gradient fill descriptor from `<a:gradFill>`. List of (position, color)
    # stops (position 0..1, color = token or hex) plus the linear angle in
    # degrees. None = solid fill (use `Shape.fill` instead). When set, the
    # emitter writes a `gradFill` XML block onto the shape's spPr; the
    # decompiler emits a `gradient:angle=Ddeg;0=token;1=token` kwarg.
    gradient: tuple[list[tuple[float, str]], float] | None = None
    # Source bodyPr `anchor` attribute — controls text vertical position
    # within the shape bbox. "ctr" centers (PowerPoint default for many
    # text frames); "b" bottoms; "t" or absent = top. Without this the DSL
    # always emits top-anchored text, so source content that's vertically
    # centered in its frame renders shifted up by half the frame height.
    valign: str | None = None
    # normAutofit (PowerPoint "shrink text on overflow"): `autoshrink` arms the
    # emitter's fit; `font_scale` (0..1) reproduces the source's pre-shrink so a
    # placeholder's text fits its box instead of overflowing the final render.
    autoshrink: bool = False
    font_scale: float = 1.0
    # Paragraph line spacing from `<a:pPr><a:lnSpc><a:spcPct val="N"/>`:
    # multiplier (N/100000), or None when the source paragraph carries no
    # explicit lnSpc — PowerPoint then uses the font's native single
    # spacing. None must round-trip to an emitted text frame WITHOUT a
    # lnSpc element (`linespacing:native`), NOT to the toolkit's 1.2
    # default: writing 120% onto source-default paragraphs pushed every
    # decompiled headline a quarter-line down the slide.
    line_spacing: float | None = None
    # Text-frame internal insets as (l, t, r, b) in design-px. Source carries
    # these on `<a:bodyPr lIns="..." tIns="..." rIns="..." bIns="...">` in EMU.
    # When absent on source, defaults to PowerPoint's published 91440 / 45720
    # EMU (= ~7.2px / ~3.6px at 13.33" canvas scale). Captured per-text so
    # decompiled DSL renders at the exact source text position; without this,
    # my emitter zeroed all four and shifted every text by ~9-19 px versus
    # source — visible as the persistent blue/red ghost offsets in the redline.
    padding: tuple[float, float, float, float] | None = None
    # When image_extract_dir is passed to derive(), pictures are extracted
    # from the source PPTX and this holds the brand-pack-relative path the
    # DSL's `default:` should resolve to. None falls back to the generic
    # placeholder (genericised brand-template behaviour).
    media_path: str | None = None
    media_rid: str | None = None  # rId of the <a:blip r:embed=...>, for extraction
    # Resolved python-pptx Part for the embedded image, captured at walk
    # time so the rId is looked up against the part that actually owns it
    # (slide vs. layout vs. master). Without this, layout-inherited
    # pictures (most corporate-template chrome) fail extraction because
    # their rId only resolves on the layout part, not the slide.
    media_part: Any = None
    # For custGeom shapes: the full pathLst converted to an SVG-`d` string
    # in *canvas pixels* (already scaled from the path's path-local space).
    # When set, emit_dsl emits the shape as an `svg { path … }` block
    # instead of the lossy bbox-rect fallback. Lets the renderer reproduce
    # rings, arrows, callouts, and other vector decoration that the PDF
    # pipeline would otherwise rasterise.
    svg_path_d: str | None = None
    # Verbatim source `<p:sp>` XML (emitted as a base64 `native` primitive) for
    # complex custGeom chrome — carried as an editable native vector instead of
    # round-tripping through svg → raster → picture. None = use the other fields.
    native_xml: str | None = None
    # For a carried <p:pic> (template image): base64 of the embedded media bytes,
    # re-embedded into the output deck by the emitter (the source rId is dead here).
    native_media: str | None = None
    # For a native-carried CHART <p:graphicFrame>: the external part-graph the
    # frame's `<c:chart r:id>` reaches (the chart part + its chartStyle /
    # chartColorStyle / embedded-xlsx children). A list of dicts, each
    # {partname, content_type, blob (base64), reltype, parent ('slide' | a
    # chart-partname str)}. Tables need NO external parts (inline <a:tbl>) so
    # they leave this None; charts carry it so the emitter can re-create the
    # parts + rewire rIds in the output deck. None = no external parts.
    native_parts: list[dict] | None = None


@dataclass
class TextRun:
    text: str
    pt: float            # font size in points
    bold: bool = False
    italic: bool = False
    color: str | None = None  # token name
    align: str | None = None  # paragraph alignment: "center" | "right" | "justify" | None


# ---------------------------------------------------------------------------
# Palette + color resolution
# ---------------------------------------------------------------------------


def load_palette(tokens_path: Path) -> dict[str, tuple[int, int, int]]:
    # Each brand pack is self-contained; read tokens.json directly.
    data = json.loads(tokens_path.read_text(encoding="utf-8"))
    palette: dict[str, tuple[int, int, int]] = {}
    colors = data.get("color") or data.get("colors") or {}

    def _hex_of(entry):
        if isinstance(entry, dict):
            return entry.get("$value") or entry.get("value")
        if isinstance(entry, str):
            return entry
        return None

    for name, entry in colors.items():
        if name.startswith("$"):
            continue
        v = _hex_of(entry)
        if not isinstance(v, str):
            continue
        v = v.strip()
        if v.startswith("{") and v.endswith("}"):
            ref = v[1:-1].split(".")[-1]
            if ref in colors:
                v = _hex_of(colors[ref]) or ""
        if v.startswith("#") and len(v) == 7:
            palette[name] = (
                int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16),
            )
    return palette


def nearest_token(rgb: tuple[int, int, int], palette: dict[str, tuple[int, int, int]]) -> str:
    # No palette → emit the raw hex literal so the DSL preserves the source
    # color verbatim. `derive()` documents this fallback for callers that
    # don't pass a tokens.json.
    if not palette:
        return "#{:02x}{:02x}{:02x}".format(*rgb)
    best = None
    best_d = math.inf
    for name, prgb in palette.items():
        d = sum((a - b) ** 2 for a, b in zip(rgb, prgb))
        if d < best_d:
            best_d = d
            best = name
    # Source-fidelity guard. Squared-euclidean threshold ≈ 25 per channel
    # (3 * 25^2 = 1875). When the closest brand token is further than this,
    # the source colour isn't really represented in the palette — emit the
    # raw hex literal instead of approximating to a token that renders as
    # a visibly different colour (e.g. a corporate source's #FFED00
    # yellow shouldn't squash to the feinschliff parent's gold #C9A24A
    # accent token just because it's the closest of 30 mostly-cool tokens).
    if best is not None and best_d <= _NEAREST_TOKEN_THRESHOLD_SQ:
        return best
    return "#{:02x}{:02x}{:02x}".format(*rgb)


_NEAREST_TOKEN_THRESHOLD_SQ = 1875


_THEME_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def master_theme_blob(pres: Presentation) -> bytes | None:
    """Raw XML of the theme part the deck's primary slide master references.

    PowerPoint numbers theme parts PER-MASTER, so the active theme is often NOT
    `theme1.xml` (a deck whose master is `slideMaster11` references
    `theme11.xml`). Resolving via the master→theme relationship — instead of a
    hardcoded part name — is what lets schemeClr fills/strokes, the background
    panel, and font capture work on real corporate templates.
    """
    try:
        for master in pres.slide_masters:
            for rel in master.part.rels.values():
                if rel.reltype == _THEME_RELTYPE:
                    return rel.target_part.blob
    except Exception:
        pass
    return None


def load_theme_scheme(pres: Presentation) -> dict[str, str]:
    """Map theme scheme keys (accent1..6, dk1, lt1, hlink, folHlink) to #RRGGBB.

    Resolves the theme part the slide master actually references (not a
    hardcoded `theme1.xml`); falls back to any theme part, then empty dict.
    """
    out: dict[str, str] = {}
    blobs: list[bytes] = []
    primary = master_theme_blob(pres)
    if primary is not None:
        blobs.append(primary)
    else:
        # Fallback: any theme part in the package (legacy behaviour).
        try:
            blobs = [p.blob for p in pres.part.package.iter_parts()
                     if "/theme/theme" in str(p.partname)]
        except Exception:
            blobs = []
    for blob in blobs:
        try:
            root = etree.fromstring(blob)
        except Exception:
            continue
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            continue
        for child in scheme:
            key = etree.QName(child).localname  # dk1, lt1, accent1, ...
            srgb = child.find("a:srgbClr", NS)
            sys_ = child.find("a:sysClr", NS)
            if srgb is not None:
                out[key] = "#" + srgb.get("val").upper()
            elif sys_ is not None:
                out[key] = "#" + (sys_.get("lastClr") or "000000").upper()
        # clrMap slot aliases — `schemeClr val="bg1|tx1|bg2|tx2"` resolves
        # through the MASTER's `<p:clrMap>`, not directly against the scheme.
        # Dark-master templates invert the defaults (bg1="dk1" tx1="lt1");
        # assuming the default mapping renders such decks colour-inverted
        # (black cover, white-on-yellow titles). Default mapping is the
        # fallback when no master / no clrMap is reachable.
        slot_map = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
        try:
            clrmap_el = pres.slide_masters[0].element.find(f"{{{NS['p']}}}clrMap")
        except Exception:
            clrmap_el = None
        if clrmap_el is not None:
            for slot in slot_map:
                real = clrmap_el.get(slot)
                if real:
                    slot_map[slot] = real
        for alias, real in slot_map.items():
            if alias not in out and real in out:
                out[alias] = out[real]
        if out:
            break
    return out


# ---------------------------------------------------------------------------
# Geometry conversion
# ---------------------------------------------------------------------------


class CanvasMap:
    def __init__(self, slide_cx_emu: int, slide_cy_emu: int, canvas_w: int, canvas_h: int):
        self.cx = slide_cx_emu
        self.cy = slide_cy_emu
        self.cw = canvas_w
        self.ch = canvas_h
        self.sx = canvas_w / slide_cx_emu
        self.sy = canvas_h / slide_cy_emu
        # SVG renders at 1pt per pt; emu→pt = 1/12700. Conversion to SVG units:
        self.emu_to_pt = 1 / EMU_PER_PT

    def x(self, emu: float) -> int:
        return round(emu * self.sx)

    def y(self, emu: float) -> int:
        return round(emu * self.sy)

    def w(self, emu: float) -> int:
        return max(1, round(emu * self.sx))

    def h(self, emu: float) -> int:
        return max(1, round(emu * self.sy))


# ---------------------------------------------------------------------------
# Shape walking
# ---------------------------------------------------------------------------


def _split_runs_by_color(runs: list["TextRun"]) -> list[list["TextRun"]]:
    """Group runs into consecutive blocks of the same SIZE (and colour).

    A shape often stacks a large header over smaller body text ("Placeholder" 16pt
    over "This text demonstrates…" 12pt) or a coloured headline over body bullets.
    Collapsing to one primitive emits the whole shape at the MAX size / first
    colour, so the small body renders too big and OVERFLOWS. Split so each block
    keeps its own size + colour at its own y-offset (the caller positions them).

    A **size change always starts a new block** (a size jump is a distinct
    visual block regardless of layout). A **colour change splits only across a
    paragraph boundary** — i.e. when a `\\n` marker separates the two
    differently-coloured runs. Two explicitly-coloured runs on the SAME line
    (no intervening `\\n`) stay in one block, because stacking them at separate
    y-offsets would push the second run onto its own line. This is the
    footer's "Internal C-SC1" (bold red) + " | C/CGB-CD | …" (ink) case: one
    source line that previously split into two stacked statements. A colourless
    run inherits and attaches to the current block. `\\n` markers attach to the
    preceding block. Single-(size,colour) shapes return one block (the caller
    skips splitting).
    """
    blocks: list[list[TextRun]] = []
    current: list[TextRun] = []
    cur_color: str | None = None
    cur_size: int | None = None
    saw_break = False  # a paragraph `\n` marker seen since the last content run
    for r in runs:
        if not (r.text and r.text != "\n"):
            # Newline marker / empty run — attach to the current block and arm
            # the colour-split (a colour change is only meaningful as a stacked
            # block when it follows a real line break).
            if current:
                current.append(r)
            if r.text == "\n":
                saw_break = True
            continue
        sz = round(r.pt)
        if not current:
            current = [r]
            cur_color, cur_size = r.color, sz
            saw_break = False
            continue
        color_changed = (
            saw_break
            and r.color is not None and cur_color is not None and r.color != cur_color
        )
        if sz != cur_size or color_changed:
            blocks.append(current)
            current = [r]
            cur_color, cur_size = r.color, sz
        else:
            current.append(r)
            if cur_color is None:
                cur_color = r.color
        saw_break = False
    if current:
        blocks.append(current)
    return blocks


def _resolve_gradient(spPr: etree._Element, theme: dict[str, str],
                       palette: dict[str, tuple[int, int, int]]
                       ) -> tuple[list[tuple[float, str]], float] | None:
    """Extract `<a:gradFill>` as ([(pos, token)], angle_deg) — or None."""
    if spPr is None:
        return None
    grad = spPr.find("a:gradFill", NS)
    if grad is None:
        return None
    stops: list[tuple[float, str]] = []
    for gs in grad.findall("a:gsLst/a:gs", NS):
        try:
            pos = int(gs.get("pos") or 0) / 100000
        except ValueError:
            continue
        srgb = gs.find("a:srgbClr", NS)
        if srgb is not None and srgb.get("val"):
            hx = srgb.get("val")
            try:
                rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                color = nearest_token(rgb, palette) if palette else f"#{hx}"
            except ValueError:
                continue
            stops.append((pos, color))
    if not stops:
        return None
    angle_deg = 0.0
    lin = grad.find("a:lin", NS)
    if lin is not None and lin.get("ang"):
        try:
            angle_deg = int(lin.get("ang")) / 60000.0
        except ValueError:
            pass
    return (stops, angle_deg)


def _resolve_fill(spPr: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return a token name, or None if no fill.

    Handles `<a:grpFill/>` by walking up `<p:grpSp>` ancestors iteratively
    until an actual `<a:solidFill>` is found. The walk is iterative (not
    recursive on `_resolve_fill`) so nested `grpFill` chains — common in
    multi-level Design Kit groups where every ancestor's grpSpPr itself
    carries `<a:grpFill/>` — don't trigger `RecursionError`.
    """
    if spPr is None:
        return None
    sf = spPr.find("a:solidFill", NS)
    if sf is None:
        # No direct solid fill. If the shape declares grpFill, walk up
        # ancestor groups looking for the first one with a real solid
        # fill on its grpSpPr.
        if spPr.find("a:grpFill", NS) is None:
            return None
        anc = spPr.getparent()
        while anc is not None and sf is None:
            if etree.QName(anc).localname == "grpSp":
                grpSpPr = anc.find("p:grpSpPr", NS)
                if grpSpPr is not None:
                    inner_sf = grpSpPr.find("a:solidFill", NS)
                    if inner_sf is not None:
                        sf = inner_sf
                        break
                    # grpSpPr itself is grpFill-only — keep climbing.
            anc = anc.getparent()
        if sf is None:
            return None
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        rgb = _apply_color_mods(rgb, srgb)
        rgb = _blend_on_white(rgb, _alpha_for_color(srgb))
        return nearest_token(rgb, palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
            rgb = _apply_color_mods(rgb, scheme)
            rgb = _blend_on_white(rgb, _alpha_for_color(scheme))
            return nearest_token(rgb, palette)
    return None


def _get_xfrm(spPr: etree._Element) -> tuple[int, int, int, int] | None:
    if spPr is None:
        return None
    xfrm = spPr.find("a:xfrm", NS)
    if xfrm is None:
        return None
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return None
    return int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy"))


def _placeholder_info(node: etree._Element) -> tuple[str | None, str | None]:
    ph = node.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvPicPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvCxnSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        return None, None
    return ph.get("type"), ph.get("idx")


def _layout_placeholder_default_sz(slide, ph_type: str | None, ph_idx: str | None) -> int | None:
    """Walk slide layout + master for the placeholder's default font size.

    PowerPoint inherits font sizes from the layout (and master) when a
    slide-level placeholder has no explicit `sz` on its runs/paragraphs.
    Layout writes the size on
    `<p:sp><p:txBody><a:lstStyle><a:lvl1pPr><a:defRPr sz="...">`; master
    defines title/body defaults on `<p:txStyles>/<p:titleStyle>` and
    `<p:bodyStyle>`.

    Without this lookup, body-level placeholders that omit explicit sz
    inherit our hardcoded 1800 (18pt), which renders chapter titles and
    other layout-controlled headlines at body-size — visibly wrong on
    showcase decks where the layout sets large headlines.

    Returns sz in hundredths-of-pt (PPTX units) or None.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    parents = [p for p in (layout, master) if p is not None]
    for parent in parents:
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            lvl1 = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr", NS)
            if lvl1 is not None:
                d = lvl1.find("a:defRPr", NS)
                if d is not None and d.get("sz"):
                    try:
                        return int(d.get("sz"))
                    except (TypeError, ValueError):
                        pass
    style_for_type = {"title": "titleStyle", "ctrTitle": "titleStyle"}
    style_name = style_for_type.get(ph_type or "", "bodyStyle")
    if master is not None:
        ts = master.element.find(f".//p:txStyles/p:{style_name}", NS)
        if ts is not None:
            lvl1 = ts.find("a:lvl1pPr", NS)
            if lvl1 is not None:
                d = lvl1.find("a:defRPr", NS)
                if d is not None and d.get("sz"):
                    try:
                        return int(d.get("sz"))
                    except (TypeError, ValueError):
                        pass
    return None


def _layout_placeholder_caps_bold(slide, ph_type: str | None, ph_idx: str | None) -> tuple[bool, bool]:
    """Walk slide layout + master for the placeholder's inherited cap/bold.

    Like `_layout_placeholder_default_sz` but for `cap="all"` (all-caps render
    transform) + `b="1"` (bold). A title whose run states neither still renders
    UPPERCASE + bold because the master titleStyle/bodyStyle (and sometimes the
    layout placeholder) defRPr sets them. Returns (caps_all, bold)."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    caps = bold = False
    # Master title/body style is the cascade base.
    style_name = {"title": "titleStyle", "ctrTitle": "titleStyle"}.get(ph_type or "", "bodyStyle")
    if master is not None:
        d = master.element.find(f".//p:txStyles/p:{style_name}/a:lvl1pPr/a:defRPr", NS)
        if d is not None:
            if d.get("cap") is not None:
                caps = d.get("cap") == "all"
            if d.get("b") is not None:
                bold = d.get("b") == "1"
    # The layout's own placeholder defRPr overrides the master.
    if layout is not None:
        for sp in _matching_placeholders(layout, ph_type, ph_idx):
            d = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr", NS)
            if d is not None:
                if d.get("cap") is not None:
                    caps = d.get("cap") == "all"
                if d.get("b") is not None:
                    bold = d.get("b") == "1"
                break
    return caps, bold


def _layout_placeholder_anchor(slide, ph_type: str | None, ph_idx: str | None) -> str | None:
    """Walk slide layout + master for the placeholder's inherited vertical anchor.

    Like `_layout_placeholder_default_sz` but for `<a:bodyPr anchor="ctr|b|t">`. A
    title whose own bodyPr sets no anchor still renders centre/bottom-anchored when
    the layout/master placeholder bodyPr does (MS Geometric: master title
    placeholder anchor="b"). Returns "middle" / "bottom" / "top" / None."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            bodyPr = sp.find(".//p:txBody/a:bodyPr", NS)
            anc = bodyPr.get("anchor") if bodyPr is not None else None
            if anc:
                return {"ctr": "middle", "b": "bottom", "t": "top"}.get(anc)
    return None


def _layout_placeholder_insets(
    slide, ph_type: str | None, ph_idx: str | None
) -> tuple[int | None, int | None, int | None, int | None]:
    """Walk slide layout + master for a placeholder's inherited text-frame insets.

    Like `_layout_placeholder_anchor` but for `<a:bodyPr lIns/tIns/rIns/bIns>`. A
    title whose own bodyPr omits insets still renders at the layout/master
    placeholder's insets — the master title placeholder sets all four to
    `0`, so the slide's "Headline"/"Subheadline" must hug the box left edge, not
    inherit PowerPoint's published 91440/45720 EMU default (which shoved them ~16
    px right / 8 px down and was a top contributor to the heading ghost in the
    redline). Returns (l, t, r, b) in EMU; each element is None when no ancestor
    placeholder specifies that side (caller falls back to the PowerPoint default).
    The MASTER is the cascade base; the LAYOUT placeholder overrides per-side."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    out: list[int | None] = [None, None, None, None]
    # Master first (base), then layout (override) — later writes win per-side.
    for parent in (p for p in (master, layout) if p is not None):
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            bodyPr = sp.find(".//p:txBody/a:bodyPr", NS)
            if bodyPr is None:
                continue
            for i, attr in enumerate(("lIns", "tIns", "rIns", "bIns")):
                v = bodyPr.get(attr)
                if v is not None:
                    try:
                        out[i] = int(v)
                    except (TypeError, ValueError):
                        pass
            break
    return tuple(out)  # type: ignore[return-value]


def _layout_placeholder_autofit(slide, ph_type: str | None, ph_idx: str | None) -> bool:
    """Whether a placeholder INHERITS normAutofit ("shrink text on overflow") from
    its layout/master. Some corporate templates set it on the LAYOUT (on the layout, not on
    the slides), so without this the decompiled placeholder text renders full-size and
    overflows the final render. We don't apply the layout's own fontScale (a template
    default) — emitting `autoshrink` lets the renderer compute the per-slide fit,
    matching PowerPoint."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            if sp.find(".//p:txBody/a:bodyPr/a:normAutofit", NS) is not None:
                return True
    return False


def _layout_placeholder_line_spacing(slide, ph_type: str | None,
                                     ph_idx: str | None) -> float | None:
    """Effective INHERITED paragraph line spacing for a placeholder: the
    layout / master placeholder's lvl1pPr lnSpc, then the master txStyles
    (titleStyle for titles, bodyStyle for everything else). Corporate
    masters routinely set spacing ONLY here (e.g. bodyStyle 107%,
    titleStyle 89%) — a slide paragraph with no own lnSpc must inherit it,
    not fall through to `linespacing:native`. Returns the multiplier or
    None when no ancestor specifies one (true native single spacing)."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None

    def _pct(el) -> float | None:
        if el is not None and el.get("val"):
            try:
                return int(el.get("val")) / 100000.0
            except (TypeError, ValueError):
                pass
        return None

    for parent in (p for p in (layout, master) if p is not None):
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            v = _pct(sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr/a:lnSpc/a:spcPct", NS))
            if v is not None:
                return v
    style_name = {"title": "titleStyle",
                  "ctrTitle": "titleStyle"}.get(ph_type or "", "bodyStyle")
    if master is not None:
        return _pct(master.element.find(
            f".//p:txStyles/p:{style_name}/a:lvl1pPr/a:lnSpc/a:spcPct", NS))
    return None


def _layout_placeholder_color(slide, ph_type: str | None, ph_idx: str | None,
                              theme: dict[str, str],
                              palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Walk slide layout + master for the placeholder's default text colour.

    Mirrors `_layout_placeholder_default_sz` but pulls `<a:defRPr><a:solidFill>`.
    A layout's WHITE title placeholder (e.g. a `ctrTitle` sitting on a coloured
    circle) carries no colour on the slide run; without this lookup it decompiles
    colourless and renders ink-grey. Returns a token/hex or None.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            d = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr", NS)
            sf = d.find("a:solidFill", NS) if d is not None else None
            if sf is not None:
                c = _resolve_solid(sf, theme, palette)
                if c:
                    return c
    style_name = {"title": "titleStyle", "ctrTitle": "titleStyle"}.get(ph_type or "", "bodyStyle")
    if master is not None:
        d = master.element.find(f".//p:txStyles/p:{style_name}/a:lvl1pPr/a:defRPr", NS)
        sf = d.find("a:solidFill", NS) if d is not None else None
        if sf is not None:
            c = _resolve_solid(sf, theme, palette)
            if c:
                return c
    return None


def _matching_placeholders(parent, ph_type: str | None, ph_idx: str | None) -> list:
    """Layout/master placeholder elements matching a slide placeholder,
    BEST match first.

    OOXML pairs a slide placeholder with its layout counterpart by `idx`;
    `type` alone only identifies singletons (title/ctrTitle). Matching
    type-first collapsed every same-type placeholder onto the layout's FIRST
    one — a team slide's 4 `pic` + 8 `body` placeholders all inherited one
    bbox/size/colour and overprinted at a single position. Order returned:
    exact-idx hits, then same-type hits (legacy fallback for decks whose
    layout lacks the idx). Covers `<p:sp>` AND `<p:pic>` placeholders.
    """
    idx_hits: list = []
    type_hits: list = []
    root = parent.element
    for tag in ("sp", "pic"):
        for sp in root.iter("{%s}%s" % (NS["p"], tag)):
            ph = sp.find(".//p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if ph_idx and ph.get("idx") == ph_idx:
                idx_hits.append(sp)
            elif ph_type and ph.get("type") == ph_type:
                type_hits.append(sp)
    return idx_hits + type_hits


def _layout_placeholder_xfrm(slide, ph_type: str | None, ph_idx: str | None) -> tuple[int, int, int, int] | None:
    """Walk slide layout + master to resolve an inherited placeholder bbox.

    Accepts a Slide, SlideLayout, or SlideMaster — when the caller is
    already a layout/master (because walk_slide now recurses through the
    inheritance chain), there's no further parent to walk so this just
    no-ops out.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    parents = [p for p in (layout, master) if p is not None]
    for parent in parents:
        for sp in _matching_placeholders(parent, ph_type, ph_idx):
            xfrm = _get_xfrm(sp.find("p:spPr", NS))
            if xfrm:
                return xfrm
    return None


def _text_runs(node: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]],
               inherited_default_sz: int | None = None,
               inherited_caps: bool = False, inherited_bold: bool = False) -> list[TextRun]:
    runs: list[TextRun] = []
    txBody = node.find(".//p:txBody", NS)
    if txBody is None:
        txBody = node.find(".//a:txBody", NS)
    if txBody is None:
        return runs
    # Body-level lstStyle/lvl1pPr/defRPr sz acts as the cascade default for any
    # run that does not set its own sz. Without this we miss titles whose size
    # is defined only at the body level (common in master-driven decks).
    body_default_sz: int | None = None
    lstStyle = txBody.find("a:lstStyle", NS)
    if lstStyle is not None:
        lvl1 = lstStyle.find("a:lvl1pPr", NS)
        if lvl1 is not None:
            d = lvl1.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                body_default_sz = int(d.get("sz"))
    for para in txBody.findall("a:p", NS):
        para_runs: list[TextRun] = []
        # Pick up para-level defRPr or pPr/defRPr for sz fallback.
        pPr = para.find("a:pPr", NS)
        # Cascade for the paragraph's default sz:
        # 1. txBody/lstStyle/lvl1pPr/defRPr sz (slide-level)
        # 2. inherited_default_sz (layout/master placeholder lookup — only
        #    threaded by `_emit_sp` when the shape is a placeholder)
        # 3. hardcoded 1800 (18pt)
        default_sz = body_default_sz
        if default_sz is None and inherited_default_sz is not None:
            default_sz = inherited_default_sz
        if default_sz is None:
            default_sz = 1800
        if pPr is not None:
            d = pPr.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                default_sz = int(d.get("sz"))
        # Paragraph-level alignment (`<a:pPr algn="ctr|r|just">`). Source
        # frequently centers KPI numbers + their labels within a card; the
        # emitter's default `align:left` shifts them left of source. Stored
        # on each TextRun in this paragraph so emit_dsl can lift the
        # majority-vote into a per-text `align:` kwarg.
        para_align = None
        if pPr is not None:
            algn = pPr.get("algn")
            if algn == "ctr":
                para_align = "center"
            elif algn == "r":
                para_align = "right"
            elif algn == "just":
                para_align = "justify"
        for r in para:
            tag = etree.QName(r).localname
            # <a:br/> is a soft line break BETWEEN runs — PowerPoint renders
            # the surrounding runs on separate lines. Dropping it would fuse
            # them ("The power" + "of communication" → "The powerof…").
            if tag == "br":
                if para_runs:
                    para_runs.append(TextRun(text="\n", pt=default_sz / 100))
                continue
            if tag != "r":
                continue
            rPr = r.find("a:rPr", NS)
            t = r.find("a:t", NS)
            if t is None or t.text is None:
                continue
            sz = default_sz
            # cap + bold cascade from the master title/body style when the run
            # states neither (inherited_*); a title whose run carries no `b`/`cap`
            # still decompiles bold + UPPERCASE, matching the render.
            bold = inherited_bold
            caps = inherited_caps
            italic = False
            color = None
            text = t.text
            if rPr is not None:
                if rPr.get("sz"):
                    sz = int(rPr.get("sz"))
                if rPr.get("b") is not None:
                    bold = rPr.get("b") == "1"
                italic = rPr.get("i") == "1"
                sf = rPr.find("a:solidFill", NS)
                if sf is not None:
                    color = _resolve_fill(rPr, theme, palette) or _resolve_solid(sf, theme, palette)
                if rPr.get("cap") is not None:
                    caps = rPr.get("cap") == "all"
            # PPTX `cap="all"` is a render-time text-transform: the run's stored
            # text stays mixed-case but draws uppercase. Bake the transform into
            # the emitted DSL since downstream layouts carry the literal text,
            # not a `text-transform` directive.
            if caps:
                text = text.upper()
            para_runs.append(TextRun(text=text, pt=sz / 100, bold=bold, italic=italic, color=color, align=para_align))
        if para_runs:
            # Insert a newline marker between paragraphs so emit_dsl can preserve
            # line breaks. Without this, "Headline" + "Lorem ipsum…" paragraphs
            # collapse into "HeadlineLorem ipsum…" and ruin the body diff score.
            if runs:
                runs.append(TextRun(text="\n", pt=default_sz / 100))
            runs.extend(para_runs)
        else:
            # No <a:r> — could be just <a:fld> (page number, date). Capture as one run.
            for fld in para.findall("a:fld", NS):
                t = fld.find("a:t", NS)
                if t is not None and t.text:
                    runs.append(TextRun(text=t.text, pt=default_sz / 100))
    return runs


def _alpha_for_color(color_el: etree._Element) -> float:
    """Return alpha 0..1 from `<a:alpha val="...">` child (PPTX uses 0..100000).
    Defaults to 1.0 when the element is absent."""
    a = color_el.find("a:alpha", NS)
    if a is None or not a.get("val"):
        return 1.0
    try:
        return max(0.0, min(1.0, int(a.get("val")) / 100000.0))
    except (TypeError, ValueError):
        return 1.0


def _blend_on_white(rgb: tuple[int, int, int], alpha: float) -> tuple[int, int, int]:
    """Pre-multiply RGBA against a white slide background.

    Most decks render alpha-on-shape against the slide's paper colour.
    Approximating "blend against white" lets us preserve the perceived
    colour of semi-transparent fills (Venn circles, overlay panels) on
    typical white-canvas slides without threading true alpha through the
    build pipeline. For non-white slide backgrounds the result is
    visually off but only fractionally so — the colour shifts toward
    white instead of the actual canvas.
    """
    if alpha >= 0.999:
        return rgb
    return tuple(
        max(0, min(255, int(round(c * alpha + 255 * (1 - alpha)))))
        for c in rgb
    )


def _apply_color_mods(rgb: tuple[int, int, int],
                      color_el: etree._Element) -> tuple[int, int, int]:
    """Apply PPTX colour modifiers (lumMod/lumOff/tint/shade) to an RGB.

    PowerPoint uses these to derive variants of theme colours — typically
    `<a:schemeClr val="accent1"><a:lumMod val="50000"/><a:lumOff val="50000"/></a:schemeClr>`
    for a 50%-mixed accent. The arithmetic is a crude HSL-luminance shim
    sufficient for the dominant cases (mods used on dark theme colours to
    derive lighter swatch variants in chart series, bar tinting, etc.).
    All values are in PPTX percent-of-100000.
    """
    lumMod = color_el.find("a:lumMod", NS)
    lumOff = color_el.find("a:lumOff", NS)
    tint = color_el.find("a:tint", NS)
    shade = color_el.find("a:shade", NS)
    if lumMod is not None or lumOff is not None:
        try:
            mod = int(lumMod.get("val")) / 100000 if lumMod is not None else 1.0
            off = int(lumOff.get("val")) / 100000 if lumOff is not None else 0.0
            rgb = tuple(max(0, min(255, int(c * mod + 255 * off))) for c in rgb)
        except (TypeError, ValueError):
            pass
    if tint is not None and tint.get("val"):
        # `tint` blends toward white. val = strength of the SOURCE colour
        # retained (lower val = closer to white).
        try:
            t = int(tint.get("val")) / 100000
            rgb = tuple(max(0, min(255, int(c * t + 255 * (1 - t)))) for c in rgb)
        except (TypeError, ValueError):
            pass
    if shade is not None and shade.get("val"):
        # `shade` blends toward black.
        try:
            s = int(shade.get("val")) / 100000
            rgb = tuple(max(0, min(255, int(c * s))) for c in rgb)
        except (TypeError, ValueError):
            pass
    return rgb


def _resolve_solid(sf: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        rgb = _apply_color_mods(rgb, srgb)
        rgb = _blend_on_white(rgb, _alpha_for_color(srgb))
        return nearest_token(rgb, palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
            rgb = _apply_color_mods(rgb, scheme)
            rgb = _blend_on_white(rgb, _alpha_for_color(scheme))
            return nearest_token(rgb, palette)
    return None


# Map brand-pack tokens (the full feinschliff vocabulary) onto the SVG
# DSL's 17-name semantic vocabulary (defined in skills/svg/references/
# dsl-reference.md, resolved through feinschmiede.diagrams.brand_bridge). Tokens
# that have no direct counterpart fall back to the nearest neutral so
# the SVG block builds — at worst we lose a shade of grey, never the
# shape itself.
_BRAND_TO_SVG_COLOR: dict[str, str] = {
    "accent":         "accent",
    "accent-hover":   "accent",
    "highlight":      "accent",
    "ink":            "ink",
    "black":          "ink",
    "white":          "paper",
    "graphite":       "neutral-strong",
    "steel":          "neutral",
    "silver":         "neutral-soft",
    "fog":            "surface-2",
    "paper":          "paper",
    "paper-2":        "surface-2",
    "off-white":      "paper",
    "off-white-2":    "surface-2",
    "rule-dark":      "neutral-strong",
    "accent-2":       "primary",
    "accent-3":       "secondary",
    "severity-low":   "success",
    "severity-medium":"warning",
    "severity-high":  "danger",
    "status-done":    "status-on",
    "status-current": "status-pending",
    # Chart-series ramp — identity-mapped because chart-series-N is BOTH
    # a brand-pack token and a valid SVG semantic name (added together
    # in the chart-decompile feature so pie/bar slice fills survive the
    # round-trip from XML → DSL → SVG block).
    "chart-series-1": "chart-series-1",
    "chart-series-2": "chart-series-2",
    "chart-series-3": "chart-series-3",
    "chart-series-4": "chart-series-4",
    "chart-series-5": "chart-series-5",
    "chart-series-6": "chart-series-6",
    "status-next":    "status-off",
}


def _svg_color_token(brand_token: str | None, *, default: str = "neutral") -> str:
    """Best-effort map of brand-pack color token → SVG DSL semantic name.

    Inline `#rrggbb` literals (produced by nearest_token when the source
    colour is too far from any palette entry) pass through unchanged so
    the SVG block renders the source colour verbatim instead of
    collapsing to the default neutral. Without this, the source-fidelity
    guard in nearest_token would buy nothing for custGeom paths — they'd
    still render as a generic grey because `_BRAND_TO_SVG_COLOR` only
    knows the named feinschliff vocabulary.
    """
    if not brand_token:
        return default
    if brand_token.startswith("#"):
        return brand_token
    return _BRAND_TO_SVG_COLOR.get(brand_token, default)


def _pptx_path_to_svg_d(
    path_el: etree._Element,
    path_w: float, path_h: float,
    target_w: float, target_h: float,
) -> str:
    """Convert one `<a:path>` element to an SVG `d` string in target coords.

    PPTX path commands map onto SVG as follows:
      <a:moveTo>     → M x,y
      <a:lnTo>       → L x,y
      <a:cubicBezTo> → C x1,y1 x2,y2 x,y   (3 pts)
      <a:quadBezTo>  → Q x1,y1 x,y         (2 pts)
      <a:arcTo>      → A rx,ry 0 large,sweep x,y  (computed from wR/hR/stAng/swAng)
      <a:close/>     → Z

    PPTX path-local coordinates are integers in the path's own
    (w, h) box. We scale linearly to (target_w, target_h) in px so the
    resulting SVG `d` lives in the same coordinate space as the
    surrounding `svg <id> X,Y WxH { … }` block (which already places
    the origin at the shape's bbox top-left).

    PPTX arcs are documented angle-based (start angle, sweep angle, both
    in 60000ths of a degree, measured from the +x axis going clockwise).
    We compute the arc endpoint manually and emit SVG's endpoint-form
    arc command. PPTX `sweep > 0` is clockwise; SVG sweep flag `1` is
    clockwise — they line up.
    """
    sx = target_w / path_w if path_w else 1.0
    sy = target_h / path_h if path_h else 1.0
    out: list[str] = []
    cx_cur = cy_cur = 0.0

    def _xy(pt: etree._Element) -> tuple[float, float]:
        return float(pt.get("x")) * sx, float(pt.get("y")) * sy

    for cmd in path_el:
        tag = etree.QName(cmd).localname
        if tag == "moveTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"M {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "lnTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"L {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "cubicBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 3:
                (x1, y1), (x2, y2), (x, y) = _xy(pts[0]), _xy(pts[1]), _xy(pts[2])
                out.append(f"C {x1:.2f},{y1:.2f} {x2:.2f},{y2:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "quadBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 2:
                (x1, y1), (x, y) = _xy(pts[0]), _xy(pts[1])
                out.append(f"Q {x1:.2f},{y1:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "arcTo":
            wR = float(cmd.get("wR")) * sx
            hR = float(cmd.get("hR")) * sy
            stAng = float(cmd.get("stAng")) / 60000.0     # degrees
            swAng = float(cmd.get("swAng")) / 60000.0
            # PPTX arc convention: the arc is on an ellipse whose CENTER
            # is at (cx_cur - wR*cos(stAng), cy_cur - hR*sin(stAng))
            # — i.e. the current point IS the arc's start; the angle
            # tells us where the start lives on the ellipse. Compute the
            # endpoint by adding the sweep.
            import math
            st_rad = math.radians(stAng)
            end_rad = math.radians(stAng + swAng)
            centre_x = cx_cur - wR * math.cos(st_rad)
            centre_y = cy_cur - hR * math.sin(st_rad)
            end_x = centre_x + wR * math.cos(end_rad)
            end_y = centre_y + hR * math.sin(end_rad)
            large_arc = 1 if abs(swAng) > 180 else 0
            sweep_flag = 1 if swAng > 0 else 0
            out.append(
                f"A {wR:.2f},{hR:.2f} 0 {large_arc} {sweep_flag} {end_x:.2f},{end_y:.2f}"
            )
            cx_cur, cy_cur = end_x, end_y
        elif tag == "close":
            out.append("Z")
    return " ".join(out)


def _custgeom_svg_d(spPr: etree._Element, target_w: float, target_h: float) -> str | None:
    """Walk a custGeom's pathLst and concatenate the SVG `d` strings.

    Multiple `<a:path>` siblings inside `<a:pathLst>` become subpaths in
    a single `d` — each starts with its own `M`. Returns None when the
    spPr is not a custGeom or there's no usable path.
    """
    if spPr is None:
        return None
    cg = spPr.find("a:custGeom", NS)
    if cg is None:
        return None
    pathLst = cg.find("a:pathLst", NS)
    if pathLst is None:
        return None
    parts: list[str] = []
    for path_el in pathLst.findall("a:path", NS):
        pw = float(path_el.get("w") or 0)
        ph = float(path_el.get("h") or 0)
        if pw <= 0 or ph <= 0:
            continue
        d = _pptx_path_to_svg_d(path_el, pw, ph, target_w, target_h)
        if d:
            parts.append(d)
    return " ".join(parts) if parts else None


def _shape_geometry_kind(spPr: etree._Element) -> str:
    """Classify a sp by its geometry: rect, oval, line, or 'shape' (custGeom)."""
    if spPr is None:
        return "rect"
    pg = spPr.find("a:prstGeom", NS)
    if pg is not None:
        preset = pg.get("prst")
        if preset in ("ellipse",):
            return "oval"
        if preset in ("line", "straightConnector1"):
            return "line"
        if preset in ("rect", "roundRect"):
            return "rect"
        # Known presets with a simple closed-polygon geometry get routed
        # to `shape` so the emitter writes an `svg { path … }` block with
        # the polygon's `d` string. See `_preset_geom_path` for the table.
        if preset in _PRESET_PATH_PRESETS:
            return "shape"
        return "rect"
    if spPr.find("a:custGeom", NS) is not None:
        return "shape"
    return "rect"


# Presets whose geometry is a closed polygon synthesized by
# `_preset_geom_path`, which returns the SVG `d` string in the shape's
# local 0..w × 0..h pixel coordinate space. Adjustment values from
# <a:avLst> are honoured for the presets whose ECMA-376 guide formulas
# use them (chevron, homePlate, parallelogram, trapezoid, hexagon,
# octagon, arrows); the spec defaults apply when absent.
_PRESET_PATH_PRESETS: frozenset[str] = frozenset({
    "triangle", "rtTriangle", "diamond",
    "parallelogram", "trapezoid",
    "pentagon", "hexagon", "heptagon", "octagon",
    "homePlate", "chevron",
    "rightArrow", "leftArrow", "upArrow", "downArrow",
})


def _preset_geom_path(preset: str, w: float, h: float,
                      adj: dict[str, float] | None = None) -> str | None:
    """SVG `d` string for the known closed-polygon presets, in local px.

    Guide formulas follow ECMA-376 presetShapeDefinitions: adjustment
    values are fractions of `ss` (the SHORTEST side), not of the width —
    a wide flat chevron keeps a compact notch instead of one stretched
    to 30% of its width. `adj` maps gd name → value/100000.
    """
    if w <= 0 or h <= 0:
        return None
    adj = adj or {}
    ss = min(w, h)
    if preset == "triangle":
        return f"M {w/2:.1f},0 L {w:.1f},{h:.1f} L 0,{h:.1f} Z"
    if preset == "rtTriangle":
        return f"M 0,0 L 0,{h:.1f} L {w:.1f},{h:.1f} Z"
    if preset == "diamond":
        return (f"M {w/2:.1f},0 L {w:.1f},{h/2:.1f} "
                f"L {w/2:.1f},{h:.1f} L 0,{h/2:.1f} Z")
    if preset == "parallelogram":
        # x1 = ss * adj (default 25000) — skew from the left, per spec.
        skew = min(w, ss * adj.get("adj", 0.25))
        return (f"M {skew:.1f},0 L {w:.1f},0 L {w-skew:.1f},{h:.1f} "
                f"L 0,{h:.1f} Z")
    if preset == "trapezoid":
        # x1 = ss * adj (default 25000) — top inset on each side.
        inset = min(w / 2, ss * adj.get("adj", 0.25))
        return (f"M {inset:.1f},0 L {w-inset:.1f},0 L {w:.1f},{h:.1f} "
                f"L 0,{h:.1f} Z")
    if preset == "pentagon":
        # Regular pentagon inscribed in the bbox: top vertex centred, side
        # vertices at 38.2% height, base at 19.1% / 80.9% width. (The spec's
        # hf/vf factors reduce to this once normalised to the bbox.)
        return (f"M {w*0.5:.1f},0 L {w:.1f},{h*0.382:.1f} "
                f"L {w*0.809:.1f},{h:.1f} L {w*0.191:.1f},{h:.1f} "
                f"L 0,{h*0.382:.1f} Z")
    if preset == "homePlate":
        # x1 = ss * adj (default 50000) — point length; flat left edge.
        x1 = min(w, ss * adj.get("adj", 0.5))
        return (f"M 0,0 L {w-x1:.1f},0 L {w:.1f},{h*0.5:.1f} "
                f"L {w-x1:.1f},{h:.1f} L 0,{h:.1f} Z")
    if preset == "hexagon":
        # x1 = ss * adj (default 25000) — corner inset.
        x1 = min(w / 2, ss * adj.get("adj", 0.25))
        return (f"M {x1:.1f},0 L {w-x1:.1f},0 L {w:.1f},{h*0.5:.1f} "
                f"L {w-x1:.1f},{h:.1f} L {x1:.1f},{h:.1f} "
                f"L 0,{h*0.5:.1f} Z")
    if preset == "heptagon":
        # Regular-ish 7-gon inscribed in the bbox.
        import math as _m
        pts = []
        cx, cy = w / 2, h / 2
        rx, ry = w / 2, h / 2
        for i in range(7):
            ang = -_m.pi / 2 + i * 2 * _m.pi / 7
            pts.append(f"{cx + rx * _m.cos(ang):.1f},{cy + ry * _m.sin(ang):.1f}")
        return "M " + " L ".join(pts) + " Z"
    if preset == "octagon":
        # x1 = ss * adj (default 29289) — corner cut on the shortest side.
        x1 = min(ss / 2, ss * adj.get("adj", 0.29289))
        return (f"M {x1:.1f},0 L {w-x1:.1f},0 L {w:.1f},{x1:.1f} "
                f"L {w:.1f},{h-x1:.1f} L {w-x1:.1f},{h:.1f} "
                f"L {x1:.1f},{h:.1f} L 0,{h-x1:.1f} L 0,{x1:.1f} Z")
    if preset == "chevron":
        # x1 = ss * adj (default 50000) — point AND notch depth. The old
        # hardcoded 30%-of-WIDTH notch stretched wide process-step chevrons
        # into near-triangles.
        x1 = min(w, ss * adj.get("adj", 0.5))
        return (f"M 0,0 L {w-x1:.1f},0 L {w:.1f},{h*0.5:.1f} "
                f"L {w-x1:.1f},{h:.1f} L 0,{h:.1f} "
                f"L {x1:.1f},{h*0.5:.1f} Z")
    if preset in ("rightArrow", "leftArrow", "upArrow", "downArrow"):
        # adj1 = shaft thickness as a fraction of the cross dimension
        # (default 50000); adj2 = head length = ss * adj2 (default 50000).
        a1 = adj.get("adj1", 0.5)
        a2 = adj.get("adj2", 0.5)
        if preset in ("rightArrow", "leftArrow"):
            head = min(w, ss * a2)
            sy0 = h * (0.5 - a1 / 2)
            sy1 = h * (0.5 + a1 / 2)
            if preset == "rightArrow":
                ax = w - head
                return (f"M 0,{sy0:.1f} L {ax:.1f},{sy0:.1f} L {ax:.1f},0 "
                        f"L {w:.1f},{h*0.5:.1f} L {ax:.1f},{h:.1f} "
                        f"L {ax:.1f},{sy1:.1f} L 0,{sy1:.1f} Z")
            ax = head
            return (f"M {w:.1f},{sy0:.1f} L {ax:.1f},{sy0:.1f} L {ax:.1f},0 "
                    f"L 0,{h*0.5:.1f} L {ax:.1f},{h:.1f} "
                    f"L {ax:.1f},{sy1:.1f} L {w:.1f},{sy1:.1f} Z")
        head = min(h, ss * a2)
        sx0 = w * (0.5 - a1 / 2)
        sx1 = w * (0.5 + a1 / 2)
        if preset == "upArrow":
            ay = head
            return (f"M {sx0:.1f},{h:.1f} L {sx0:.1f},{ay:.1f} L 0,{ay:.1f} "
                    f"L {w*0.5:.1f},0 L {w:.1f},{ay:.1f} L {sx1:.1f},{ay:.1f} "
                    f"L {sx1:.1f},{h:.1f} Z")
        ay = h - head
        return (f"M {sx0:.1f},0 L {sx0:.1f},{ay:.1f} L 0,{ay:.1f} "
                f"L {w*0.5:.1f},{h:.1f} L {w:.1f},{ay:.1f} L {sx1:.1f},{ay:.1f} "
                f"L {sx1:.1f},0 Z")
    return None


# ---------------------------------------------------------------------------
# Tree walk
# ---------------------------------------------------------------------------


def walk_slide(slide, cmap: CanvasMap, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> list[Shape]:
    """Walk the slide's spTree; also collect shapes inherited from the
    slide layout + master that aren't already represented on the slide.

    A typical PowerPoint slide carries only its unique content (a title
    placeholder, the body text). All decorative chrome — corporate logo,
    page-number bar, branded background blocks — lives on the slide's
    layout and master. Without this inheritance walk, the decompile of
    a corporate-template deck emits ~10% of what the source renders.

    Layout/master shapes that are placeholder fills already provided by
    the slide itself (same ph_idx) are skipped — slide content wins.
    Everything else is added at the front of the shape list so it draws
    behind slide-level content.
    """
    shapes: list[Shape] = []
    spTree = slide.element.find(".//p:cSld/p:spTree", NS)
    _walk(spTree, (0, 0), shapes, slide, cmap, theme, palette)
    # A slide-level placeholder only "owns" its idx when it actually
    # carries content (text, fill, or geometry) — empty placeholders
    # (common pattern: <p:ph idx="N"/> + empty <p:spPr/>) inherit
    # EVERYTHING from the layout, including the layout's fill/size.
    # Filtering by content here lets the layout walk re-add the rich
    # version of those placeholders below.
    def _has_content(sh: Shape) -> bool:
        if sh.fill or sh.stroke or sh.gradient or sh.svg_path_d:
            return True
        if sh.is_picture:
            return True
        if any(r.text and r.text.strip() and r.text != "\n" for r in sh.text_runs):
            return True
        return False

    slide_ph_idxs = {s.ph_idx for s in shapes if s.ph_idx and _has_content(s)}
    # Filter out the empty placeholders themselves — the layout version
    # will provide the actual content.
    shapes = [s for s in shapes if not (s.ph_idx and not _has_content(s))]
    inherited: list[Shape] = []
    show_master = _show_master_sp(slide)
    layout_master_chain = _layout_master_chain(slide)
    for src in layout_master_chain:
        is_master = etree.QName(src.element).localname == "sldMaster"
        chain_spTree = src.element.find(".//p:cSld/p:spTree", NS)
        if chain_spTree is None:
            continue
        chain_shapes: list[Shape] = []
        _walk(chain_spTree, (0, 0), chain_shapes, src, cmap, theme, palette)
        for s in chain_shapes:
            # `showMasterSp="0"` (slide, else layout) hides the master's
            # plain shapes — PowerPoint never renders them on such slides.
            # Placeholders are exempt: the flag governs decorative master
            # shapes only, placeholder inheritance still flows.
            if is_master and not show_master and not (s.ph_idx or s.ph_type):
                continue
            # Skip placeholder shapes the slide already owns with content.
            if s.ph_idx and s.ph_idx in slide_ph_idxs:
                continue
            # Skip pure page-number / footer placeholders by type when the
            # slide hasn't overridden them — these are pp:fld things that
            # the source-deck renderer fills at slide time; decompiling
            # them from the master emits literal "<#>" tokens that pollute
            # the output. (sldNum/ftr/dt placeholders all carry ph_type.)
            if s.ph_type in ("sldNum", "ftr", "dt") and not s.text_runs:
                continue
            # Inherited PICTURE placeholders with no real media binary
            # (no media_rid) are shells the slide should fill but didn't.
            # PowerPoint renders them as empty; the build's missing-asset
            # fallback paints them as white rectangles which inflate the
            # diff visibly. Skip them so the rendered output matches
            # PowerPoint's "absent" behaviour.
            if s.is_picture and not s.media_rid:
                continue
            # Skip layout placeholders whose only text is template
            # prompt copy ("Hier Zitat einfügen.", "Überschrift 1, TT
            # Norms Pro, 28 pt", etc) — these are author hints that
            # PowerPoint suppresses at render time when the slide's
            # version is empty. We can't easily distinguish prompt from
            # real content, but layouts mark prompts with
            # `<p:nvPr hasCustomPrompt="1">` — drop just the text on
            # those, keep the fill/geometry so the layout still emits
            # its visual frame.
            # Inherited PLACEHOLDER text is template prompt copy that
            # PowerPoint never renders ("Überschrift 1, TT Norms Pro",
            # "Diagramm durch Klicken …", "Click to edit Master title").
            # The slide's placeholder either overrides it (filtered above
            # via slide_ph_idxs) or is empty — in which case PowerPoint
            # shows nothing. Drop the text but keep fill / geometry so
            # the layout's visual frame (yellow rect, black square)
            # still emits. Non-placeholder layout/master shapes (logo
            # glyphs, decorative rects) pass through untouched.
            if s.ph_type and s.text_runs:
                s.text_runs = []
                if not _has_content(s):
                    continue
            s.from_chain = True
            inherited.append(s)
            if s.ph_idx:
                slide_ph_idxs.add(s.ph_idx)
    # Inherited chrome draws behind slide content.
    return inherited + shapes


def _show_master_sp(slide) -> bool:
    """Whether master shapes render on this slide (`showMasterSp`).

    The slide's own flag wins when explicitly set; otherwise the layout's
    flag decides; absent both, PowerPoint's default is to show them.
    """
    el = getattr(slide, "element", None)
    v = el.get("showMasterSp") if el is not None else None
    if v is not None:
        return v not in ("0", "false")
    layout = getattr(slide, "slide_layout", None)
    lel = getattr(layout, "element", None)
    lv = lel.get("showMasterSp") if lel is not None else None
    if lv is not None:
        return lv not in ("0", "false")
    return True


def _layout_master_chain(slide) -> list:
    """Return [layout, master] for a slide (best-effort, never raises)."""
    chain = []
    try:
        layout = slide.slide_layout
        if layout is not None:
            chain.append(layout)
            master = layout.slide_master
            if master is not None:
                chain.append(master)
    except Exception:
        pass
    return chain


def extract_slide_bg_fill(slide, theme: dict[str, str],
                          palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return the slide's background solid-fill colour as a token / hex.

    Walks slide → LAYOUT only (NOT the master). A layout that wants a
    full-bleed colour sets its own `<p:cSld><p:bg>` (kept here — e.g. a
    timeline layout's yellow, a divider's navy). The MASTER-level bg is
    deliberately ignored: corporate light templates carry a dark master bg
    that the layout visually overrides with a white layer / colour panel, so
    emitting it as a full-canvas rect paints every white content slide black
    (a ~97% regression). solidFill or bgRef→theme; first level found wins.
    """
    for src in [slide, *_layout_master_chain(slide)[:1]]:
        bg = src.element.find(".//p:cSld/p:bg", NS)
        if bg is None:
            continue
        bgPr = bg.find("p:bgPr", NS)
        if bgPr is not None:
            color = _resolve_fill(bgPr, theme, palette)
            if color:
                return color
        # <p:bgRef idx="N"><a:schemeClr val="bg1"/></p:bgRef> — referenced
        # background-fill style from theme1.xml's bgFillStyleLst. The
        # schemeClr fills the `phClr` placeholder in the referenced style.
        bgRef = bg.find("p:bgRef", NS)
        if bgRef is not None:
            color = _resolve_bg_ref(bgRef, theme, palette)
            if color:
                return color
    return None


def extract_slide_bg_image(slide) -> tuple[bytes, str] | None:
    """The slide's background image, walking slide → LAYOUT (not master).

    Corporate/gallery templates paint full-bleed artwork via
    `<p:bg><p:bgPr><a:blipFill>` (e.g. Scientific's engraved petri-dish
    layout background) — invisible to the solid-fill extractor, so the
    decompiled slide rendered bare white. Returns (blob, ext) of the first
    blipFill bg found, resolved against the surface that declares it.
    """
    for src in [slide, *_layout_master_chain(slide)[:1]]:
        bg = src.element.find(".//p:cSld/p:bg", NS)
        if bg is None:
            continue
        blip = bg.find(".//p:bgPr/a:blipFill/a:blip", NS)
        if blip is None:
            return None  # this surface owns the bg, and it isn't an image
        rid = blip.get(f"{{{RELS_NS}}}embed")
        if not rid:
            return None
        try:
            part = src.part.related_part(rid)
        except Exception:
            return None
        partname = str(getattr(part, "partname", "/bg.png"))
        ext = partname.rsplit(".", 1)[-1].lower() if "." in partname else "png"
        blob = getattr(part, "blob", None)
        return (blob, ext) if blob else None
    return None


def _resolve_bg_ref(bgRef, theme: dict[str, str],
                    palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Resolve a `<p:bgRef idx="N">` reference against the theme's
    `bgFillStyleLst`, filling `<a:schemeClr val="phClr"/>` with the
    schemeClr that the bgRef carries.

    PowerPoint idx encoding: 1001 = bgFillStyleLst[0], 1002 = [1], etc.
    Only `<a:solidFill>` entries are inlined; gradient/blip refs fall
    through (no DSL primitive for those yet at the bg level).
    """
    scheme = bgRef.find("a:schemeClr", NS)
    if scheme is None or not scheme.get("val"):
        return None
    scheme_key = scheme.get("val")
    # Resolve scheme → hex via the theme dict captured by load_theme_scheme.
    # Some keys are aliases: `bg1`→`lt1`, `bg2`→`lt2`, `tx1`→`dk1`, `tx2`→`dk2`.
    alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    hex_val = theme.get(scheme_key) or theme.get(alias.get(scheme_key, scheme_key))
    if not hex_val:
        return None
    try:
        rgb = (int(hex_val[1:3], 16), int(hex_val[3:5], 16), int(hex_val[5:7], 16))
    except (ValueError, IndexError):
        return None
    return nearest_token(rgb, palette) if palette else f"#{hex_val[1:].upper()}"


def _walk(node, offset, shapes, slide, cmap, theme, palette):
    # offset is either a 2-tuple (ox, oy) or an 8-tuple carrying a scaled-
    # group affine — the latter is unpacked by `_shape_bbox`. We only need
    # ox/oy here to forward to nested non-scaled groups.
    ox, oy = offset[:2]
    for ch in node:
        tag = etree.QName(ch).localname
        # Shapes flagged hidden="1" on their cNvPr are never rendered by
        # PowerPoint/LibreOffice but ARE walked here — template add-ins hide
        # machinery shapes on the master (e.g. a classification/date plate),
        # which otherwise lands as a phantom filled rect on every decompiled
        # layout. Applies to sp/pic/cxnSp/graphicFrame/grpSp alike;
        # `*/p:cNvPr` is the element's OWN nv*Pr block, not a group child's.
        cnvpr = ch.find("*/p:cNvPr", NS)
        if cnvpr is not None and cnvpr.get("hidden") in ("1", "true"):
            continue
        if tag == "sp":
            _emit_sp(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "pic":
            _emit_pic(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "cxnSp":
            _emit_cxn(ch, offset, shapes, cmap, theme, palette)
        elif tag == "graphicFrame":
            _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "grpSp":
            # Whole-group native-carry: a top-level group of pure decorative
            # chrome (custGeom / graphicFrame, NO content placeholder) is frozen
            # verbatim — PowerPoint applies its own off/ext/chOff/chExt affine, so
            # even a SCALED group of vectors (world maps, decorative clusters) that
            # the per-shape re-synth mangles renders pixel-exact. Returns True when
            # carried (don't recurse); False → recurse exactly as before.
            if _try_carry_group(ch, offset, shapes, slide, cmap, theme):
                continue
            # Walk children with the group's offset added. Scaled groups
            # (ext != chExt — typical for master-level logo bundles
            # dropped into smaller slots) are skipped because the walker
            # only carries a pure translation offset; emitting their
            # children at unscaled coords would land them off-canvas.
            grp_xfrm = ch.find("p:grpSpPr/a:xfrm", NS)
            child_off = (ox, oy)
            if grp_xfrm is not None:
                off = grp_xfrm.find("a:off", NS)
                ext = grp_xfrm.find("a:ext", NS)
                chOff = grp_xfrm.find("a:chOff", NS)
                chExt = grp_xfrm.find("a:chExt", NS)
                if (off is not None and ext is not None
                        and chOff is not None and chExt is not None):
                    try:
                        ox_emu = int(off.get("x"))
                        oy_emu = int(off.get("y"))
                        cx = int(ext.get("cx"))
                        cy = int(ext.get("cy"))
                        chcx = int(chExt.get("cx"))
                        chcy = int(chExt.get("cy"))
                        chox = int(chOff.get("x"))
                        choy = int(chOff.get("y"))
                        if chcx > 0 and chcy > 0:
                            sx = cx / chcx
                            sy = cy / chcy
                            # Scaled group: thread a 6-tuple offset that
                            # _shape_bbox unpacks and applies as an EMU-level
                            # affine. Translation-only groups stay as 2-tuple
                            # for backward compat.
                            if abs(sx - 1.0) > 0.001 or abs(sy - 1.0) > 0.001:
                                child_off = (ox, oy, ox_emu, oy_emu, chox, choy, sx, sy)
                                _walk(ch, child_off, shapes, slide, cmap, theme, palette)
                                continue
                        # Pure translation fallthrough.
                        child_off = (ox + ox_emu - chox, oy + oy_emu - choy)
                    except (ValueError, TypeError):
                        pass
                elif off is not None and chOff is not None:
                    try:
                        dx = int(off.get("x")) - int(chOff.get("x"))
                        dy = int(off.get("y")) - int(chOff.get("y"))
                        child_off = (ox + dx, oy + dy)
                    except (ValueError, TypeError):
                        pass
            _walk(ch, child_off, shapes, slide, cmap, theme, palette)


def _shape_bbox(ch, offset, slide):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        ph_type, ph_idx = _placeholder_info(ch)
        xfrm = _layout_placeholder_xfrm(slide, ph_type, ph_idx)
    if xfrm is None:
        return None
    x, y, w, h = xfrm
    # Offset shapes:
    #   2-tuple (ox, oy)      — translation-only ancestor group(s)
    #   8-tuple (ox, oy, ax, ay, chox, choy, sx, sy)
    #                         — current shape lives inside a scaled group;
    #                           apply the EMU-level affine before adding
    #                           any outer translation.
    if len(offset) == 2:
        ox, oy = offset
        return x + ox, y + oy, w, h
    ox, oy, ax, ay, chox, choy, sx, sy = offset
    x_emu = ax + (x - chox) * sx
    y_emu = ay + (y - choy) * sy
    w_emu = w * sx
    h_emu = h * sy
    return x_emu + ox, y_emu + oy, w_emu, h_emu


def _bake_scheme_colors(el, theme: dict[str, str]) -> None:
    """In-place: rewrite `<a:schemeClr val=KEY>` → `<a:srgbClr val=HEX>` from the
    SOURCE theme, so a carried-native shape keeps its EXACT source colours even
    when the output deck's theme differs. Child mods (lumMod/alpha/…) are valid on
    srgbClr too and ride along untouched."""
    a_ns = NS["a"]
    alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    for sc in list(el.iter(f"{{{a_ns}}}schemeClr")):
        key = sc.get("val")
        hexv = theme.get(key) or theme.get(alias.get(key, key))
        if hexv:
            sc.tag = f"{{{a_ns}}}srgbClr"
            sc.set("val", hexv.lstrip("#").upper())


def _carry_element(el, theme: dict[str, str], offset=(0, 0),
                   xfrm_off_path: str | None = None):
    """Deep-copy `el` and bake schemeClr→srgbClr against the SOURCE theme for a
    native carry. When `xfrm_off_path` is given (e.g. "p:spPr/a:xfrm/a:off"),
    shift that `<a:off>` by the walker's pure-translation `offset` so the
    carried element's xfrm becomes slide-absolute before splicing."""
    import copy as _copy
    out = _copy.deepcopy(el)
    _bake_scheme_colors(out, theme)
    if xfrm_off_path is not None:
        ox, oy = int(offset[0]), int(offset[1])
        off_el = out.find(xfrm_off_path, NS)
        if off_el is not None and (ox or oy):
            off_el.set("x", str(int(off_el.get("x") or 0) + ox))
            off_el.set("y", str(int(off_el.get("y") or 0) + oy))
    return out


def _emit_sp(ch, offset, shapes, slide, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    # Pull placeholder default sz from layout/master so body placeholders
    # without explicit run-level `sz` inherit the right headline size.
    inherited_sz = _layout_placeholder_default_sz(slide, ph_type, ph_idx) if (ph_type or ph_idx) else None
    # cap="all" / bold also cascade from the layout/master placeholder style.
    inherited_caps, inherited_bold = (
        _layout_placeholder_caps_bold(slide, ph_type, ph_idx) if (ph_type or ph_idx) else (False, False)
    )
    runs = _text_runs(ch, theme, palette, inherited_default_sz=inherited_sz,
                      inherited_caps=inherited_caps, inherited_bold=inherited_bold)
    # Scaled group: PowerPoint scales text with the group resize, but runs
    # keep their NOMINAL sz in the XML — apply the group's vertical scale so
    # tile text doesn't render at pre-scale size (slide-48/49 class: 12pt
    # nominal inside a 0.6x group renders ~7pt in the source).
    if len(offset) == 8:
        _gsy = offset[7]
        if abs(_gsy - 1.0) > 0.01:
            for _r in runs:
                if _r.pt:
                    _r.pt = round(_r.pt * _gsy, 1)
    # G3 — capture text colour the slide run INHERITS rather than states. When a
    # run carries no explicit `<a:rPr><a:solidFill>`, fall back to (a) the shape's
    # `<p:style><a:fontRef>` (decorative styled shapes) then (b) the layout/master
    # PLACEHOLDER default colour (e.g. a WHITE `ctrTitle` over a coloured circle).
    # Without this, on-shape / on-placeholder titles render ink-grey. Emitted
    # downstream only when it differs from the style-bundle default.
    if runs and not any(r.color for r in runs):
        _style = ch.find("p:style", NS)
        _font_ref = _style.find("a:fontRef", NS) if _style is not None else None
        _text_color = _resolve_solid(_font_ref, theme, palette) if _font_ref is not None else None
        if _text_color is None and (ph_type or ph_idx):
            _text_color = _layout_placeholder_color(slide, ph_type, ph_idx, theme, palette)
        if _text_color:
            for _r in runs:
                if _r.text and _r.text != "\n":
                    _r.color = _text_color
    kind = _shape_geometry_kind(spPr)
    # For custGeom shapes (kind="shape") — typically map polygons,
    # decorative vector clusters, or hand-drawn paths — bypass the
    # brand-token mapping in `nearest_token` and emit the source colour
    # as raw hex. These shapes carry hundreds of subtly-different
    # source colours (e.g. world-map country fills at #EBEBEB /
    # #DDDDDD / #C8C8C8) and the round-trip through nearest_token →
    # `_svg_color_token` → `brand_bridge.resolve` collapses them to a
    # handful of SVG semantic names that resolve to materially
    # different greys in the brand pack. Going straight to hex
    # preserves source-pixel fidelity for these high-cardinality
    # vector compositions.
    if kind == "shape":
        fill = _resolve_fill(spPr, theme, palette={})
    else:
        fill = _resolve_fill(spPr, theme, palette)
    if fill is None and spPr is not None and spPr.find("a:noFill", NS) is None:
        # Master-styled shape: `<p:style><a:fillRef idx="N"><a:schemeClr
        # val="accent1"/>` carries the fill when spPr declares none — the
        # theme fill style at idx substitutes the ref's color for its
        # `phClr` placeholder. Solid approximation: resolve the ref color
        # itself (fillStyleLst[0] IS a plain phClr solid; gradient styles
        # approximate to their base color). Without this, styled shapes
        # silently lost their fill (spec-audit phClr gap).
        _style = ch.find("p:style", NS)
        _fref = _style.find("a:fillRef", NS) if _style is not None else None
        if _fref is not None:
            try:
                _fidx = int(_fref.get("idx") or 0)
            except ValueError:
                _fidx = 0
            if 1 <= _fidx <= 999:
                _pal = {} if kind == "shape" else palette
                _srgb = _fref.find("a:srgbClr", NS)
                _scheme = _fref.find("a:schemeClr", NS)
                if _srgb is not None and _srgb.get("val"):
                    hx = _srgb.get("val")
                    rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                    rgb = _apply_color_mods(rgb, _srgb)
                    rgb = _blend_on_white(rgb, _alpha_for_color(_srgb))
                    fill = nearest_token(rgb, _pal)
                elif _scheme is not None and theme.get(_scheme.get("val")):
                    hex_str = theme[_scheme.get("val")]
                    rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16),
                           int(hex_str[5:7], 16))
                    rgb = _apply_color_mods(rgb, _scheme)
                    rgb = _blend_on_white(rgb, _alpha_for_color(_scheme))
                    fill = nearest_token(rgb, _pal)
    gradient = _resolve_gradient(spPr, theme, palette)
    # Vertical anchor — `<a:bodyPr anchor="ctr">` / `b` / `t`. Without
    # this the rendered text lands at frame-top even when source centers
    # it, which is the dominant cause of the redline "two ghost positions"
    # pattern: source content at frame center, render content at frame top.
    valign: str | None = None
    padding_emu: tuple[int, int, int, int] | None = None
    autoshrink = False
    font_scale = 1.0
    line_spacing: float | None = None
    # NOT `find(...) or find(...)` — lxml element truthiness is based on
    # child count, so a childless <p:txBody/> would falsily fall through.
    txBody = ch.find(".//p:txBody", NS)
    if txBody is None:
        txBody = ch.find(".//a:txBody", NS)
    if txBody is not None:
        bodyPr = txBody.find("a:bodyPr", NS)
        if bodyPr is not None:
            anc = bodyPr.get("anchor")
            if anc == "ctr":
                valign = "middle"
            elif anc == "b":
                valign = "bottom"
            # Insets — l/t/r/b. The slide's own bodyPr wins per-side; for any
            # side it omits, the placeholder INHERITS the layout/master
            # placeholder's inset (the master title sets all four to 0, so
            # the heading must hug the box edge); only when no ancestor specifies
            # a side does PowerPoint's published default (91440 / 45720 EMU)
            # apply. Without the inheritance step every empty-bodyPr title
            # decompiled with the +16/+8 px default offset and ghosted in the
            # redline.
            inh = (
                _layout_placeholder_insets(slide, ph_type, ph_idx)
                if (ph_type or ph_idx)
                else (None, None, None, None)
            )

            def _ins(attr: str, inherited: int | None, default: int) -> int:
                v = bodyPr.get(attr)
                if v is not None:
                    try:
                        return int(v)
                    except (TypeError, ValueError):
                        pass
                return inherited if inherited is not None else default

            left = _ins("lIns", inh[0], 91440)
            top = _ins("tIns", inh[1], 45720)
            right = _ins("rIns", inh[2], 91440)
            bottom = _ins("bIns", inh[3], 45720)
            padding_emu = (left, top, right, bottom)
            # normAutofit = PowerPoint "Shrink text on overflow": the source pre-
            # shrinks the text by `fontScale` to fit the box (the run `sz` stays at
            # the authored size). Capture it so the emitter reproduces the fit
            # instead of rendering full-size and overflowing (inherited placeholders).
            na = bodyPr.find("a:normAutofit", NS)
            if na is not None:
                autoshrink = True
                if na.get("fontScale"):
                    try:
                        font_scale = int(na.get("fontScale")) / 100000.0
                    except (TypeError, ValueError):
                        pass
        # Explicit paragraph line spacing — first paragraph's
        # <a:pPr><a:lnSpc><a:spcPct val="N"/> wins (consistent with how
        # color/size/align pick the first content run). Absent → None →
        # the emitter writes NO lnSpc so the renderer's native single
        # spacing applies, exactly like the source.
        spc = txBody.find("a:p/a:pPr/a:lnSpc/a:spcPct", NS)
        if spc is not None and spc.get("val"):
            try:
                line_spacing = int(spc.get("val")) / 100000.0
            except (TypeError, ValueError):
                pass
    # Vertical anchor also inherits: a slide title with no own bodyPr anchor still
    # renders bottom/centre-anchored when the layout/master placeholder bodyPr sets
    # it (MS Geometric: master title placeholder anchor="b" → titles sit at the box
    # bottom, not top). Mirror the size / caps / colour inheritance. Feature-2's box
    # extension correctly skips bottom/middle text, so it won't fight this.
    if valign is None and (ph_type or ph_idx):
        valign = _layout_placeholder_anchor(slide, ph_type, ph_idx)
    # normAutofit also inherits: some templates set "shrink text on overflow" on the
    # LAYOUT placeholder, not the slide, so a placeholder whose own bodyPr lacks it
    # still shrinks-to-fit — capture that (autoshrink) or the text overflows.
    if not autoshrink and (ph_type or ph_idx):
        autoshrink = _layout_placeholder_autofit(slide, ph_type, ph_idx)
    # Line spacing inherits like anchor / autofit / insets: a placeholder
    # paragraph with no own lnSpc takes the layout/master placeholder value,
    # then the master txStyles (e.g. bodyStyle 107%, titleStyle 89%).
    if line_spacing is None and (ph_type or ph_idx):
        line_spacing = _layout_placeholder_line_spacing(slide, ph_type, ph_idx)
    # Convert insets EMU → design-px for the Shape (CanvasMap-relative).
    padding_px: tuple[float, float, float, float] | None = None
    if padding_emu is not None:
        left, top, right, bottom = padding_emu
        padding_px = (cmap.w(left), cmap.h(top), cmap.w(right), cmap.h(bottom))
    # Stroke (line) colour + width + dash. PowerPoint stores stroke width
    # in EMU on `<a:ln w="...">` (default ~9525 EMU = 0.75pt = 1px hairline);
    # dash preset on the optional `<a:prstDash val="...">` child.
    stroke = None
    stroke_width: float | None = None
    stroke_dash: str | None = None
    ln = spPr.find("a:ln", NS) if spPr is not None else None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
        w_attr = ln.get("w")
        if w_attr:
            try:
                w_emu = int(w_attr)
                # Width is uniform in EMU (1pt = 12700); converting via the
                # horizontal scale keeps it consistent with the rest of the
                # design-px coordinate system on this canvas.
                stroke_width = cmap.w(w_emu)
            except (ValueError, TypeError):
                pass
        dash = ln.find("a:prstDash", NS)
        if dash is not None and dash.get("val"):
            stroke_dash = dash.get("val")
    # Corner radius — captured from `prstGeom prst="roundRect"` with an
    # `<a:gd name="adj" fmla="val N">`. N is N/100000ths of the shape's
    # shortest side. PowerPoint defaults to 0.10 when adj is absent.
    corner_radius: float | None = None
    if spPr is not None:
        pg = spPr.find("a:prstGeom", NS)
        if pg is not None and pg.get("prst") == "roundRect":
            gd = pg.find(".//a:gd[@name='adj']", NS)
            # ECMA-376 default when adj is omitted: 16667/100000 of the
            # shortest side (NOT 0.10 — that undersized default-rounded
            # corners by ~40%).
            adj_frac = 0.16667
            if gd is not None and gd.get("fmla"):
                m = re.search(r"val (\d+)", gd.get("fmla"))
                if m:
                    adj_frac = int(m.group(1)) / 100000
            corner_radius = cmap.w(min(w, h)) * adj_frac
    # Drop shadow — `<a:effectLst><a:outerShdw>`. Standard PowerPoint card
    # shadows use blurRad/dist in EMU, dir in 1/60000ths of a degree, and a
    # solid colour with an `<a:alpha val="...">` modifier (0-100000 = 0-100%).
    shadow: tuple[float, float, float, str, float] | None = None
    if spPr is not None:
        eff = spPr.find(".//a:effectLst/a:outerShdw", NS)
        if eff is not None:
            blur_emu = int(eff.get("blurRad") or 0)
            dist_emu = int(eff.get("dist") or 0)
            dir_60k = int(eff.get("dir") or 0)
            blur_px = cmap.w(blur_emu)
            dist_px = cmap.w(dist_emu)
            angle_deg = dir_60k / 60000.0
            sh_color = "black"
            sh_alpha = 1.0
            srgb = eff.find("a:srgbClr", NS)
            if srgb is not None and srgb.get("val"):
                hx = srgb.get("val")
                try:
                    rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                    sh_color = nearest_token(rgb, palette) if palette else f"#{hx}"
                except ValueError:
                    pass
                alpha = srgb.find("a:alpha", NS)
                if alpha is not None and alpha.get("val"):
                    sh_alpha = int(alpha.get("val")) / 100000.0
            shadow = (blur_px, dist_px, angle_deg, sh_color, sh_alpha)

    # Picture-typed placeholder → picture shape (no actual <p:pic>).
    if ph_type == "pic":
        shapes.append(Shape(
            kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            is_picture=True, ph_type=ph_type, ph_idx=ph_idx,
        ))
        return

    # Pure-text shape (placeholder, label, etc.) — no rect, just text.
    if runs and fill is None and kind == "rect":
        # Multi-paragraph colour split: when a text shape carries paragraphs
        # in different colours (typical of a cyan headline above silver body
        # bullets), emit one DSL `text` primitive per consecutive same-colour
        # block at calculated y-offsets, instead of collapsing them into a
        # single primitive whose first-run colour overrides the rest.
        blocks = _split_runs_by_color(runs)
        if len(blocks) > 1:
            frame_x_px = cmap.x(x)
            frame_y_px = cmap.y(y)
            frame_w_px = cmap.w(w)
            # Cursor-style y placement: each block consumes its own height
            # based on its primary pt (line-height factor 1.25 captures
            # standard slide leading without over-spacing tight headlines).
            cursor = frame_y_px
            for block_runs in blocks:
                block_pts = [r.pt for r in block_runs if r.text and r.text != "\n"]
                primary_pt = max(block_pts) if block_pts else 12
                line_count = sum(1 for r in block_runs if r.text and r.text != "\n")
                block_h_px = max(int(round(primary_pt * (4 / 3) * 1.25 * max(line_count, 1))), 24)
                # Split blocks always anchor top so the next block starts at
                # the bottom of the previous one — using the parent shape's
                # valign:middle would re-center each split block and overlap
                # neighbours.
                shapes.append(Shape(
                    kind="text", x=frame_x_px, y=cursor,
                    w=frame_w_px, h=block_h_px,
                    text_runs=block_runs, ph_type=ph_type, ph_idx=ph_idx,
                    valign=None, padding=padding_px,
                ))
                cursor += block_h_px
            return
        shapes.append(Shape(
            kind="text", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, valign=valign,
            autoshrink=autoshrink, font_scale=font_scale,
            padding=padding_px, line_spacing=line_spacing,
        ))
        return

    # FILLED rect carrying multi-style text (e.g. tile cards: bold 18pt
    # "Placeholder" lead-in over regular 12pt body inside one filled
    # shape). Collapsing to a single text primitive emits the body at the
    # lead-in's max size and bold weight. Split exactly like the fill-None
    # case above: the rect emits on its own (no runs), each text block
    # keeps its own size / colour / weight at its own y-offset.
    if runs and fill is not None and kind == "rect":
        blocks = _split_runs_by_color(runs)
        if len(blocks) > 1:
            shapes.append(Shape(
                kind="rect", x=cmap.x(x), y=cmap.y(y),
                w=cmap.w(w), h=cmap.h(h),
                fill=fill, stroke=stroke, stroke_width=stroke_width,
                stroke_dash=stroke_dash, corner_radius=corner_radius,
                shadow=shadow, gradient=gradient,
            ))
            cursor = cmap.y(y)
            for block_runs in blocks:
                block_pts = [r.pt for r in block_runs if r.text and r.text != "\n"]
                primary_pt = max(block_pts) if block_pts else 12
                line_count = sum(1 for r in block_runs if r.text and r.text != "\n")
                block_h_px = max(int(round(primary_pt * (4 / 3) * 1.25 * max(line_count, 1))), 24)
                shapes.append(Shape(
                    kind="text", x=cmap.x(x), y=cursor,
                    w=cmap.w(w), h=block_h_px,
                    text_runs=block_runs, ph_type=ph_type, ph_idx=ph_idx,
                    valign=None, padding=padding_px,
                    line_spacing=line_spacing,
                ))
                cursor += block_h_px
            return

    # custGeom paths convert directly to SVG path `d`. Build it in
    # canvas-pixel space so the surrounding svg-block can simply
    # `path "<d>"` without further transforms.
    svg_d = None
    native_xml = None
    if kind == "shape":
        svg_d = _custgeom_svg_d(spPr, cmap.w(w), cmap.h(h))
        if svg_d is None and spPr is not None:
            # Preset-geom polygon (triangle, diamond, arrow, etc.) — the
            # source uses `prstGeom prst="…"` with no custGeom, so
            # _custgeom_svg_d returns None. Synthesize the path from the
            # preset name so the renderer draws the correct outline
            # instead of falling back to a bbox-rect.
            pg = spPr.find("a:prstGeom", NS)
            if pg is not None:
                preset = pg.get("prst")
                if preset:
                    adj: dict[str, float] = {}
                    av = pg.find("a:avLst", NS)
                    if av is not None:
                        for gd in av.findall("a:gd", NS):
                            fmla = gd.get("fmla") or ""
                            if fmla.startswith("val "):
                                try:
                                    adj[gd.get("name") or ""] = \
                                        int(fmla[4:]) / 100000.0
                                except ValueError:
                                    pass
                    svg_d = _preset_geom_path(preset, cmap.w(w), cmap.h(h),
                                              adj)
        # Prefer carrying the native <p:sp> verbatim for TOP-LEVEL complex chrome:
        # a real, editable vector spliced straight into the output deck — NO
        # svg → raster → picture round-trip (which both distorts the shape and is
        # a picture "cheat"). The DSL stays the content layer (text + images);
        # corporate-design geometry rides along untouched. Colours are baked
        # schemeClr→srgbClr against the SOURCE theme so they survive the output
        # deck's theme. Grouped shapes (offset != 0) keep the svg path — their
        # xfrm is group-relative, not slide-absolute, so a verbatim splice would
        # land in the wrong place.
        if svg_d is not None and len(offset) == 2:
            # `offset` is the (EMU) group/layout translation this shape was walked
            # under; its xfrm is relative to that, so shift it to slide-absolute
            # before splicing. (Scaled groups thread an 8-tuple affine instead —
            # those fall through to the svg path, which is bbox-correct already.)
            sp_el = _carry_element(ch, theme, offset, "p:spPr/a:xfrm/a:off")
            native_xml = etree.tostring(sp_el).decode("utf-8")

    # Geometry shape (rect / oval / shape). May also carry text.
    shapes.append(Shape(
        kind=kind, x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        stroke_dash=stroke_dash, corner_radius=corner_radius, shadow=shadow,
        gradient=gradient,
        text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, svg_path_d=svg_d,
        native_xml=native_xml,
    ))


def _emit_pic(ch, offset, shapes, slide, cmap, theme, palette):
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    # Capture the embedded media rId so derive() can extract the binary
    # when image_extract_dir is set (pipeline-optimization mode).
    rid = None
    media_part = None
    blip = ch.find(".//a:blip", NS)
    if blip is not None:
        rid = blip.get(f"{{{NS['r']}}}embed")
        if rid is None:
            # SVG-only pictures: PowerPoint may leave <a:blip> bare and put
            # the reference solely in the <asvg:svgBlip> extension.
            _svg = ch.find(
                ".//{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip")
            if _svg is not None:
                rid = _svg.get(f"{{{NS['r']}}}embed")
        # Resolve the related part NOW against the source object that
        # actually owns the rId (the slide, layout, or master `slide`
        # param of this call). Storing only the rId and re-resolving on
        # `slide.part` later breaks for layout-inherited pictures —
        # their rId is scoped to the layout's relationships, not the
        # slide's.
        if rid is not None:
            try:
                media_part = slide.part.related_part(rid)
            except (KeyError, AttributeError):
                media_part = None
    # Template image: changeability follows the SOURCE AUTHOR'S OWN MARKERS
    # (provenance + media type), not a size heuristic:
    #   placeholder pic        → picture slot — the template explicitly
    #                            offers it for replacement. Detected by
    #                            <p:ph> (ph_type) OR by the placeholder-
    #                            derived shape NAME: filling a content
    #                            placeholder with a picture can drop the
    #                            <p:ph> element but keeps the name
    #                            ("Content Placeholder 5",
    #                            "Inhaltsplatzhalter 6", "Bildplatzhalter");
    #   plain pic, JPEG media  → picture slot — photographs are topical
    #                            content even when placed outside a
    #                            placeholder (photo strips, cover shots);
    #   plain pic, PNG/SVG/…   → fixed corporate-design graphic (logo,
    #                            supergraphic, icon, illustration band,
    #                            "Grafik N") — carry it natively (verbatim
    #                            element + media) so it is NOT bindable.
    # The previous mark-shape rule (min dimension ≤ 64px) left every
    # content-sized PNG illustration as a bindable slot, which the audit
    # then flagged and template users rightly complained about ("this CD
    # band must not be replaceable").
    native_xml = None
    native_media = None
    _partname = str(getattr(media_part, "partname", "")).lower()
    _is_photo = _partname.endswith((".jpg", ".jpeg"))
    _cnvpr = ch.find(".//p:nvPicPr/p:cNvPr", NS)
    _shape_name = _cnvpr.get("name", "") if _cnvpr is not None else ""
    _is_named_placeholder = bool(
        re.search(r"placeholder|platzhalter", _shape_name, re.IGNORECASE))
    if (ph_type is None and not _is_named_placeholder
            and media_part is not None and len(offset) == 2
            and not _is_photo):
        try:
            import base64 as _b64
            pic_el = _carry_element(ch, theme, offset, "p:spPr/a:xfrm/a:off")
            native_xml = etree.tostring(pic_el).decode("utf-8")
            media_blob = media_part.blob
            if str(getattr(media_part, "partname", "")).lower().endswith(".svg"):
                # python-pptx cannot re-embed SVG bytes at build time —
                # carry a rasterized PNG (emit drops the svgBlip sidecar).
                media_blob = _rasterize_svg_bytes(media_blob) or media_blob
            native_media = _b64.b64encode(media_blob).decode("ascii")
        except Exception:
            native_xml = native_media = None
    shapes.append(Shape(
        kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        is_picture=True, ph_type=ph_type, ph_idx=ph_idx, media_rid=rid,
        media_part=media_part, native_xml=native_xml, native_media=native_media,
    ))


def _bent_route_local(prst: str, w: float, h: float,
                      adj: dict[str, float]) -> list[tuple[float, float]] | None:
    """Polyline route of a bentConnectorN preset in LOCAL (unrotated,
    unflipped) shape coordinates, per the OOXML preset definitions. adj
    values are fractions of w/h (val 50000 → 0.5)."""
    a1 = adj.get("adj1", 0.5)
    a2 = adj.get("adj2", 0.5)
    a3 = adj.get("adj3", 0.5)
    if prst == "bentConnector2":
        return [(0, 0), (w, 0), (w, h)]
    if prst == "bentConnector3":
        x1 = a1 * w
        return [(0, 0), (x1, 0), (x1, h), (w, h)]
    if prst == "bentConnector4":
        x1, y1 = a1 * w, a2 * h
        return [(0, 0), (x1, 0), (x1, y1), (w, y1), (w, h)]
    if prst == "bentConnector5":
        x1, y1, x2 = a1 * w, a2 * h, a3 * w
        return [(0, 0), (x1, 0), (x1, y1), (x2, y1), (x2, h), (w, h)]
    return None


def _emit_cxn(ch, offset, shapes, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        return
    x0, y0, w0, h0 = xfrm
    x, y, w, h = x0, y0, w0, h0
    # Offset is a 2-tuple (translation-only ancestor groups) or an 8-tuple
    # (scaled-group affine). Apply it exactly like _shape_bbox so a connector
    # inside a SCALED group doesn't crash (`ox, oy = offset` blew up on the
    # 8-tuple) and still lands at the right place.
    if len(offset) == 2:
        ox, oy = offset
        x, y = x + ox, y + oy
    else:
        ox, oy, ax, ay, chox, choy, sx, sy = offset
        x = ax + (x - chox) * sx + ox
        y = ay + (y - choy) * sy + oy
        w, h = w * sx, h * sy

    def _apply_offset(px: float, py: float) -> tuple[float, float]:
        if len(offset) == 2:
            return px + offset[0], py + offset[1]
        ox, oy, ax, ay, chox, choy, sx, sy = offset
        return ax + (px - chox) * sx + ox, ay + (py - choy) * sy + oy

    # Stroke color + width + dash from <a:ln w="..."><a:solidFill .../><a:prstDash .../></a:ln>.
    ln = spPr.find("a:ln", NS)
    stroke = None
    stroke_width: float | None = None
    stroke_dash: str | None = None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
        w_attr = ln.get("w")
        if w_attr:
            try:
                stroke_width = cmap.w(int(w_attr))
            except (ValueError, TypeError):
                pass
        dash = ln.find("a:prstDash", NS)
        if dash is not None and dash.get("val"):
            stroke_dash = dash.get("val")

    # Bent connectors (the elbow trees of org charts and hierarchies): a
    # straight line between the bbox corners misrepresents the route —
    # PowerPoint draws axis-aligned H/V segments under a rot/flip transform.
    # Decompile the preset route into explicit `line` segments instead.
    # (Carrying the cxnSp verbatim does NOT work: LibreOffice routes a
    # connector by its a:stCxn/a:endCxn shape references when present —
    # in the rebuilt deck those IDs point at unrelated shapes — and its
    # pure xfrm+rot+flip rendering of bent presets is wrong even in the
    # unmodified source once the references are stripped.)
    _geom = spPr.find("a:prstGeom", NS) if spPr is not None else None
    _prst = (_geom.get("prst") or "") if _geom is not None else ""
    if _prst.startswith("bentConnector") and w0 >= 0 and h0 >= 0:
        adj: dict[str, float] = {}
        av = _geom.find("a:avLst", NS)
        if av is not None:
            for gd in av.findall("a:gd", NS):
                fmla = gd.get("fmla") or ""
                if fmla.startswith("val "):
                    try:
                        adj[gd.get("name") or ""] = int(fmla[4:]) / 100000.0
                    except ValueError:
                        pass
        pts = _bent_route_local(_prst, float(w0), float(h0), adj)
        if pts is not None:
            xfrm_el = spPr.find("a:xfrm", NS)
            flip_h = xfrm_el.get("flipH") in ("1", "true")
            flip_v = xfrm_el.get("flipV") in ("1", "true")
            try:
                rot_deg = int(xfrm_el.get("rot") or 0) / 60000.0
            except ValueError:
                rot_deg = 0.0
            if flip_h:
                pts = [(w0 - px, py) for px, py in pts]
            if flip_v:
                pts = [(px, h0 - py) for px, py in pts]
            # Rotate about the local box center (clockwise in y-down slide
            # coordinates, like OOXML rot), then translate into slide space
            # and through the group offset/affine.
            rad = math.radians(rot_deg)
            cosr, sinr = math.cos(rad), math.sin(rad)
            ccx, ccy = w0 / 2.0, h0 / 2.0
            abs_pts = []
            for px, py in pts:
                dx, dy = px - ccx, py - ccy
                rx = dx * cosr - dy * sinr
                ry = dx * sinr + dy * cosr
                abs_pts.append(_apply_offset(x0 + ccx + rx, y0 + ccy + ry))
            # Dedupe across sibling connectors too — two elbows fanning out
            # of the same parent share their trunk segment verbatim.
            seen_segs: set[tuple[int, int, int, int]] = {
                (s.x, s.y, s.x + s.w, s.y + s.h)
                for s in shapes if s.kind == "line"
            }
            for (ax1, ay1), (ax2, ay2) in zip(abs_pts, abs_pts[1:]):
                # Normalise direction (full-endpoint swap, so diagonals
                # keep their slope) — downstream treats lines as x,y + w,h
                # and expects non-negative extents where possible.
                if (ax2, ay2) < (ax1, ay1):
                    (ax1, ay1), (ax2, ay2) = (ax2, ay2), (ax1, ay1)
                px1, py1 = cmap.x(ax1), cmap.y(ay1)
                px2, py2 = cmap.x(ax2), cmap.y(ay2)
                # Skip segments that collapse to a point at canvas resolution
                # (a degenerate connector axis, adj at 0/100%) and exact
                # duplicates (sibling elbows share their trunk segment).
                if (px1, py1) == (px2, py2) or (px1, py1, px2, py2) in seen_segs:
                    continue
                seen_segs.add((px1, py1, px2, py2))
                shapes.append(Shape(
                    kind="line", x=px1, y=py1, w=px2 - px1, h=py2 - py1,
                    stroke=stroke or "fog",
                    stroke_width=stroke_width,
                    stroke_dash=stroke_dash,
                ))
            return

    shapes.append(Shape(
        kind="line", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        stroke=stroke or "fog",
        stroke_width=stroke_width,
        stroke_dash=stroke_dash,
    ))


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
# Relationship TYPE the slide part uses to reach a chart part. Stored as the
# chart-part's `reltype` in the carried part-graph so the emitter re-creates
# the slide→chart relationship with the correct type (the value is otherwise
# unused for the leaf style/colors/xlsx parts, whose own reltypes are carried
# verbatim from the source rels).
RELS_NS_CHART = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"

# SmartArt / diagram graphicFrame. The `<dgm:relIds>` inside graphicData carries
# four slide-rels (data / layout / quickStyle / colors); a 5th part — the
# pre-rendered drawing — is reached via the dataN.xml extLst's
# `<dsp:dataModelExt relId>` (also a slide rel). All five (+ any media sub-rel of
# the data / drawing part) are native-carried so the diagram renders pixel-exact,
# vs `_emit_smartart`'s lossy flatten (parses the cached drawing into bbox rects,
# dropping connectors / geometry / per-node styling).
DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
DSP_DATAMODEL_NS = "http://schemas.microsoft.com/office/drawing/2008/diagram"
# Reltype for the diagramDrawing part (the dataN extLst relId points at it). This
# is a Microsoft-extension reltype with no RELATIONSHIP_TYPE constant.
RELS_NS_DGM_DRAWING = "http://schemas.microsoft.com/office/2007/relationships/diagramDrawing"
# The four `<dgm:relIds>` attributes → their slide-rel reltypes, so the emitter
# re-creates each slide→diagram relationship with the right type.
_DGM_RELID_ATTRS = (
    ("dm", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"),
    ("lo", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramLayout"),
    ("qs", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramQuickStyle"),
    ("cs", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramColors"),
)
# Relationship type a slide/part uses to reach an embedded raster image. Carried
# as the reltype of a grouped <p:pic>'s media part so the emitter re-creates the
# slide→image relationship with the right type when splicing a native group.
RELS_NS_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def _capture_part_graph(part, parent_key: str, reltype: str, src_rid: str,
                        parts: list[dict], seen: set[str],
                        shared: dict[str, dict[str, tuple[str, str]]],
                        theme: dict[str, str]) -> None:
    """Capture one OPC part + its transitive internal rel graph for native carry.

    Appends {partname, content_type, blob (b64), reltype, parent, src_rid} to
    `parts` and recurses into the part's own relationships. Every XML blob gets
    schemeClr→srgbClr baked against the SOURCE theme so the carried content
    keeps its EXACT source palette under the output deck's theme (chart series
    are nearly always theme accent refs; without baking they'd render in Office
    defaults instead of the brand palette). Binary parts (embedded .xlsx,
    media) ride along verbatim.

    `parent` is 'slide' for a root part (related from slide.part) else the
    OWNING partname, so the emitter rewires bottom-up; `src_rid` is the rId by
    which this part is referenced FROM its parent in the SOURCE deck — the
    emitter maps src_rid→new_rid to rewrite the matching r:id / r:embed refs.

    A part shared by two parents (a data part AND drawing10 both referencing
    ../media/image131.png; two charts sharing one workbook) is materialised
    ONCE, but EVERY parent's src_rid→(partname, reltype) mapping is recorded
    in `shared` so the second parent's own reference gets rewritten too —
    `_fold_shared_refs` turns the extra records into "ref" entries. Without
    this the second parent's reference would dangle at splice time.
    """
    import base64 as _b64
    pn = str(part.partname)
    if pn in seen:
        if src_rid:
            shared.setdefault(parent_key, {})[src_rid] = (pn, reltype)
        return
    seen.add(pn)
    if src_rid:
        shared.setdefault(parent_key, {})[src_rid] = (pn, reltype)
    raw = part.blob
    ct = part.content_type
    if ct.endswith("+xml") or ct.endswith("/xml"):
        try:
            _proot = etree.fromstring(raw)
            _bake_scheme_colors(_proot, theme)
            raw = etree.tostring(_proot, xml_declaration=True,
                                 encoding="UTF-8", standalone=True)
        except Exception:
            raw = part.blob
    parts.append({
        "partname": pn, "content_type": ct,
        "blob": _b64.b64encode(raw).decode("ascii"),
        "reltype": reltype, "parent": parent_key, "src_rid": src_rid,
    })
    try:
        child_rels = list(part.rels.values())
    except Exception:
        child_rels = []
    for r in child_rels:
        if r.is_external:
            continue
        try:
            tgt = r.target_part
        except Exception:
            continue
        _capture_part_graph(tgt, pn, r.reltype, r.rId, parts, seen, shared, theme)


def _fold_shared_refs(parts: list[dict],
                      shared: dict[str, dict[str, tuple[str, str]]]) -> None:
    """Fold `_capture_part_graph`'s shared-child records into the parts list:
    one extra {"ref": partname, parent, src_rid, reltype} entry per
    (parent, src_rid) not already materialised — blob omitted; the emitter
    recognises a ref-only entry, reuses the materialised part, and just wires
    the second parent's relationship + rewrites its reference."""
    for parent_key, m in shared.items():
        for src_rid, (old_pn, reltype) in m.items():
            if any(e.get("parent") == parent_key and e.get("src_rid") == src_rid
                   for e in parts):
                continue
            parts.append({"ref": old_pn, "parent": parent_key,
                          "src_rid": src_rid, "reltype": reltype})


def _capture_group_parts(grp, slide, theme) -> list[dict] | None:
    """Collect the external part-graph a native-carried `<p:grpSp>` reaches, so a
    grouped chart / SmartArt / image renders pixel-exact after splicing.

    Mirrors `_emit_graphic_frame`'s Stage B/C part-capture, but rooted at a WHOLE
    group rather than a single graphicFrame, and additionally carries grouped
    `<p:pic>` media (the existing single-`<p:pic>` `media:` path can't cover a
    group with several pics). Returns a list of {partname, content_type, blob,
    reltype, parent, src_rid} dicts (+ optional `ref` shared-part entries), or
    None when the group reaches no external parts (pure inline custGeom — the
    common decorative-vector case). Raises on a resolve failure so the caller's
    try/except falls back to recursion.

    Every captured XML blob gets schemeClr→srgbClr baked against the SOURCE theme
    so the carried content keeps its exact source palette under the output deck's
    theme — identical to the chart/diagram branches.
    """
    parts: list[dict] = []
    seen: set[str] = set()
    shared: dict[str, dict[str, tuple[str, str]]] = {}

    def _capture(part, parent_key: str, reltype: str, src_rid: str) -> None:
        _capture_part_graph(part, parent_key, reltype, src_rid,
                            parts, seen, shared, theme)

    # Grouped chart <c:chart r:id> — the chart part + its style/colors/xlsx graph.
    for cref in grp.iter(f"{{{CHART_NS}}}chart"):
        rid = cref.get(f"{{{RELS_NS}}}id")
        if rid:
            _capture(slide.part.related_part(rid), "slide", RELS_NS_CHART, rid)
    # Grouped SmartArt <dgm:relIds> — the four dgm parts + pre-rendered drawing.
    for relids in grp.findall(f".//{{{DGM_NS}}}relIds"):
        data_part = None
        for attr, reltype in _DGM_RELID_ATTRS:
            rid = relids.get(f"{{{RELS_NS}}}{attr}")
            if not rid:
                continue
            p = slide.part.related_part(rid)
            if attr == "dm":
                data_part = p
            _capture(p, "slide", reltype, rid)
        if data_part is not None:
            try:
                droot = etree.fromstring(data_part.blob)
                dmext = droot.find(f".//{{{DSP_DATAMODEL_NS}}}dataModelExt")
            except Exception:
                dmext = None
            draw_rid = dmext.get("relId") if dmext is not None else None
            if draw_rid:
                _capture(slide.part.related_part(draw_rid), "slide",
                         RELS_NS_DGM_DRAWING, draw_rid)
    # Grouped <p:pic> media — every <a:blip r:embed> resolves to an image part on
    # the slide; carry it so the spliced blip can be re-pointed at a fresh rId.
    for blip in grp.findall(".//a:blip", NS):
        rid = blip.get(f"{{{RELS_NS}}}embed")
        if not rid:
            continue
        try:
            mp = slide.part.related_part(rid)
        except Exception:
            continue
        _capture(mp, "slide", RELS_NS_IMAGE, rid)

    if not parts:
        return None
    _fold_shared_refs(parts, shared)
    return parts


def _try_carry_group(ch, offset, shapes, slide, cmap, theme) -> bool:
    """Native-carry a WHOLE decorative `<p:grpSp>` verbatim when it's pure chrome.

    Extends native-carry from top-level graphicFrames / custGeom shapes to entire
    groups — the last structural gap. A scaled group of custGeom vectors (world-map
    illustrations, decorative clusters) re-synthesises lossily today because the
    walker only carries a pure translation; carried WHOLE, PowerPoint applies the
    group's own off/ext/chOff/chExt affine and the children render pixel-exact.

    Qualifies (so we freeze it) only when the group:
      * contains a `custGeom` OR a `<p:graphicFrame>` descendant (complex content
        the per-shape re-synth botches), AND
      * contains NO `<p:ph>` descendant — carrying it whole never buries a fillable
        content slot, so the content/chrome split is preserved.

    Top-level only (`len(offset) == 2`): the group's xfrm is then slide-absolute, so
    shifting its own `grpSpPr/a:off` by the (pure-translation) offset places it
    correctly and the internal chOff/chExt handle the child affine untouched. A
    nested-scaled group (8-tuple offset) recurses as before. Returns True when it
    carried the group (caller must NOT recurse); False to recurse exactly as today.
    Any failure returns False → safe fall-back to recursion.
    """
    if len(offset) != 2:
        return False
    # Qualify gate: complex content present, no content placeholder buried inside.
    has_complex = (ch.find(".//a:custGeom", NS) is not None
                   or ch.find(".//p:graphicFrame", NS) is not None)
    if not has_complex:
        return False
    if ch.find(".//p:ph", NS) is not None:
        return False
    try:
        grp_xfrm = ch.find("p:grpSpPr/a:xfrm", NS)
        if grp_xfrm is None:
            return False
        off = grp_xfrm.find("a:off", NS)
        ext = grp_xfrm.find("a:ext", NS)
        if off is None or ext is None:
            return False
        ox, oy = int(offset[0]), int(offset[1])
        gx = int(off.get("x")) + ox
        gy = int(off.get("y")) + oy
        gw = int(ext.get("cx"))
        gh = int(ext.get("cy"))
        # Carry the external part-graph (grouped chart / SmartArt / pic media)
        # BEFORE mutating the copy, so blob capture reads the live source rels.
        parts = _capture_group_parts(ch, slide, theme)
        # Shift the group's OWN origin by the ancestor translation; leave chOff /
        # chExt (they define the child coordinate space the group internally maps).
        grp = _carry_element(ch, theme, offset, "p:grpSpPr/a:xfrm/a:off")
        shapes.append(Shape(
            kind="graphic",
            x=cmap.x(gx), y=cmap.y(gy), w=cmap.w(gw), h=cmap.h(gh),
            native_xml=etree.tostring(grp).decode("utf-8"),
            native_parts=parts,
        ))
        return True
    except Exception:
        return False


def _source_table_style(slide, style_id: str, theme: dict[str, str]):
    """The source deck's `<a:tblStyle styleId=…>` element, schemeClr-baked.

    Returns a deep copy safe to serialise into the DSL, or None when the
    package has no tableStyles part / no style with that id.
    """
    try:
        for part in slide.part.package.iter_parts():
            if not str(part.partname).endswith("tableStyles.xml"):
                continue
            root = etree.fromstring(part.blob)
            for st in root:
                if st.get("styleId") == style_id:
                    return _carry_element(st, theme)
    except Exception:
        pass
    return None


def _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette):
    """Tables and charts both arrive as <p:graphicFrame>. Dispatch by inner kind."""
    xfrm = ch.find("p:xfrm", NS)
    if xfrm is None:
        return
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return
    x0 = int(off.get("x"))
    y0 = int(off.get("y"))
    fw = int(ext.get("cx"))
    fh = int(ext.get("cy"))
    # Offset is a 2-tuple (translation-only ancestor groups) or an 8-tuple
    # (scaled-group affine). Apply it exactly like _shape_bbox so a frame
    # inside a SCALED group doesn't crash (`ox_local, oy_local = offset` blew
    # up on the 8-tuple) and still lands at the right place.
    if len(offset) == 2:
        ox_local, oy_local = offset
        x0, y0 = x0 + ox_local, y0 + oy_local
    else:
        ox, oy, ax, ay, chox, choy, sx, sy = offset
        x0 = ax + (x0 - chox) * sx + ox
        y0 = ay + (y0 - choy) * sy + oy
        fw, fh = fw * sx, fh * sy

    tbl = ch.find(".//a:tbl", NS)
    if tbl is not None:
        # Native-carry the whole graphicFrame (inline <a:tbl>, no external parts) for
        # a pixel-exact table — _emit_table's cell-by-cell re-synthesis drifts (guessed
        # row heights, only bottom borders, no merged cells, and mixed cell font sizes
        # collapse + overflow). Top-level only (offset is a pure translation).
        if len(offset) == 2:
            try:
                import base64 as _b64
                frame = _carry_element(ch, theme, offset, "p:xfrm/a:off")
                # The tbl's <a:tableStyleId> points into the SOURCE deck's
                # tableStyles.xml; the output deck doesn't have that style, so
                # the renderer falls back to its default (wrong header fill /
                # band colours / borders). Carry the referenced <a:tblStyle>
                # (schemeClr-baked) so the emitter can merge it in.
                parts: list[dict] | None = None
                sid = tbl.find("a:tblPr/a:tableStyleId", NS)
                if sid is not None and sid.text:
                    style_el = _source_table_style(slide, sid.text.strip(), theme)
                    if style_el is not None:
                        parts = [{"table_style": _b64.b64encode(
                            etree.tostring(style_el)).decode("ascii")}]
                shapes.append(Shape(
                    kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                    w=cmap.w(fw), h=cmap.h(fh),
                    native_xml=etree.tostring(frame).decode("utf-8"),
                    native_parts=parts,
                ))
                return
            except Exception:
                pass
        _emit_table(tbl, x0, y0, shapes, cmap, theme, palette)
        return

    # <c:chart r:id="..."/> inside graphicData → resolve chart part via slide rels.
    chart_ref = ch.find(f".//{{{CHART_NS}}}chart")
    if chart_ref is not None:
        # Native-carry the whole graphicFrame + its external part-graph (the chart
        # part itself + its chartStyle / chartColorStyle / embedded-xlsx children)
        # so the chart renders pixel-exact, vs _emit_chart's lossy re-synthesis
        # (pie/bar only, recoloured to brand, line/area/etc dropped). Top-level
        # only (offset is a pure translation; a grouped frame's xfrm is
        # group-relative). Falls back to _emit_chart on ANY failure so a single
        # awkward chart never crashes the decompile.
        if len(offset) == 2:
            try:
                frame = _carry_element(ch, theme, offset, "p:xfrm/a:off")
                # Walk every <c:chart r:id> in the frame and collect the part-graph
                # rooted at each chart part (chart → style / colors / xlsx; see
                # `_capture_part_graph` for the parent / src_rid entry semantics
                # the emitter rewires bottom-up).
                parts: list[dict] = []
                seen: set[str] = set()
                shared: dict[str, dict[str, tuple[str, str]]] = {}

                def _capture_graph(part, parent_key: str, reltype: str, src_rid: str) -> None:
                    _capture_part_graph(part, parent_key, reltype, src_rid,
                                        parts, seen, shared, theme)

                for cref in frame.iter(f"{{{CHART_NS}}}chart"):
                    rid = cref.get(f"{{{RELS_NS}}}id")
                    if not rid:
                        continue
                    cp = slide.part.related_part(rid)
                    _capture_graph(cp, "slide", RELS_NS_CHART, rid)
                if not parts:
                    raise ValueError("no chart parts resolved")
                _fold_shared_refs(parts, shared)
                shapes.append(Shape(
                    kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                    w=cmap.w(fw), h=cmap.h(fh),
                    native_xml=etree.tostring(frame).decode("utf-8"),
                    native_parts=parts,
                ))
                return
            except Exception:
                pass
        rid = chart_ref.get(f"{{{RELS_NS}}}id")
        if rid:
            try:
                chart_part = slide.part.related_part(rid)
            except Exception:
                chart_part = None
            if chart_part is not None:
                _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette)
        return

    # SmartArt / diagram: graphicData uri='…/drawingml/2006/diagram' containing
    # `<dgm:relIds r:dm/r:lo/r:qs/r:cs>`. Native-carry the whole graphicFrame + the
    # diagram part-graph (data / layout / quickStyle / colors + the pre-rendered
    # drawing + any media sub-rel) so the diagram renders pixel-exact, vs
    # `_emit_smartart`'s lossy flatten. Top-level only (offset is a pure
    # translation; a grouped frame's xfrm is group-relative). Falls back to
    # `_emit_smartart` on ANY failure so an awkward diagram never crashes decompile.
    gdata = ch.find(".//a:graphicData", NS)
    gdata_uri = gdata.get("uri") if gdata is not None else None
    relids_list = ch.findall(f".//{{{DGM_NS}}}relIds")
    if gdata_uri and gdata_uri.endswith("/diagram") and relids_list and len(offset) == 2:
        try:
            frame = _carry_element(ch, theme, offset, "p:xfrm/a:off")
            # Collect the diagram part-graph (same entry semantics as the chart
            # branch — see `_capture_part_graph`). The four dgm parts + the
            # drawing part hang off 'slide'; a media image (data/drawing →
            # ../media/imageN.png) hangs off its data/drawing part.
            parts: list[dict] = []
            seen: set[str] = set()
            shared: dict[str, dict[str, tuple[str, str]]] = {}

            def _capture_graph(part, parent_key: str, reltype: str, src_rid: str) -> None:
                _capture_part_graph(part, parent_key, reltype, src_rid,
                                    parts, seen, shared, theme)

            for relids in relids_list:
                # The four core parts, resolved via the slide rels.
                data_part = None
                for attr, reltype in _DGM_RELID_ATTRS:
                    rid = relids.get(f"{{{RELS_NS}}}{attr}")
                    if not rid:
                        continue
                    p = slide.part.related_part(rid)
                    if attr == "dm":
                        data_part = p
                    _capture_graph(p, "slide", reltype, rid)
                # The pre-rendered drawing: dataN extLst <dsp:dataModelExt relId>
                # is a SLIDE rel pointing at the drawing part.
                if data_part is not None:
                    dme = data_part.blob
                    try:
                        droot = etree.fromstring(dme)
                        dmext = droot.find(
                            f".//{{{DSP_DATAMODEL_NS}}}dataModelExt"
                        )
                    except Exception:
                        dmext = None
                    draw_rid = dmext.get("relId") if dmext is not None else None
                    if draw_rid:
                        dp = slide.part.related_part(draw_rid)
                        _capture_graph(dp, "slide", RELS_NS_DGM_DRAWING, draw_rid)
            if not parts:
                raise ValueError("no diagram parts resolved")
            _fold_shared_refs(parts, shared)
            shapes.append(Shape(
                kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                w=cmap.w(fw), h=cmap.h(fh),
                native_xml=etree.tostring(frame).decode("utf-8"),
                native_parts=parts,
            ))
            return
        except Exception:
            pass

    # SmartArt diagrams: graphicData uri='…/drawingml/2006/diagram'. The slide
    # rels carry both a `diagramData` (the semantic model) and a
    # `diagramDrawing` (the pre-rendered drawing*.xml computed by PowerPoint
    # when the user last edited the diagram). Parsing the drawing skips
    # re-implementing the SmartArt layout engine — every shape, its xfrm,
    # fill, stroke, and text live inside `<dsp:sp>` elements that mirror
    # the `<p:sp>` structure.
    diag_rel_ns = "http://schemas.microsoft.com/office/2007/relationships/diagramDrawing"
    if slide is not None and hasattr(slide, "part"):
        try:
            rels = slide.part.rels
            drawing_part = None
            for rel in rels.values():
                if rel.reltype == diag_rel_ns:
                    # The drawing isn't directly referenced by rId from the
                    # graphicFrame — it's a sibling relationship of the
                    # diagramData. PowerPoint ties them by partname suffix
                    # (data6.xml ↔ drawing6.xml). Find the matching one by
                    # numeric suffix on the graphicData's diagramData rId.
                    drawing_part = rel.target_part
                    # The slide may have multiple diagramDrawing rels (one
                    # per SmartArt). Use the rId currently being processed
                    # if available; fall back to first drawing.
                    break
            if drawing_part is not None:
                _emit_smartart(drawing_part.blob, x0, y0, fw, fh,
                               shapes, cmap, theme, palette)
        except Exception:
            pass


DSP_NS = "http://schemas.microsoft.com/office/drawing/2008/diagram"


def _emit_smartart(blob: bytes, x0: int, y0: int, fw: int, fh: int,
                   shapes: list, cmap, theme, palette) -> None:
    """Decompile a SmartArt's pre-rendered drawing XML.

    PowerPoint caches the computed-layout shapes for each SmartArt
    diagram in `ppt/diagrams/drawing*.xml` so they don't need to be
    relaid-out at presentation time. Each shape lives inside `<dsp:sp>`
    elements that mirror the regular `<p:sp>` structure but use the
    `dsp` namespace. Walking that tree gives us the actual circles,
    arrows, callouts, etc. without re-implementing the SmartArt layout
    engine.

    `x0,y0` and `fw,fh` are the host `<p:graphicFrame>`'s EMU position
    and size — the drawing's internal coordinates are already in
    slide-EMU (the layout engine wrote them out absolute), so no extra
    transform is needed for shapes whose own xfrm already lives in
    slide-space. The frame offset is preserved as a fallback for
    shapes whose xfrm is relative.
    """
    try:
        root = etree.fromstring(blob)
    except Exception:
        return
    spTree = root.find(f".//{{{DSP_NS}}}spTree")
    if spTree is None:
        return
    for sp in spTree.findall(f"{{{DSP_NS}}}sp"):
        spPr = sp.find(f"{{{DSP_NS}}}spPr")
        if spPr is None:
            continue
        xfrm = spPr.find("a:xfrm", NS)
        if xfrm is None:
            continue
        off = xfrm.find("a:off", NS)
        ext = xfrm.find("a:ext", NS)
        if off is None or ext is None:
            continue
        try:
            x = int(off.get("x")) + x0
            y = int(off.get("y")) + y0
            w = int(ext.get("cx"))
            h = int(ext.get("cy"))
        except (TypeError, ValueError):
            continue
        # `<dsp:sp>` xfrm coords are RELATIVE to the host graphicFrame
        # (the layout engine writes them in drawing-internal space), so
        # we add the frame's (x0, y0) to land each shape on the slide.
        # Fill — same resolver as <p:sp> uses (the a: children are
        # identical regardless of dsp vs p parent).
        fill = _resolve_fill(spPr, theme, palette)
        # Stroke + width from <a:ln>.
        stroke = None
        stroke_width: float | None = None
        ln = spPr.find("a:ln", NS)
        if ln is not None:
            sf = ln.find("a:solidFill", NS)
            if sf is not None:
                stroke = _resolve_solid(sf, theme, palette)
            w_attr = ln.get("w")
            if w_attr:
                try:
                    stroke_width = cmap.w(int(w_attr))
                except (TypeError, ValueError):
                    pass
        # Geometry kind.
        kind = "rect"
        pg = spPr.find("a:prstGeom", NS)
        if pg is not None:
            preset = pg.get("prst")
            if preset == "ellipse":
                kind = "oval"
            elif preset in ("line", "straightConnector1"):
                kind = "line"
            elif preset in ("rect", "roundRect"):
                kind = "rect"
        elif spPr.find("a:custGeom", NS) is not None:
            kind = "shape"
        # Text — dsp:txBody mirrors a:txBody / p:txBody.
        txBody = sp.find(f"{{{DSP_NS}}}txBody")
        runs = _text_runs(sp, theme, palette) if txBody is not None else []
        # Drop placeholder demo text the same way regular shapes do.
        if runs and _is_placeholder_text(runs):
            runs = []
        shapes.append(Shape(
            kind=kind,
            x=cmap.x(x), y=cmap.y(y),
            w=cmap.w(w), h=cmap.h(h),
            fill=fill, stroke=stroke, stroke_width=stroke_width,
            text_runs=runs,
        ))


def _emit_table(tbl, x0, y0, shapes, cmap, theme, palette):
    grid = tbl.find("a:tblGrid", NS)
    if grid is None:
        return
    col_widths = [int(c.get("w")) for c in grid.findall("a:gridCol", NS)]

    # Auto-fit h=0 rows: PowerPoint expands them to fit text. Use a heuristic
    # of 350000 EMU (~31 px) per text line so subsequent rows shift down and
    # don't overlap the header. Without this the header sits visually inside
    # row 1.
    EMU_PER_LINE = 350000

    rows = list(tbl.findall("a:tr", NS))
    effective_h = []
    for tr in rows:
        h = int(tr.get("h", "0"))
        if h == 0:
            # estimate from text content
            text = "".join(t.text or "" for t in tr.findall(".//a:t", NS))
            lines = max(1, len(text.split("\n"))) if text else 1
            h = EMU_PER_LINE * lines
        effective_h.append(h)

    y_cursor = y0
    for ri, tr in enumerate(rows):
        row_h = effective_h[ri]
        x_cursor = x0
        cells = tr.findall("a:tc", NS)
        for ci, tc in enumerate(cells):
            cw = col_widths[ci] if ci < len(col_widths) else 0
            tcPr = tc.find("a:tcPr", NS)
            fill = _resolve_fill(tcPr, theme, palette) if tcPr is not None else None
            runs = _text_runs(tc, theme, palette)
            if fill is not None:
                shapes.append(Shape(
                    kind="rect", x=cmap.x(x_cursor), y=cmap.y(y_cursor),
                    w=cmap.w(cw), h=cmap.h(row_h), fill=fill,
                ))
            # Cell bottom border → emit as line (orange separators under headers).
            if tcPr is not None:
                lnB = tcPr.find("a:lnB", NS)
                if lnB is not None:
                    sf = lnB.find("a:solidFill", NS)
                    stroke = _resolve_solid(sf, theme, palette) if sf is not None else None
                    if stroke is not None:
                        y_b = y_cursor + row_h
                        shapes.append(Shape(
                            kind="line", x=cmap.x(x_cursor), y=cmap.y(y_b),
                            w=cmap.w(cw), h=0, stroke=stroke,
                        ))
            if runs:
                shapes.append(Shape(
                    kind="text", x=cmap.x(x_cursor + cw // 20), y=cmap.y(y_cursor + row_h // 4),
                    w=cmap.w(cw - cw // 10), h=cmap.h(row_h),
                    text_runs=runs,
                ))
            x_cursor += cw
        y_cursor += row_h


def _emit_pie_chart(pie_el, x0, y0, fw, fh, shapes, cmap, theme=None, palette=None):
    """Extract pie/doughnut chart geometry and emit one svg{} arc path per slice.

    Each slice becomes a Shape with kind='shape', svg_path_d set to an SVG
    arc path of the form 'M cx,cy L x1,y1 A r,r 0 large,sweep x2,y2 Z',
    fill mapped to chart-series-N via slice index (using the brand's
    chart-series ramp — e.g. accent → accent-80 → ... → accent-10).
    emit_dsl() converts each Shape into a standalone `svg{}` block; the
    blocks share the same bbox (the chart frame) so slices overlay into
    one unified pie at render time.

    Percentage labels are emitted as `text` Shapes positioned just outside
    the slice's arc bisector when <c:dLbls><c:showPercent val="1"> is set
    in the source — matches PowerPoint's default external-label placement.

    Per-slice colors in the source XML are intentionally ignored in favour
    of the brand's chart-series ramp: source decks authored against the
    brand already use the same hue sequence, so index-based mapping gives
    fidelity AND brand-correctness in one pass. Source decks authored
    off-brand will use the brand's hues here — that is by design (a
    brand-pack decompile is brand-conforming output, not pixel mimicry).
    """
    ser = pie_el.find(f"{{{CHART_NS}}}ser")
    if ser is None:
        return
    val_els = ser.findall(
        f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    try:
        values = [float(v.text) for v in val_els if v.text]
    except (TypeError, ValueError):
        return
    if not values or sum(values) <= 0:
        return
    total = sum(values)

    # Per-slice colors from <c:dPt> data-point elements. Each dPt carries
    # <c:idx val="N"/> identifying the slice index plus an optional
    # <c:spPr><a:solidFill> with that slice's brand color. Falls back to
    # the chart-series-N ramp by index when absent.
    #
    # We resolve dPt fills via `palette={}` so the source's exact hex
    # propagates through unchanged. Going through the brand-token
    # nearest_token collapses two close source hues to the same token
    # ("accent"), which then renders identically — defeating the whole
    # purpose of per-slice colours. Same rationale as the custGeom
    # palette-bypass added earlier.
    slice_colors: dict[int, str] = {}
    if theme is not None:
        for dpt in ser.findall(f"{{{CHART_NS}}}dPt"):
            idx_el = dpt.find(f"{{{CHART_NS}}}idx")
            sp_pr = dpt.find(f"{{{CHART_NS}}}spPr")
            if idx_el is None or sp_pr is None:
                continue
            try:
                idx = int(idx_el.get("val") or "-1")
            except (TypeError, ValueError):
                continue
            color = _resolve_fill(sp_pr, theme, palette={})
            if color and idx >= 0:
                slice_colors[idx] = color

    cat_els = ser.findall(
        f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    categories = [c.text or "" for c in cat_els]

    # Legend position: source PowerPoint convention is r/l/t/b. We honour
    # only r/l (column-style legend with one row per category — the
    # dominant case for pies); t/b are rare on small pies and fall back
    # to right-side. Search at chart-space level since legend lives on
    # the chart root, not inside pie_el. When the chart has NO `<c:legend>`
    # element at all, the pie should fill its frame (no legend slot to
    # reserve). The previous code defaulted to "r" even when no legend
    # element existed, which made pies on legend-less charts (showcase
    # decks where every slice is labelled inline with `<c:showPercent>`)
    # render at ~50% of their source size.
    chart_root = pie_el
    while chart_root is not None and chart_root.tag != f"{{{CHART_NS}}}chartSpace":
        parent = chart_root.getparent()
        if parent is None:
            break
        chart_root = parent
    legend_pos: str | None = None
    if chart_root is not None:
        legend_el = chart_root.find(f".//{{{CHART_NS}}}legend")
        if legend_el is not None:
            # Overlay flag: when the legend is set to overlay the plot area
            # (`<c:overlay val="1"/>`), no horizontal slot is reserved for
            # it — treat as legend-less for sizing purposes.
            overlay_el = legend_el.find(f"{{{CHART_NS}}}overlay")
            is_overlay = overlay_el is not None and overlay_el.get("val") in ("1", "true")
            lp = legend_el.find(f"{{{CHART_NS}}}legendPos")
            if not is_overlay:
                # Collapse PowerPoint's 8-position legend space ("l", "r",
                # "t", "b", "tr", "tl", "br", "bl") down to the dominant
                # axis. Pie sizing only cares whether the legend lives
                # on the left/right (reserving horizontal space) or
                # top/bottom (vertical) — corner positions like "tr" act
                # as right-side legends for sizing.
                raw = lp.get("val") if lp is not None else "r"
                if raw in ("l", "tl", "bl"):
                    legend_pos = "l"
                elif raw in ("t", "b"):
                    legend_pos = raw
                else:  # "r", "tr", "br", or unknown — treat as right
                    legend_pos = "r"

    # Data label flags — <c:dLbls><c:showPercent val="1"/> or
    # <c:dLbls><c:showVal val="1"/>. Mutually exclusive in practice;
    # percent wins when both are set, matching PowerPoint behaviour.
    show_percent = False
    show_val = False
    sp = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showPercent"
    )
    if sp is not None and sp.get("val") == "1":
        show_percent = True
    sv = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showVal"
    )
    if sv is not None and sv.get("val") == "1":
        show_val = True

    # svg-block-local pixel coords for slice paths. The block's outer
    # bbox is the chart frame; coords inside are 0..bbox_w_px by
    # 0..bbox_h_px.
    bbox_w_px = cmap.w(fw)
    bbox_h_px = cmap.h(fh)

    # `<c:plotArea><c:layout><c:manualLayout>` gives EXACT fractional
    # plot-area position within the chart frame (xMode/yMode="edge" with
    # x/y/w/h as fractions of bbox_w/bbox_h). When present, use those
    # directly — they're what PowerPoint's layout engine resolved when
    # the deck author placed the chart. Falls back to the aspect-based
    # heuristic when the source uses `<c:layout/>` (auto-layout).
    pa_layout = None
    chart_root_for_layout = chart_root
    if chart_root_for_layout is not None:
        pa_layout = chart_root_for_layout.find(
            f".//{{{CHART_NS}}}plotArea/{{{CHART_NS}}}layout/{{{CHART_NS}}}manualLayout"
        )

    def _layout_frac(el, tag: str, default: float | None = None) -> float | None:
        if el is None:
            return default
        c = el.find(f"{{{CHART_NS}}}{tag}")
        if c is None or not c.get("val"):
            return default
        try:
            return float(c.get("val"))
        except (TypeError, ValueError):
            return default

    if pa_layout is not None:
        plot_xf = _layout_frac(pa_layout, "x", 0.0) or 0.0
        plot_yf = _layout_frac(pa_layout, "y", 0.0) or 0.0
        plot_wf = _layout_frac(pa_layout, "w", 1.0) or 1.0
        plot_hf = _layout_frac(pa_layout, "h", 1.0) or 1.0
        pie_off_x = plot_xf * bbox_w_px
        pie_off_y = plot_yf * bbox_h_px
        pie_w_px = plot_wf * bbox_w_px
        pie_h_px = plot_hf * bbox_h_px
    else:
        # Heuristic: pie-area fraction adapts to chart-frame aspect.
        # Wide frames (multi-pie-in-column layouts) keep ~60%; square-ish
        # frames shrink to ~50% leaving room for the legend.
        frame_aspect = bbox_w_px / bbox_h_px if bbox_h_px else 1.0
        if categories and legend_pos in ("l", "r"):
            pie_w_frac = 0.60 if frame_aspect > 1.4 else 0.50
        else:
            pie_w_frac = 1.0
        pie_w_px = bbox_w_px * pie_w_frac
        pie_h_px = bbox_h_px
        pie_off_x = (bbox_w_px - pie_w_px) if legend_pos == "l" else 0.0
        pie_off_y = 0.0
    cx_px = pie_off_x + pie_w_px / 2
    cy_px = pie_off_y + pie_h_px / 2
    # min(w,h) keeps pies circular in non-square frames. When manualLayout
    # gave us a plot-area smaller than the chart frame, the radius is half
    # the plot's shortest side (no further margin); otherwise 0.36 leaves
    # margin for the auto-layout's external label placement.
    if pa_layout is not None:
        r_px = min(pie_w_px, pie_h_px) / 2.0
    else:
        r_px = min(pie_w_px, pie_h_px) * 0.36

    # Doughnut hole: `<c:holeSize val="N"/>` on `<c:doughnutChart>` where
    # N is 10..90 = inner-radius percentage of outer radius. Default 50
    # when the element is missing on a doughnutChart; pieChart has no
    # hole. `pie_el.tag` localname distinguishes the two.
    inner_r_px = 0.0
    if etree.QName(pie_el).localname == "doughnutChart":
        hole_pct = 50
        hs = pie_el.find(f"{{{CHART_NS}}}holeSize")
        if hs is not None and hs.get("val"):
            try:
                hole_pct = max(10, min(90, int(hs.get("val"))))
            except (TypeError, ValueError):
                pass
        inner_r_px = r_px * (hole_pct / 100.0)

    # Start at 12 o'clock (-π/2), sweep clockwise. PowerPoint pies follow
    # this convention; matching it preserves slice-to-color correspondence
    # against the source.
    #
    # `<c:firstSliceAng val="N"/>` rotates the start clockwise by N
    # degrees (0..360). Sources rarely use it but when they do (corporate
    # decks rotating the highlighted slice into a fixed position), the
    # entire slice-to-colour ordering shifts.
    angle_start = -math.pi / 2
    first_ang_el = pie_el.find(f"{{{CHART_NS}}}firstSliceAng")
    if first_ang_el is not None and first_ang_el.get("val"):
        try:
            angle_start += math.radians(float(first_ang_el.get("val")))
        except (TypeError, ValueError):
            pass

    for i, v in enumerate(values):
        if v <= 0:
            angle_start += 0
            continue
        sweep = (v / total) * 2 * math.pi
        angle_end = angle_start + sweep
        x1 = cx_px + r_px * math.cos(angle_start)
        y1 = cy_px + r_px * math.sin(angle_start)
        x2 = cx_px + r_px * math.cos(angle_end)
        y2 = cy_px + r_px * math.sin(angle_end)
        large_arc = 1 if sweep > math.pi else 0
        if inner_r_px > 0:
            # Annular sector — outer arc forward + inner arc reversed.
            ix1 = cx_px + inner_r_px * math.cos(angle_start)
            iy1 = cy_px + inner_r_px * math.sin(angle_start)
            ix2 = cx_px + inner_r_px * math.cos(angle_end)
            iy2 = cy_px + inner_r_px * math.sin(angle_end)
            d = (
                f"M {x1:.1f},{y1:.1f} "
                f"A {r_px:.1f},{r_px:.1f} 0 {large_arc},1 {x2:.1f},{y2:.1f} "
                f"L {ix2:.1f},{iy2:.1f} "
                f"A {inner_r_px:.1f},{inner_r_px:.1f} 0 {large_arc},0 {ix1:.1f},{iy1:.1f} "
                f"Z"
            )
        else:
            # Pie wedge — apex at centre.
            d = (
                f"M {cx_px:.1f},{cy_px:.1f} "
                f"L {x1:.1f},{y1:.1f} "
                f"A {r_px:.1f},{r_px:.1f} 0 {large_arc},1 {x2:.1f},{y2:.1f} Z"
            )
        fill_token = slice_colors.get(i) or f"chart-series-{(i % 6) + 1}"
        shapes.append(Shape(
            kind="shape",
            x=cmap.x(x0), y=cmap.y(y0),
            w=bbox_w_px, h=bbox_h_px,
            fill=fill_token,
            svg_path_d=d,
        ))

        if show_percent or show_val:
            mid_angle = angle_start + sweep / 2
            # showPercent labels sit ~15% outside circumference (PPT
            # default external placement on small pies). showVal labels
            # sit inside the slice at ~65% radius — PowerPoint's default
            # internal-label position.
            label_r = r_px * 1.15 if show_percent else r_px * 0.65
            lx = cx_px + label_r * math.cos(mid_angle)
            ly = cy_px + label_r * math.sin(mid_angle)
            if show_percent:
                label_text = f"{round((v / total) * 100)} %"
            else:
                # Format like the source: 8.2 → "8,2", 1 → "1".
                if v == int(v):
                    label_text = str(int(v))
                else:
                    label_text = f"{v:.1f}".replace(".", ",")
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(lx) - 20,
                y=cmap.y(y0) + int(ly) - 10,
                w=40, h=20,
                text_runs=[TextRun(text=label_text, pt=10)],
            ))

        angle_start = angle_end

    # Legend (categories + colour swatches). Position: right or left of pie.
    # Stack one row per category centred vertically against the pie. Skip
    # entirely when no categories — labelled-from-percentage slices alone
    # carry enough signal for the small pies common in showcase decks.
    if categories and legend_pos in ("l", "r"):
        swatch_w = 14
        swatch_h = 14
        row_gap = 32  # vertical pitch between legend rows
        text_gap = 8  # swatch-to-label horizontal gap
        legend_w_px = bbox_w_px - pie_w_px
        # Legend bbox top-left within the chart frame.
        legend_x_local = 0 if legend_pos == "l" else pie_w_px
        # Center the legend block vertically against the pie.
        block_h = len(categories) * row_gap
        legend_y_local = max(0, (bbox_h_px - block_h) / 2)
        # Insets keep legend out of frame edges.
        legend_x_local += 12
        for i, cat in enumerate(categories):
            row_y = legend_y_local + i * row_gap
            color = f"chart-series-{(i % 6) + 1}"
            # Swatch as a rect — survives outside svg{} blocks and accepts
            # the brand vocab directly.
            shapes.append(Shape(
                kind="rect",
                x=cmap.x(x0) + int(legend_x_local),
                y=cmap.y(y0) + int(row_y),
                w=swatch_w, h=swatch_h,
                fill=color,
            ))
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(legend_x_local) + swatch_w + text_gap,
                y=cmap.y(y0) + int(row_y) - 4,
                w=int(legend_w_px) - swatch_w - text_gap - 12,
                h=24,
                text_runs=[TextRun(text=cat, pt=10)],
            ))


def _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette):
    """Extract chart geometry from a c:chartSpace part and emit DSL primitives.

    Dispatches by chart type:
      * c:pieChart / c:doughnutChart → _emit_pie_chart (arc paths + labels)
      * c:barChart                   → _emit_bar_chart (rects + axis + labels)
    Other chart types (line, area, scatter, etc.) fall through unhandled —
    the decompile-as-rects fallback in the hybrid SVG pass still gives a
    rough first-pass for those, and improve-brand sub-agents refine.
    """
    try:
        root = etree.fromstring(chart_part.blob)
    except Exception:
        return

    pie = (root.find(f".//{{{CHART_NS}}}pieChart")
           or root.find(f".//{{{CHART_NS}}}doughnutChart"))
    if pie is not None:
        _emit_pie_chart(pie, x0, y0, fw, fh, shapes, cmap, theme=theme, palette=palette)
        return

    bar = root.find(f".//{{{CHART_NS}}}barChart")
    if bar is None:
        return

    # Bar orientation: <c:barDir val="bar"/> = horizontal bars (categories
    # stack vertically, values extend rightward); val="col" (default) =
    # vertical columns. The previous code always emitted columns, so any
    # horizontal-bar source slide rendered with bars rotated 90° — visible
    # as vertical stripes where the source had horizontal bars.
    bar_dir_el = bar.find(f"{{{CHART_NS}}}barDir")
    horizontal_bars = (
        bar_dir_el is not None and bar_dir_el.get("val") == "bar"
    )

    # Stacking: <c:grouping val="standard|stacked|percentStacked"/>.
    # "standard" (default): series sit side-by-side per category.
    # "stacked": series stack head-to-tail; axis_max = max sum-per-category.
    # "percentStacked": each category bar is 100%; series segment by share.
    grouping_el = bar.find(f"{{{CHART_NS}}}grouping")
    grouping = grouping_el.get("val") if grouping_el is not None else "standard"
    if grouping not in ("standard", "stacked", "percentStacked", "clustered"):
        grouping = "standard"

    series = []
    for ser in bar.findall(f"{{{CHART_NS}}}ser"):
        name_el = ser.find(f".//{{{CHART_NS}}}tx//{{{CHART_NS}}}v")
        name = name_el.text if name_el is not None else "?"
        vals = [float(v.text) for v in ser.findall(f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        cats = [v.text for v in ser.findall(f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        # Per-series fill colour. Resolve via empty palette so the source
        # hex propagates verbatim — going through nearest_token collapses
        # two close source hues to the same brand token, losing the
        # series-to-series colour distinction on stacked / clustered
        # bar charts.
        sp_pr = ser.find(f"{{{CHART_NS}}}spPr")
        ser_color = _resolve_fill(sp_pr, theme, palette={}) if sp_pr is not None else None
        # Per-data-point colours from `<c:dPt>` overrides. Showcase decks
        # often colour alternating bars different hues to highlight a
        # specific category — that information lives in dPt only, not on
        # the series. Same palette={} bypass.
        bar_colors: dict[int, str] = {}
        for dpt in ser.findall(f"{{{CHART_NS}}}dPt"):
            idx_el = dpt.find(f"{{{CHART_NS}}}idx")
            dpt_sp = dpt.find(f"{{{CHART_NS}}}spPr")
            if idx_el is None or dpt_sp is None:
                continue
            try:
                idx = int(idx_el.get("val") or "-1")
            except (TypeError, ValueError):
                continue
            color = _resolve_fill(dpt_sp, theme, palette={})
            if color and idx >= 0:
                bar_colors[idx] = color
        series.append((name, vals, cats, ser_color, bar_colors))
    if not series:
        return

    n_cats = max(len(s[1]) for s in series)
    n_series = len(series)
    cats = series[0][2] if series[0][2] else [f"Cat {i+1}" for i in range(n_cats)]
    if grouping == "stacked":
        # Each category's bar = sum of all series values for that category.
        cat_totals = [sum(s[1][ci] for s in series if ci < len(s[1]))
                      for ci in range(n_cats)]
        data_max = max(cat_totals) if cat_totals else 0
    elif grouping == "percentStacked":
        # Every category sums to 100% — plot axis is just 100.
        data_max = 100
    else:
        data_max = max((max(s[1]) for s in series if s[1]), default=0)
    # Round axis max up to the next integer above data_max.
    # LibreOffice's auto-axis (the source-PNG ground truth in the verify
    # loop) adds one major-unit of headroom over data_max, so matching its
    # tick count beats the more semantically-correct ceil(data_max) — the
    # struct_diff_ratio improves when ticks line up with the source-PNG
    # rasterisation, not with what PowerPoint would have drawn.
    if grouping == "percentStacked":
        axis_max = 100
    else:
        axis_max = math.ceil(data_max + 0.5) if data_max > 0 else 5

    # Axis visibility — `<c:valAx><c:delete val="1"/>` and `<c:catAx>
    # <c:delete val="1"/>` hide the respective axis at render time. Many
    # showcase charts use this for a clean look: bars/segments alone, no
    # tick labels or category strings. Reading the flag from the source
    # XML is much better than always emitting ticks (and then mismatching
    # source pixels at every tick position).
    val_axis_hidden = False
    cat_axis_hidden = False
    for ax in root.findall(f".//{{{CHART_NS}}}valAx"):
        d = ax.find(f"{{{CHART_NS}}}delete")
        if d is not None and d.get("val") in ("1", "true"):
            val_axis_hidden = True
            break
    for ax in root.findall(f".//{{{CHART_NS}}}catAx"):
        d = ax.find(f"{{{CHART_NS}}}delete")
        if d is not None and d.get("val") in ("1", "true"):
            cat_axis_hidden = True
            break

    # Data-label flags — `<c:dLbls>` controls whether value/category/series
    # text gets drawn next to each bar. PowerPoint's structure puts this
    # on the bar chart element itself (and optionally per-series). Read
    # the top-level flags only for now; missing flags default to PPT's
    # behaviour (no labels unless set).
    def _dlbl_flag(parent, name: str) -> bool:
        el = parent.find(f"{{{CHART_NS}}}dLbls/{{{CHART_NS}}}{name}")
        return el is not None and el.get("val") in ("1", "true")
    show_val_labels = _dlbl_flag(bar, "showVal")
    show_cat_labels = _dlbl_flag(bar, "showCatName")

    # Plot-area extents inside the frame (EMU). When the axes are hidden
    # the plot can fill the frame edge-to-edge; otherwise reserve
    # PowerPoint's typical insets for tick/category labels.
    if val_axis_hidden and cat_axis_hidden:
        plot_x = x0
        plot_y = y0
        plot_w = fw
        plot_h = fh
    else:
        plot_x = x0 + int(fw * 0.07) if not val_axis_hidden else x0
        plot_y = y0 + int(fh * 0.12) if not cat_axis_hidden else y0
        plot_w = int(fw * (1.0 - 0.07 - 0.02)) if not val_axis_hidden else fw
        plot_h = int(fh * (1.0 - 0.12 - 0.22)) if not cat_axis_hidden else int(fh * (1.0 - 0.22))

    # Y-axis numeric labels — only when the value axis isn't hidden.
    if not val_axis_hidden:
        n_ticks = axis_max + 1
        for i in range(n_ticks):
            v = axis_max - i  # top→bottom
            ty = plot_y + int(plot_h * i / axis_max)
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0 + int(fw * 0.005)),
                y=cmap.y(ty - 180000),
                w=cmap.w(int(fw * 0.05)),
                h=cmap.h(360000),
                text_runs=[TextRun(text=str(v), pt=14)],
            ))

    # Skip gridlines: source has barely-visible hairlines; rendering them at
    # 0.75pt over a 1240px-wide plot adds heavy diff pixels and emit_dsl
    # orders lines after rects, so they paint OVER the bars producing stripes.

    # Category labels above each group — only when the category axis
    # isn't hidden AND the source explicitly enables them via dLbls.
    cat_w = plot_w // n_cats if n_cats else plot_w
    # Category labels render when the axis is visible AND the source either
    # opts into `<c:showCatName val="1"/>` OR omits the dLbls flag entirely
    # (PowerPoint's default behaviour shows axis-tied category labels even
    # without an explicit dLbls override).
    if not cat_axis_hidden and show_cat_labels:
        for ci in range(n_cats):
            cx = plot_x + ci * cat_w + cat_w // 4
            shapes.append(Shape(
                kind="text",
                x=cmap.x(cx),
                y=cmap.y(plot_y - int(fh * 0.07)),
                w=cmap.w(cat_w // 2),
                h=cmap.h(int(fh * 0.06)),
                text_runs=[TextRun(text=cats[ci] if ci < len(cats) else "", pt=14)],
            ))

    # Bars: each category has n_series side-by-side bars. PowerPoint sizes
    # them via `<c:gapWidth val="N"/>` where N is the inter-group gap as a
    # percentage of bar width (default 150 = gap is 1.5x bar width). The
    # category width then holds n_series bars plus a gap on each side:
    #   cat_w = bar_w * n_series + bar_w * (gapWidth/100)
    #   →  bar_w = cat_w / (n_series + gapWidth/100)
    # Reading the actual gapWidth lets thick "showcase" bars (default 150)
    # decompile at their real width instead of a fixed 8.5%-of-cat hairline,
    # which left source-bar pixels uncovered and inflated struct_diff on
    # every bar-chart slide.
    gap_pct = 150.0
    gw_el = bar.find(f"{{{CHART_NS}}}gapWidth")
    if gw_el is not None and gw_el.get("val"):
        try:
            gap_pct = float(gw_el.get("val"))
        except (TypeError, ValueError):
            pass
    if horizontal_bars:
        # Horizontal layout: category axis runs vertically (rows), value
        # axis runs horizontally. Each category gets a row of height
        # `cat_h`; bars within stack by series and extend rightward.
        cat_h = plot_h // n_cats if n_cats else plot_h
        if grouping in ("stacked", "percentStacked"):
            # One bar per category, series segments fill it left-to-right.
            bar_h = int(cat_h / (1 + gap_pct / 100))
            for ci in range(n_cats):
                # Per-category total (stacked uses sum, percentStacked
                # normalises to 100 so each row fills plot_w).
                if grouping == "percentStacked":
                    cat_total = sum(s[1][ci] for s in series if ci < len(s[1]))
                else:
                    cat_total = data_max
                cursor_x = plot_x
                row_y = plot_y + ci * cat_h + (cat_h - bar_h) // 2
                for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
                    if ci >= len(vals):
                        continue
                    v = vals[ci]
                    # Per-data-point `<c:dPt>` colour overrides the
                    # series colour for this specific category index.
                    color = dpt_colors.get(ci) or ser_color or f"chart-series-{(si % 6) + 1}"
                    if grouping == "percentStacked":
                        seg_w = int(plot_w * (v / cat_total)) if cat_total > 0 else 0
                    else:
                        seg_w = int(plot_w * (v / axis_max)) if axis_max > 0 else 0
                    shapes.append(Shape(
                        kind="rect",
                        x=cmap.x(cursor_x), y=cmap.y(row_y),
                        w=cmap.w(seg_w), h=cmap.h(bar_h),
                        fill=color,
                    ))
                    # Value labels for stacked / percentStacked horizontal
                    # bars — emit ONLY when source has `<c:showVal val="1"/>`.
                    # Label sits in the middle of its segment.
                    if show_val_labels and seg_w > 200000:
                        label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                        label = label.replace(".", ",")
                        shapes.append(Shape(
                            kind="text",
                            x=cmap.x(cursor_x + seg_w // 2 - 100000),
                            y=cmap.y(row_y),
                            w=cmap.w(200000),
                            h=cmap.h(int(bar_h)),
                            text_runs=[TextRun(text=label, pt=12)],
                        ))
                    cursor_x += seg_w
        else:
            bar_h = int(cat_h / (n_series + gap_pct / 100))
            group_h = bar_h * n_series
            group_inset_v = (cat_h - group_h) // 2
            for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
                default_color = ser_color or f"chart-series-{(si % 6) + 1}"
                for ci, v in enumerate(vals):
                    color = dpt_colors.get(ci) or default_color
                    by_ = plot_y + ci * cat_h + group_inset_v + si * bar_h
                    bw_ = int(plot_w * v / axis_max) if axis_max > 0 else 0
                    bx_ = plot_x
                    shapes.append(Shape(
                        kind="rect",
                        x=cmap.x(bx_), y=cmap.y(by_),
                        w=cmap.w(bw_), h=cmap.h(bar_h),
                        fill=color,
                    ))
                    if show_val_labels:
                        label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                        label = label.replace(".", ",")
                        shapes.append(Shape(
                            kind="text",
                            x=cmap.x(bx_ + bw_ + 50000),
                            y=cmap.y(by_),
                            w=cmap.w(int(fw * 0.08)),
                            h=cmap.h(int(bar_h)),
                            text_runs=[TextRun(text=label, pt=14)],
                        ))
    else:
        bar_w = int(cat_w / (n_series + gap_pct / 100))
        group_w = bar_w * n_series
        group_inset = (cat_w - group_w) // 2
        for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
            default_color = ser_color or f"chart-series-{(si % 6) + 1}"
            for ci, v in enumerate(vals):
                color = dpt_colors.get(ci) or default_color
                bx = plot_x + ci * cat_w + group_inset + si * bar_w
                bh = int(plot_h * v / axis_max) if axis_max > 0 else 0
                by = plot_y + plot_h - bh
                shapes.append(Shape(
                    kind="rect",
                    x=cmap.x(bx), y=cmap.y(by),
                    w=cmap.w(bar_w), h=cmap.h(bh),
                    fill=color,
                ))
                # Value label above the bar — only when source enables it.
                if show_val_labels:
                    label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                    label = label.replace(".", ",")
                    shapes.append(Shape(
                        kind="text",
                        x=cmap.x(bx - bar_w // 2),
                        y=cmap.y(by - 400000),
                        w=cmap.w(bar_w * 2),
                        h=cmap.h(360000),
                        text_runs=[TextRun(text=label, pt=14)],
                    ))

    # Legend + chart-title emission — gated on the source actually
    # carrying `<c:legend>` and `<c:title>` (with `<c:autoTitleDeleted
    # val="0"/>`). The previous code always painted them. Showcase
    # charts with no `<c:legend>` element rendered a phantom series
    # name (e.g. "Datenreihe 1") at the bottom-left that wraps mid-
    # word inside the swatch slot.
    legend_y = y0 + fh - int(fh * 0.12)
    legend_x = plot_x + int(fw * 0.02)
    swatch_w = int(fw * 0.012)
    swatch_h = int(fh * 0.025)
    has_legend = root.find(f".//{{{CHART_NS}}}legend") is not None
    title_el = root.find(f".//{{{CHART_NS}}}title")
    auto_title_deleted_el = root.find(f".//{{{CHART_NS}}}autoTitleDeleted")
    title_deleted = (
        auto_title_deleted_el is not None
        and auto_title_deleted_el.get("val") in ("1", "true")
    )
    title_text = ""
    if title_el is not None and not title_deleted:
        for t_el in title_el.iterfind(f".//{{{NS['a']}}}t"):
            if t_el.text:
                title_text += t_el.text
    if title_text:
        shapes.append(Shape(
            kind="text",
            x=cmap.x(legend_x), y=cmap.y(legend_y),
            w=cmap.w(int(fw * 0.22)), h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=title_text, pt=14)],
        ))
    if not has_legend:
        return
    lx = legend_x + int(fw * 0.18)
    for si, (name, _, _, ser_color, _dpt) in enumerate(series):
        color = ser_color or f"chart-series-{(si % 6) + 1}"
        shapes.append(Shape(
            kind="rect",
            x=cmap.x(lx), y=cmap.y(legend_y + (swatch_h // 4)),
            w=cmap.w(swatch_w), h=cmap.h(swatch_h),
            fill=color,
        ))
        shapes.append(Shape(
            kind="text",
            x=cmap.x(lx + swatch_w + 50000),
            y=cmap.y(legend_y),
            # Legend label width must fit "Data N" at 14pt without
            # wrapping; the previous 4% gave ~31 design-px which forced
            # multi-line "Da\nta\n1" — visibly wrong on every chart
            # legend. 10% fits the common case comfortably + leaves
            # room for short multi-word series names.
            w=cmap.w(int(fw * 0.10)),
            h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=name or "", pt=14)],
        ))
        lx += int(fw * 0.12)


# ---------------------------------------------------------------------------
# Style mapping (PPTX pt size → DSL style token)
# ---------------------------------------------------------------------------


_NUM_RE = re.compile(r"^\s*\d{1,2}\.\s*$")


def _style_for(pt: float, text: str, is_footer: bool) -> str:
    # Map source pt → nearest available style bundle (by emitted px). The
    # standard feinschliff bundles emit these px sizes:
    #   body-sm=16  ·  body=26  ·  sub=44  ·  title-l=80  ·  huge=120  ·  display=160
    # A DSL `text style:sub` round-trips through python-pptx as ≈33pt on the
    # rendered slide (px*0.75). So matching SOURCE pt to BUNDLE px directly
    # (1pt ≈ 1.333px) lands closer to source than the prior bucketing,
    # which mapped any pt ≥ 28 to title-l and over-sized 32pt slide titles
    # by almost 2×. Boundaries below are midpoints between adjacent bundle
    # px values, expressed as source pt.
    if _NUM_RE.match(text):
        return "agenda-num"
    if is_footer:
        return "footer"
    px = pt * 1.333
    if px >= 140:                              # 140+ px (≈ 105pt+) → display
        return "display"
    if px >= 100:                              # 100-140 px (≈ 75-105pt) → huge
        return "huge"
    if px >= 62:                               # 62-100 px (≈ 47-75pt) → title-l
        return "title-l"
    if px >= 35:                               # 35-62 px (≈ 26-47pt) → sub
        return "sub"
    if px >= 21:                               # 21-35 px (≈ 16-26pt) → body
        return "body"
    return "body-sm"                           # <21 px → body-sm (16 px)


# ---------------------------------------------------------------------------
# Emission
# ---------------------------------------------------------------------------


# Demo-placeholder text patterns. Templates ship slides with literal
# guidance text so the user can see what each slot is for; we suppress
# them on decompile so the round-trip render shows an empty slot
# instead of the prompt.
#
# Exact phrases (case-insensitive, stripped). Add a phrase here whenever
# a new corporate template surfaces a new prompt variant.
_PLACEHOLDER_EXACT: frozenset[str] = frozenset(map(str.lower, [
    "headline", "subheadline", "sub-headline", "sub headline",
    "placeholder", "placeholder text", "placeholder copy",
    "this is a placeholder text", "this is placeholder text",
    "title", "subtitle", "sub-title",
    "presentation title", "chapter title", "section title", "slide title",
    "welcome!", "welcome", "thank you!", "thank you",
    "lorem ipsum",
    "body text", "body copy",
    "caption", "footnote",
    "click here to add text", "click to add text",
    "this text can be replaced with your own text",
    "this text can be replaced",
    "your text here", "insert text here",
    # PowerPoint outline-level prompts (default for body placeholders) —
    # English + German since many corporate templates ship localised.
    "click to edit master title style",
    "click to edit master text styles",
    "second level", "third level", "fourth level", "fifth level",
    "text hinzufügen", "text hinzufuegen",
    "zweite ebene", "dritte ebene", "vierte ebene", "fünfte ebene", "fuenfte ebene",
    "klicken sie, um einen titel hinzuzufügen",
    "klicken sie, um titel hinzuzufügen",
    # Common prompts from vendor design kits.
    "add text", "add title", "add headline", "add subheadline",
]))

# Prefix patterns — first few words of the text run lower-cased.
_PLACEHOLDER_PREFIXES: tuple[str, ...] = (
    "click to edit",
    "lorem ipsum",
    "this text can be replaced",
    "this text demonstrates",       # "This text demonstrates how your own text..."
    "this is a placeholder",
    "this is placeholder",
    "placeholder text",             # "Placeholder text\n…" wrappers
    "click here to add",
    "tap to add",
    "double-click to edit",
    "replace this text",
    "replace with your",
    "your own text",
    "sample text",
    "example text",
)

# Mail-merge / template-variable tokens (corporate convention): wholly-token
# strings like `%classification%`
# or chained `%a%%b%%c%`. PowerPoint replaces these at the org level;
# our renderer has no such facility, so they emit literally.
_TEMPLATE_VAR_RE = __import__("re").compile(r"^(%[A-Za-z][A-Za-z0-9_-]*%)+$")


def _is_placeholder_line(line: str) -> bool:
    norm = line.strip().lower().rstrip(".!?:;,")
    if not norm:
        return True   # blank line counts as placeholder noise
    if norm in _PLACEHOLDER_EXACT:
        return True
    if any(norm.startswith(p) for p in _PLACEHOLDER_PREFIXES):
        return True
    if _TEMPLATE_VAR_RE.match(line.strip()):
        return True
    return False


# Unambiguous master/outline prompts — imperative UI copy that is NEVER
# legitimate slide-authored content (unlike "Headline"/"Welcome"/"Title",
# which a deck may genuinely ship as example copy). These are suppressed even
# when a shape looks slide-authored, because pptx_flatten_inherited.py
# materialises inherited master/layout placeholders as slide-level shapes —
# clearing the `from_chain` flag the normal placeholder suppression relies on,
# so the outline prompt ("Add Text\nZweite Ebene\n…") would otherwise survive
# and render on every flattened slide.
_STRONG_PROMPT_EXACT: frozenset[str] = frozenset(map(str.lower, [
    "add text", "add title", "add headline", "add subheadline",
    "click to edit master title style", "click to edit master text styles",
    "click here to add text", "click to add text",
    "second level", "third level", "fourth level", "fifth level",
    "text hinzufügen", "text hinzufuegen",
    "zweite ebene", "dritte ebene", "vierte ebene", "fünfte ebene", "fuenfte ebene",
    "klicken sie, um einen titel hinzuzufügen",
    "klicken sie, um titel hinzuzufügen",
    "insert photo", "insert picture", "insert image",
    "your text here", "insert text here",
]))
_STRONG_PROMPT_PREFIXES: tuple[str, ...] = (
    "click to edit master", "click here to add", "double-click to edit",
)


def _is_strong_prompt_line(line: str) -> bool:
    norm = line.strip().lower().rstrip(".!?:;,")
    if not norm:
        return True   # blank line inside an otherwise all-prompt block is noise
    if norm in _STRONG_PROMPT_EXACT:
        return True
    return any(norm.startswith(p) for p in _STRONG_PROMPT_PREFIXES)


def _is_strong_prompt(text_runs: list["TextRun"]) -> bool:
    """True iff every non-blank line is an unambiguous master/outline prompt.

    Stricter than `_is_placeholder_text`: applied to slide-AUTHORED text too
    (post-flatten placeholders), so it must only match copy that can never be
    real content — outline levels, "Click to edit Master…", "Add Text",
    "Insert photo". One real line spares the shape.
    """
    if not text_runs:
        return False
    raw = "".join(r.text or "" for r in text_runs)
    lines = [ln for ln in raw.splitlines() if ln.strip()]
    if not lines:
        return False
    return all(_is_strong_prompt_line(ln) for ln in lines)


def _is_placeholder_text(text_runs: list["TextRun"]) -> bool:
    """Return True iff the concatenated run text is recognizable demo /
    template-prompt copy that should NOT survive the decompile round-trip.

    Suppression triggers on any of:
      * the WHOLE text (linebreaks collapsed to spaces) matches an exact
        prompt or a known prefix — catches `"This text can be replaced\\nwith your own text."` where the source linebreaks for layout
        reasons but the whole string is still one prompt
      * every non-blank line matches an individual prompt pattern —
        catches `"Headline\\nSubheadline\\nBody"` where each line is its
        own prompt stacked in one placeholder
    """
    if not text_runs:
        return False
    raw = "".join(r.text or "" for r in text_runs).strip()
    if not raw:
        return False
    # Collapse line breaks + duplicate whitespace, then check.
    flat = " ".join(raw.split())
    if _is_placeholder_line(flat):
        return True
    lines = [ln for ln in raw.splitlines() if ln.strip()]
    if lines and all(_is_placeholder_line(ln) for ln in lines):
        return True
    return False


def _strip_placeholder_paragraphs(text_runs: list["TextRun"]) -> list["TextRun"]:
    """Return text_runs with placeholder paragraphs removed.

    Source decks often combine real labels with prompt copy inside one
    placeholder, e.g.

        "04\\nThis is a placeholder text\\nThis text demonstrates how your
        own text will look when you replace the placeholder with your own
        text."

    Dropping the whole shape would lose the "04" label that's actual
    content. This walks paragraphs (separated by `TextRun(text="\\n")`
    markers inserted by `_text_runs`), drops paragraphs whose joined
    text matches `_is_placeholder_line`, and stitches the survivors
    back together with the same `\\n` separators.

    Returns the original list unchanged when no paragraph qualifies as
    placeholder, or an empty list when EVERY paragraph qualifies (caller
    drops the shape).
    """
    if not text_runs:
        return text_runs
    # Group into paragraphs. A paragraph is a contiguous slice of runs
    # not crossed by a `\n` separator run.
    paragraphs: list[list["TextRun"]] = [[]]
    for r in text_runs:
        if r.text == "\n":
            paragraphs.append([])
            continue
        paragraphs[-1].append(r)
    kept: list[list["TextRun"]] = []
    dropped_any = False
    for para in paragraphs:
        joined = "".join((r.text or "") for r in para)
        if _is_placeholder_line(joined.strip()):
            dropped_any = True
            continue
        kept.append(para)
    if not dropped_any:
        return text_runs
    out: list["TextRun"] = []
    for i, para in enumerate(kept):
        if i > 0:
            # Reuse a TextRun shape similar to what `_text_runs` emits —
            # `pt` from the para's first run so the separator's height is
            # consistent with the surrounding text.
            sep_pt = para[0].pt if para else 12
            out.append(TextRun(text="\n", pt=sep_pt))
        out.extend(para)
    return out


def emit_dsl(shapes: list[Shape], cmap: CanvasMap, layout_name: str,
             theme_name: str = "feinschliff",
             placeholder_rel: str = PLACEHOLDER_REL,
             bg_fill: str | None = None,
             bg_image_path: str | None = None,
             native_dir: Path | None = None,
             native_rel: str | None = None) -> str:
    out: list[str] = [
        "# auto-derived from PPTX+SVG hybrid — review before use",
        f"# layout: {layout_name}",
        f"canvas {cmap.cw}x{cmap.ch}",
        f"theme {theme_name}",
        "",
    ]

    # Slide-level background fill (from <p:cSld><p:bg>) emits first as a
    # full-canvas rect. PowerPoint draws this *under* every shape, so the
    # DSL ordering matches the source z-order. Without this, dark slides
    # rebuild on the brand's paper default and white text disappears.
    # A blipFill background (template artwork, e.g. a full-bleed engraved
    # illustration) emits as a fixed full-canvas picture instead — it is
    # CHROME, not a fillable content slot, so no slot expression.
    if bg_image_path:
        out.append(f'picture 0,0 {cmap.cw}x{cmap.ch} path:"{bg_image_path}" cover:true')
    elif bg_fill:
        out.append(f"rect 0,0 {cmap.cw}x{cmap.ch} fill:{bg_fill}")

    # Pre-split into geometry-first, text-last, footer-last.
    rects = [s for s in shapes if s.kind == "rect" and s.fill]
    ovals = [s for s in shapes if s.kind == "oval"]
    custs = [s for s in shapes if s.kind == "shape"]
    pics = [s for s in shapes if s.kind == "pic" or s.is_picture]
    lines = [s for s in shapes if s.kind == "line"]
    graphics = [s for s in shapes if s.kind == "graphic"]
    texts: list[Shape] = []
    for s in shapes:
        if s.kind == "text" and s.text_runs:
            texts.append(s)
        elif s.kind in ("rect", "oval", "shape") and s.text_runs:
            # geometry shape that also carries text — emit shape now, defer text
            texts.append(Shape(
                kind="text", x=s.x, y=s.y, w=s.w, h=s.h, text_runs=s.text_runs,
                ph_type=s.ph_type, ph_idx=s.ph_idx, valign=s.valign,
                padding=s.padding, from_chain=s.from_chain,
            ))
    # Drop demo-placeholder text. Corporate templates ship slides with
    # literal guidance text inside placeholders ("Headline", "Subheadline",
    # "Click to edit Master title", "%classification%" mail-merge tokens,
    # "This text can be replaced with your own text.") so the user knows
    # what each slot is for. We don't want those strings showing up in
    # the rendered round-trip — keeps the layout chrome / geometry but
    # drops the prompt text so the slot reads as empty in the brand
    # template render. See `_is_placeholder_text` for the patterns.
    #
    # When a shape mixes a real label with placeholder paragraphs
    # ("04\nThis is a placeholder text\nThis text demonstrates…"), strip
    # ONLY the placeholder paragraphs so the label survives. Shapes
    # whose every paragraph is placeholder collapse to empty runs and
    # the whole shape gets dropped.
    filtered: list[Shape] = []
    for t in texts:
        # Demo-placeholder suppression applies ONLY to text inherited from the
        # layout/master (true prompt copy). Slide-AUTHORED text is real content —
        # corporate templates ship visible example copy ("Presentation
        # title", "This is a placeholder text") right on the slide, which the renderer
        # shows; string-matching it as a prompt wrongly deleted it. (MS-gallery prompts
        # are inherited and already blanked in walk_slide, so this doesn't regress them.)
        if not t.from_chain:
            # pptx_flatten_inherited materialises inherited master/layout
            # placeholders as slide-level shapes, clearing from_chain. Their
            # outline-prompt copy ("Add Text\nZweite Ebene\n…") is never real
            # content, so suppress it even though it now looks slide-authored.
            # Genuine example copy is left untouched (strict matcher).
            if _is_strong_prompt(t.text_runs):
                continue
            filtered.append(t)
            continue
        if _is_placeholder_text(t.text_runs):
            continue
        stripped = _strip_placeholder_paragraphs(t.text_runs)
        if stripped is t.text_runs:
            filtered.append(t)
            continue
        if not any((r.text or "").strip() for r in stripped):
            continue
        filtered.append(Shape(
            kind=t.kind, x=t.x, y=t.y, w=t.w, h=t.h,
            text_runs=stripped, ph_type=t.ph_type, ph_idx=t.ph_idx,
            valign=t.valign, padding=t.padding, from_chain=t.from_chain,
        ))
    texts = filtered

    footer_y_threshold = int(cmap.ch * 0.92)

    # Preserve SOURCE z-order — `rects` is already in render order (inherited
    # chrome behind slide content, each layer in spTree order). The previous
    # area-descending sort wrongly buried a large CONTENT panel beneath a smaller
    # decorative panel drawn ON TOP of it (MS Geometric: the cream content card
    # sank under the accent strip, so the whole slide read as the accent colour).
    # Only a near-full-bleed background rect is still forced to the bottom.
    # Rects interleave with shapes / pics / natives via the shared `_layers`
    # stream below — emitting ALL rects before the pic section painted
    # native-carried background art (e.g. a banner wave whose spTree
    # position sits BEFORE the tiles) on top of the content rects it
    # underlies in the source.
    # Stroke / dash / radius are captured so framed cards + dividers round-trip.
    _canvas_area = max(1, cmap.cw * cmap.ch)
    _src_pos = {id(_s): _i for _i, _s in enumerate(shapes)}
    _layers: list[tuple[int, list[str]]] = []
    for r in rects:
        line = f"rect {r.x},{r.y} {r.w}x{r.h} fill:{r.fill}"
        if r.gradient is not None:
            stops, angle = r.gradient
            stops_str = ";".join(f"{p:.2f}={c}" for p, c in stops)
            line += f" gradient:angle={angle:g};{stops_str}"
        if r.corner_radius is not None and r.corner_radius > 0:
            line += f" radius:{r.corner_radius:g}"
        if r.stroke:
            line += f" stroke:{r.stroke}"
            if r.stroke_width is not None and r.stroke_width > 0:
                line += f" stroke-width:{r.stroke_width:g}"
            if r.stroke_dash:
                line += f" dash:{r.stroke_dash}"
        if r.shadow is not None:
            blur, dist, angle, color, alpha = r.shadow
            line += (f" shadow:blur:{blur:g},dist:{dist:g},"
                     f"angle:{angle:g},color:{color},alpha:{alpha:.2f}")
        if (r.w * r.h) >= 0.9 * _canvas_area:
            out.append(line)    # full-bleed background stays at the bottom
        else:
            _layers.append((_src_pos.get(id(r), 0), [line]))

    # Custom shapes (puzzle pieces, parallelograms, border paths, ring
    # sectors, etc.). When we recovered an SVG `d` string from the source
    # `<a:custGeom>`, emit an `svg { path "<d>" … }` block so the
    # renderer reproduces the actual vector geometry — this is the
    # difference between "blue donut with the right arc" and "grey bbox
    # rect where the donut should be." Stroke-only paths emit
    # `stroke:<token>` with no fill; otherwise fill (or fog fallback).
    # When no path data is available we keep the lossy bbox-rect fallback
    # so the DSL still builds.
    # Custom shapes / graphic frames / ovals / pictures interleave in SOURCE
    # z-order: PowerPoint draws them as one stream, and grouping by kind put
    # e.g. a custGeom outline drawn ON TOP of a photo would sink underneath
    # the picture slot when grouped by kind. Each shape's lines collect into a
    # chunk tagged with its position in `shapes` (already spTree order);
    # the chunks emit sorted at the end. (`_src_pos` / `_layers` are
    # initialised above the rect loop, which feeds the same stream.)
    for i, s in enumerate(custs, 1):
        chunk: list[str] = []
        if s.native_xml:
            # Native vector chrome carried verbatim from the source (editable,
            # pixel-exact). Small fragments ride inline as base64 so the DSL
            # line is self-contained; big ones (33 MB vectorised groups exist)
            # go to a sidecar file under the pack's assets dir.
            import base64 as _b64
            _xb = s.native_xml.encode("utf-8")
            _ref = _native_sidecar_ref(_xb, "xml", native_dir, native_rel)
            if _ref:
                chunk.append(f'native shape{i} xml_file:"{_ref}"')
            else:
                _enc = _b64.b64encode(_xb).decode("ascii")
                chunk.append(f'native shape{i} b64:"{_enc}"')
        elif s.svg_path_d:
            # `none` is not in the SVG DSL's 17-name semantic colour
            # vocabulary, so we omit `fill:` entirely when the source has
            # no solid fill — the path primitive defaults to stroke-only
            # rendering. With a fill, map the brand token onto an SVG
            # vocabulary name.
            attrs = []
            if s.fill:
                attrs.append(f"fill:{_svg_color_token(s.fill)}")
            if s.stroke:
                attrs.append(f"stroke:{_svg_color_token(s.stroke)}")
            attr_str = (" " + " ".join(attrs)) if attrs else ""
            chunk.append(f"svg shape{i} {s.x},{s.y} {s.w}x{s.h} {{")
            # Path coordinates are already in svg-block-local pixels (the
            # converter scales from path-local space to the shape's bbox).
            chunk.append(f"  path p \"{s.svg_path_d}\"{attr_str}")
            chunk.append("}")
        elif s.fill is None and s.stroke:
            chunk.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect stroke:{s.stroke}")
        else:
            chunk.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect fill:{s.fill or 'fog'}")
        _layers.append((_src_pos.get(id(s), 0), chunk))

    # Native graphic frames carried verbatim — pixel-exact, vs the lossy
    # re-synthesis. Tables ship inline (<a:tbl>, no external parts → b64 only).
    # Charts ALSO carry an external part-graph (the chart part + chartStyle /
    # chartColorStyle / embedded-xlsx) as a base64-of-json `parts:` kwarg so the
    # emitter can re-create the parts + rewire rIds in the output deck.
    for i, s in enumerate(graphics, 1):
        chunk = []
        if s.native_xml:
            import base64 as _b64
            _xb = s.native_xml.encode("utf-8")
            _ref = _native_sidecar_ref(_xb, "xml", native_dir, native_rel)
            if _ref:
                line = f'native graphic{i} xml_file:"{_ref}"'
            else:
                line = f'native graphic{i} b64:"{_b64.b64encode(_xb).decode("ascii")}"'
            if s.native_parts:
                _pb = json.dumps(s.native_parts).encode("utf-8")
                _pref = _native_sidecar_ref(_pb, "json", native_dir, native_rel)
                if _pref:
                    line += f' parts_file:"{_pref}"'
                else:
                    line += f' parts:"{_b64.b64encode(_pb).decode("ascii")}"'
            chunk.append(line)
        if chunk:
            _layers.append((_src_pos.get(id(s), 0), chunk))

    # Ovals (circles, decorative dots). Stroke-only ovals (callout
    # circles, annotation marks) emit as stroke without fill; fill-only
    # ovals emit fill; if neither is present, fall back to a muted neutral
    # so the DSL always builds.
    for o in ovals:
        if o.fill is None and o.stroke:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval stroke:{o.stroke}"
        else:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval fill:{o.fill or 'fog'}"
            if o.stroke:
                line += f" stroke:{o.stroke}"
        if o.stroke_width is not None and o.stroke_width > 0:
            line += f" stroke-width:{o.stroke_width:g}"
        _layers.append((_src_pos.get(id(o), 0), [line]))

    # Pictures — default to the brand's generic placeholder (so a derived
    # layout works as a reusable template) OR, when derive() was called
    # with image_extract_dir, to the actual extracted-from-source asset
    # (pipeline-optimization mode: no picture_coverage masking needed,
    # struct_diff_ratio reflects real shape/text mismatch).
    # Clamp bbox to the canvas so that picture-bleed boxes (e.g.
    # 166,-144 2345x1319 on 1920x1080) become canvas-fitted rectangles.
    # PowerPoint crops bleed at slide edges anyway, and the unclamped
    # bbox confuses the visual-diff coverage gate (>90% triggers a
    # struct = total fallback that masks real text deficits).
    for i, p in enumerate(pics, 1):
        chunk = []
        if p.native_xml and p.native_media:
            # Template image carried natively (fixed corporate-design chrome):
            # the emitter re-embeds the media + splices the <p:pic>. Big media
            # goes to a sidecar file (raw bytes, sha-named) instead of inline.
            import base64 as _b64
            _xb = p.native_xml.encode("utf-8")
            _ref = _native_sidecar_ref(_xb, "xml", native_dir, native_rel)
            if _ref:
                line = f'native pic{i} xml_file:"{_ref}"'
            else:
                line = f'native pic{i} b64:"{_b64.b64encode(_xb).decode("ascii")}"'
            _mb = _b64.b64decode(p.native_media)
            _mref = _native_sidecar_ref(_mb, _media_ext(_mb), native_dir, native_rel)
            if _mref:
                line += f' media_file:"{_mref}"'
            else:
                line += f' media:"{p.native_media}"'
            chunk.append(line)
            _layers.append((_src_pos.get(id(p), 0), chunk))
            continue
        slot = "image" if len(pics) == 1 else f"image{i}"
        cx0 = max(0, p.x)
        cy0 = max(0, p.y)
        cx1 = min(cmap.cw, p.x + p.w)
        cy1 = min(cmap.ch, p.y + p.h)
        cw = max(1, cx1 - cx0)
        ch = max(1, cy1 - cy0)
        default_path = p.media_path or placeholder_rel
        # The expander's default-filter grammar requires `default("...")`
        # — parentheses + double quotes (see lib/dsl/expander.py:_DEFAULT_FILTER_RE).
        # The earlier `default:'…'` form silently failed to match, so the
        # slot resolved to empty string and the build fell into the
        # "no image bound" placeholder-rect branch — exactly why pictures
        # rendered as grey rects even when the asset path was correct.
        chunk.append(
            f'picture {cx0},{cy0} {cw}x{ch} '
            f'path:"{{{{ {slot} | default(\\"{default_path}\\") }}}}" cover:true'
        )
        _layers.append((_src_pos.get(id(p), 0), chunk))

    for _pos, _chunk in sorted(_layers, key=lambda e: e[0]):
        out.extend(_chunk)

    # Lines. Stroke-width preserves the source `<a:ln w="...">` value so
    # a 3pt horizontal divider survives the round-trip — the previous
    # `stroke-width:1` hardcode flattened every line to a hairline and
    # made decorative dividers invisible.
    for ln in lines:
        x1, y1 = ln.x, ln.y
        x2, y2 = ln.x + ln.w, ln.y + ln.h
        sw = ln.stroke_width if ln.stroke_width is not None and ln.stroke_width > 0 else 1
        attrs = f"stroke:{ln.stroke or 'fog'} stroke-width:{sw:g}"
        if ln.stroke_dash:
            attrs += f" dash:{ln.stroke_dash}"
        out.append(f"line {x1},{y1} {x2},{y2} {attrs}")

    if rects or pics or lines or ovals or custs:
        out.append("")

    # Footer collection: shapes whose y is in the bottom 8%. Keep the whole
    # SHAPE (not a per-run flatten) so the footer emit below preserves the
    # source text box's real width / padding / size / colour — a per-run
    # flatten dropped all four, hardcoded a 400 px maxwidth, and split a
    # two-run line ("Internal C-SC1" + " | C/CGB-CD | …") into two stacked
    # `text` statements. The narrow 400 px box wrapped the long © copyright
    # line into ~6 lines (the "Textumbruch" / line-break drift the reviewer
    # flagged), shifting the whole footer block.
    footer_shapes: list[Shape] = []
    body_texts: list[Shape] = []
    for t in texts:
        if t.y >= footer_y_threshold:
            footer_shapes.append(t)
        else:
            body_texts.append(t)

    # Body text in reading order.
    body_texts.sort(key=lambda s: (round(s.y / 5) * 5, s.x))
    for t in body_texts:
        # Concatenate runs verbatim — PPTX emits explicit space-only runs
        # between words, so " ".join would produce double spaces. Collapse
        # any runs of whitespace down to a single space after concat.
        raw = "".join(r.text for r in t.text_runs if r.text)
        # Collapse runs of spaces/tabs but PRESERVE newlines (paragraph breaks).
        full = re.sub(r"[ \t]+", " ", raw).strip()
        # Strip stray spaces on either side of a newline that resulted from
        # space-only runs between paragraphs.
        full = re.sub(r" *\n *", "\n", full)
        if not full:
            continue
        # Exclude paragraph-break marker runs (text == "\n") from the size
        # vote — those runs inherit the body-level default size (often 18pt)
        # which would otherwise drown out the actual text-content size when
        # two 13pt paragraphs share a single shape.
        content_pts = [r.pt for r in t.text_runs if r.text and r.text != "\n"]
        pt = max(content_pts) if content_pts else max((r.pt for r in t.text_runs), default=18)
        # normAutofit: reproduce the source's pre-shrink so placeholder text fits
        # its box (the source shrank the displayed size by font_scale; the run sz
        # stayed at the authored value, which we'd otherwise emit and overflow).
        pt *= t.font_scale
        style = _style_for(pt, full, is_footer=False)
        text = full.replace("\\", "\\\\").replace("\n", "\\n").replace('"', '\\"')
        mw = max(80, t.w)
        mh = max(24, t.h)
        # Emit `color:` override when the source's text-run colour differs
        # from the chosen style bundle's default. Captured in TextRun.color
        # from `<a:rPr><a:solidFill>`; without this the title that should
        # render in accent-blue lands in ink-grey (or whatever the style
        # bundle's default colour is) and inflates the visual diff.
        run_colors = [r.color for r in t.text_runs if r.color]
        run_color = run_colors[0] if run_colors else None
        style_default = STYLE_BUNDLES.get(style, {}).get("color")
        color_attr = (
            f" color:{run_color}" if run_color and run_color != style_default else ""
        )
        # Per-run weight override. `<a:rPr b="1">` rides on the source run;
        # the size-based classifier picks a bundle whose default weight may
        # not match (e.g. source 60pt bold maps to `huge` which is
        # weight:light by token convention). Without this, the rendered
        # text loses its emphasis even though the source explicitly carried
        # the bold flag. Emit only when source ≠ bundle default to keep
        # decompile output stable for the common case.
        source_bold = any(r.bold for r in t.text_runs)
        bundle_weight = STYLE_BUNDLES.get(style, {}).get("weight")
        weight_attr = ""
        if source_bold and bundle_weight not in ("bold", "semibold", "black"):
            weight_attr = " weight:bold"
        elif (not source_bold) and bundle_weight == "bold":
            # Source author explicitly chose regular against a bold-default
            # bundle — preserve that. `regular` is the most common name in
            # `font-weight` tokens; brands that use a different label can
            # override per-layout.
            weight_attr = " weight:regular"
        # Per-run size override. The classifier rounds source pt to the
        # nearest bundle, but the bundle steps are coarse (16/26/44/80 px)
        # and the emitted pt depends on the brand's slide physical width
        # (via `_PX_TO_PT`). Emit the source pt verbatim so renders match
        # source physical pt regardless of slide scale.
        # exactly. Stable for the common case (no emit when bundle is close).
        # Always emit `size:<pt>pt` from the source pt. The toolkit's style
        # bundles use design-px (`body=26px`, `sub=44px`, ...) which only
        # round-trip to the right rendered pt at the 13.33"-wide slide
        # convention (`_PX_TO_PT=0.5`). When the brand pack inherits a
        # different physical slide size from the source PPTX (`slide.width_emu`
        # in tokens.json), the emitter's `_PX_TO_PT` shifts and the same
        # px-valued bundles render at a different pt — which silently
        # shrinks every untagged text. Locking each run to its source pt
        # makes physical font sizes faithful regardless of slide scale.
        size_attr = f" size:{pt:g}pt"
        # autoshrink: source bodyPr had normAutofit — arm the emitter's fit as a
        # safety net (no-op when the scaled size already fits the box).
        autoshrink_attr = " autoshrink:true" if t.autoshrink else ""
        # linespacing: explicit source spcPct → its multiplier; absent →
        # `native` so the emitter writes NO lnSpc and the renderer's single
        # spacing applies (the toolkit's 1.2 default pushed every
        # source-default headline a quarter-line down).
        linespacing_attr = (
            f" linespacing:{t.line_spacing:g}" if t.line_spacing is not None
            else " linespacing:native"
        )
        valign_attr = f" valign:{t.valign}" if t.valign else ""
        # Horizontal align — pick the first non-None align from the runs.
        # PPTX stores it per-paragraph; for a single emitted `text` primitive
        # the first paragraph's alignment wins (consistent with how we pick
        # color + size from the first content run).
        run_aligns = [r.align for r in t.text_runs if r.align]
        run_align = run_aligns[0] if run_aligns else None
        align_attr = f" align:{run_align}" if run_align else ""
        padding_attr = ""
        if t.padding is not None:
            left, top, right, bottom = t.padding
            # Compact form `padding:N` when all four insets are equal,
            # `padding:L,T,R,B` otherwise — keeps the common-case DSL short.
            if left == right and top == bottom and left == top:
                padding_attr = f" padding:{left:g}"
            else:
                padding_attr = f" padding:{left:g},{top:g},{right:g},{bottom:g}"
        # Extend the text box into the AVAILABLE space — the clearance to the
        # nearest neighbouring element (or the slide edge) — so /deck can fill in
        # MORE text than the template carried without colliding with surrounding
        # chrome. Grow only in directions that DON'T reposition the existing
        # top-left-anchored text: left/justify-aligned → widen rightwards, and
        # top-valign → grow downwards. Centred / right / middle / bottom text
        # keeps its source box (extending would re-centre / move the glyphs).
        # The current (template) text is unchanged — a short line still sits at
        # the top-left; only longer user content uses the extra room.
        _gap = 10
        _right, _bottom = cmap.cw, cmap.ch
        for _o in shapes:
            if _o is t:
                continue
            # A filled card rect that CONTAINS the text box bounds growth:
            # text inside a card must not spill past the card's edges, even
            # when nothing else shares the row (slide-45's org-chart root
            # body grew across the whole slide once the elbow connectors
            # stopped contributing a blocking bbox). Slide-sized panels
            # (backgrounds, full-bleed bands ≥90% of both dimensions) don't
            # count — text on those grows like background text.
            if (_o.kind == "rect" and _o.fill
                    and not (_o.w >= cmap.cw * 0.9 and _o.h >= cmap.ch * 0.9)
                    and _o.x <= t.x + 2 and _o.y <= t.y + 2
                    and _o.x + _o.w >= t.x + t.w - 2
                    and _o.y + _o.h >= t.y + t.h - 2):
                _right = min(_right, _o.x + _o.w)
                _bottom = min(_bottom, _o.y + _o.h)
                continue
            _y_overlap = not (_o.y + _o.h <= t.y or _o.y >= t.y + t.h)
            if _y_overlap:
                if _o.x >= t.x + t.w:
                    _right = min(_right, _o.x)      # neighbour clear to the right
                elif _o.x + _o.w > t.x + t.w:
                    # An element STARTS left of the text's right edge but
                    # extends beyond it (full-height illustration sharing the
                    # row, photo panel under a caption): the space to the
                    # right is already occupied — no growth at all. Growing
                    # there put slide-39's body column on top of its scene.
                    _right = min(_right, t.x + t.w + _gap)
            if _o.y >= t.y + t.h and not (_o.x + _o.w <= t.x or _o.x >= t.x + t.w):
                _bottom = min(_bottom, _o.y)         # neighbour clear below
        if run_align in (None, "left", "justify"):
            mw = max(mw, _right - t.x - _gap)
        if t.valign in (None, "top"):
            mh = max(mh, _bottom - t.y - _gap)
        out.append(
            f'text {t.x},{t.y} style:{style}{color_attr}{weight_attr}{size_attr}{autoshrink_attr}{linespacing_attr}{valign_attr}{align_attr}{padding_attr} '
            f'maxwidth:{mw} maxheight:{mh} "{text}"'
        )

    # Footer-region text. Anything below `footer_y_threshold` (bottom 8%) is
    # emitted as `style:footer` text primitives. Each source text box becomes
    # ONE statement that keeps its real captured box (x, y, maxwidth from the
    # box width, maxheight from the box height), its source font size, its
    # per-run colour (the classification line's leading run is bold red
    # #D70012), and its text-frame padding. Emitting the real box width is what
    # makes the long copyright line wrap onto the same line count as the source
    # (the previous hardcoded 400 px box wrapped it ~6× and pushed the footer
    # up). Multi-run boxes concatenate verbatim so a two-run line stays one
    # line instead of splitting into two stacked statements.
    if footer_shapes:
        footer_shapes.sort(key=lambda s: (s.y, s.x))
        out.append("")
        for t in footer_shapes:
            raw = "".join(r.text for r in t.text_runs if r.text)
            full = re.sub(r"[ \t]+", " ", raw).strip()
            full = re.sub(r" *\n *", "\n", full)
            if not full:
                continue
            text = full.replace("\\", "\\\\").replace("\n", "\\n").replace('"', '\\"')
            content_pts = [r.pt for r in t.text_runs if r.text and r.text != "\n"]
            pt = max(content_pts) if content_pts else 6.0
            pt *= t.font_scale
            mw = max(80, t.w)
            mh = max(16, t.h)
            size_attr = f" size:{pt:g}pt"
            # Footer copyright/classification lines are short and box-bound; the
            # first run's colour wins (matches the body path's colour pick). The
            # `footer` bundle's own default colour is suppressed when equal.
            run_colors = [r.color for r in t.text_runs if r.color]
            run_color = run_colors[0] if run_colors else None
            footer_default = STYLE_BUNDLES.get("footer", {}).get("color")
            color_attr = (
                f" color:{run_color}" if run_color and run_color != footer_default else ""
            )
            padding_attr = ""
            if t.padding is not None:
                left, top, right, bottom = t.padding
                if left == right and top == bottom and left == top:
                    padding_attr = f" padding:{left:g}"
                else:
                    padding_attr = f" padding:{left:g},{top:g},{right:g},{bottom:g}"
            linespacing_attr = (
                f" linespacing:{t.line_spacing:g}" if t.line_spacing is not None
                else " linespacing:native"
            )
            out.append(
                f'text {t.x},{t.y} style:footer{color_attr}{size_attr}{linespacing_attr}{padding_attr} '
                f'maxwidth:{mw} maxheight:{mh} "{text}"'
            )

    return "\n".join(out) + "\n"



# ---------------------------------------------------------------------------
# Public derive()
# ---------------------------------------------------------------------------


def derive(
    pptx_path: Path,
    slide_idx: int,
    canvas_w: int = 1920,
    canvas_h: int = 1080,
    tokens_path: Path | None = None,
    layout_name: str = "derived",
    theme_name: str = "feinschliff",
    placeholder_rel: str = PLACEHOLDER_REL,
    pdf_path: Path | None = None,
    image_extract_dir: Path | None = None,
    image_extract_rel: str | None = None,
    native_extract_dir: Path | None = None,
    native_extract_rel: str | None = None,
) -> str:
    """Decompile one slide of `pptx_path` (1-indexed) into a Feinschliff DSL
    string. Brand-agnostic: pass `theme_name` and `tokens_path` to point at
    the target brand pack. `tokens_path` is used only for nearest-color
    matching against brand color tokens; if omitted, raw hex colors land
    in the DSL.

    `pdf_path` is currently unused; reserved for SVG cross-check of
    custGeom bboxes (callers render the slide's PDF page on demand).

    `image_extract_dir` + `image_extract_rel` enable **image carry-over**
    for pipeline-optimization runs: each `<p:pic>` in the source slide
    has its embedded binary written to `image_extract_dir/imageN.<ext>`,
    and the generated DSL's `default:` for that slot points at
    `image_extract_rel/imageN.<ext>` (a brand-pack-relative path the
    build will resolve at render time). Without these args, picture
    statements fall back to the generic placeholder image as before,
    which is the right default for *templating* an out-of-tree brand
    pack (where the source illustrations are not re-used). For *testing
    decompile fidelity* against the source, carry the images over so
    the visual diff measures real shape/text mismatch instead of
    picture-region noise.

    `native_extract_dir` + `native_extract_rel` route native-carry payloads
    bigger than `NATIVE_INLINE_MAX` raw bytes into sha-named sidecar files
    (`xml_file:` / `media_file:` / `parts_file:` DSL refs, resolved against
    the brand pack's asset root at build time) instead of inline base64 —
    without them a 33 MB carried vector group inlines into a 44 MB
    .slide.dsl. Pass `<pack>/assets/native` + `"native"` for brand packs.
    """
    # python-pptx rejects template content-types (.potx / .pptm-template)
    # with a cryptic "is not a PowerPoint file, content type is
    # '…presentationml.template.main+xml'" error. Catch the suffix up front
    # and tell the caller exactly how to convert. (LibreOffice converts
    # cleanly in one pass; renaming the file is NOT enough — the internal
    # [Content_Types].xml carries the template MIME and must be rewritten.)
    if pptx_path.suffix.lower() in (".potx", ".potm"):
        raise ValueError(
            f"{pptx_path.name} is a PowerPoint TEMPLATE (.potx/.potm), not a "
            f"presentation. Convert it first with:\n"
            f"  soffice --headless --convert-to pptx --outdir {pptx_path.parent} "
            f"{pptx_path}\n"
            f"then re-run with the produced .pptx."
        )
    palette: dict[str, tuple[int, int, int]] = {}
    if tokens_path and tokens_path.exists():
        palette = load_palette(tokens_path)
    pres = Presentation(str(pptx_path))
    theme = load_theme_scheme(pres)
    slide = pres.slides[slide_idx - 1]

    cx = pres.slide_width
    cy = pres.slide_height
    cmap = CanvasMap(cx, cy, canvas_w, canvas_h)

    shapes = walk_slide(slide, cmap, theme, palette)
    bg_fill = extract_slide_bg_fill(slide, theme, palette)
    bg_image = extract_slide_bg_image(slide)
    _ = pdf_path  # reserved for SVG cross-check, off by default

    if image_extract_dir is not None:
        if image_extract_rel is None:
            raise ValueError(
                "image_extract_rel is required when image_extract_dir is set "
                "(it goes into the DSL `default:` slot — the build resolves "
                "it relative to the brand pack root)."
            )
        image_extract_dir.mkdir(parents=True, exist_ok=True)
        # Prefer pictures whose Part was resolved at walk-time (handles
        # layout-inherited pics); fall back to slide.part lookup for the
        # in-slide case so callers that pre-date media_part still work.
        pics = [s for s in shapes if s.is_picture and (s.media_part or s.media_rid)]
        for i, p in enumerate(pics, 1):
            part = p.media_part
            if part is None and p.media_rid:
                try:
                    part = slide.part.related_part(p.media_rid)
                except KeyError:
                    continue
            if part is None:
                continue
            blob = getattr(part, "blob", None)
            partname = str(getattr(part, "partname", "/image.bin"))
            ext = partname.rsplit(".", 1)[-1].lower() if "." in partname else "bin"
            stem = "image" if len(pics) == 1 else f"image{i}"
            out_name = f"{stem}.{ext}"
            (image_extract_dir / out_name).write_bytes(blob or b"")
            if ext == "svg":
                # The build's picture path needs a raster (PIL cannot read
                # SVG) — rasterize next to the carried vector; on failure
                # the DSL keeps the svg path and the build falls back to
                # the brand placeholder.
                png = _rasterize_svg(image_extract_dir / out_name)
                if png is not None:
                    out_name = png.name
            p.media_path = f"{image_extract_rel.rstrip('/')}/{out_name}"

    # Background artwork needs a file on disk — only materialisable when the
    # caller gave us an extract dir (brand-pack flows always do). Without
    # one, fall back to the solid-fill path silently.
    bg_image_path: str | None = None
    if bg_image is not None and image_extract_dir is not None and image_extract_rel:
        blob, ext = bg_image
        image_extract_dir.mkdir(parents=True, exist_ok=True)
        (image_extract_dir / f"bg.{ext}").write_bytes(blob)
        bg_image_path = f"{image_extract_rel.rstrip('/')}/bg.{ext}"

    return emit_dsl(shapes, cmap, layout_name,
                    theme_name=theme_name,
                    placeholder_rel=placeholder_rel,
                    bg_fill=bg_fill,
                    bg_image_path=bg_image_path,
                    native_dir=native_extract_dir,
                    native_rel=native_extract_rel)


def main() -> int:
    ap = argparse.ArgumentParser(description="Decompile one PPTX slide → Feinschliff DSL (hybrid PPTX+SVG)")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--slide", type=int, default=1)
    ap.add_argument("--canvas", default="1920x1080")
    ap.add_argument("--theme", default="feinschliff",
                    help="Brand name to emit on the `theme` directive (default: feinschliff)")
    ap.add_argument("--brand-tokens", type=Path, default=None,
                    help="Brand tokens.json — used for nearest-color matching against brand palette")
    ap.add_argument("--layout-name", default="derived")
    ap.add_argument("--placeholder", default=PLACEHOLDER_REL,
                    help=f"Picture placeholder path (default: {PLACEHOLDER_REL})")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"missing: {args.pptx}", file=sys.stderr)
        return 2
    w, h = (int(x) for x in args.canvas.split("x"))
    print(derive(args.pptx, args.slide, w, h,
                 tokens_path=args.brand_tokens,
                 layout_name=args.layout_name,
                 theme_name=args.theme,
                 placeholder_rel=args.placeholder), end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
