"""Build the feinschliff house brand master.pptx.

Walks every layout under `feinschliff/brands/feinschliff/layouts/`, runs
`feinschliff build` with the matching content fixture, then merges every
rendered .pptx into a single master deck via deep XML clone — the same
mechanism feinschmiede.master_template.ClonePlan uses at render time.

Output:
    feinschliff/brands/feinschliff/master.pptx
    feinschliff/brands/feinschliff/snippets.yaml (one entry per layout)

This is a one-time build script. After PR4 ships, the feinschliff house
brand routes through feinschliff.master_template, not the DSL pipeline,
and the DSL/ module is deleted.

Run:
    uv run python feinschliff-builder/scripts/build_feinschliff_master.py
    uv run python feinschliff-builder/scripts/build_feinschliff_master.py --workers 8
"""
from __future__ import annotations

import argparse
import concurrent.futures as futures
import copy
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[2]

# Brand-pack roots known to this builder. Add new entries as DSL brand packs
# migrate to the master-template format. Each entry resolves layouts +
# content fixtures from convention.
_BRANDS: dict[str, dict[str, Path]] = {
    "feinschliff": {
        "brand_dir": REPO_ROOT / "feinschliff/brands/feinschliff",
        "fixtures":  REPO_ROOT / "tests/feinschliff/fixtures/layouts",
    },
    "gs-ramspau": {
        "brand_dir": REPO_ROOT / "feinschliff-extra/brands/gs-ramspau",
        "fixtures":  REPO_ROOT / "feinschliff-extra/brands/gs-ramspau/tests/fixtures/layouts",
    },
}

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS = {"a": _A}


# Deck-shape order: opening → agenda → headline narrative → text content →
# photo/visual → charts/data → diagrams/process → matrices → quote → close.
# Layouts not listed get appended after the curated order (alphabetical).
_ORDER = [
    # 1. Opening
    "title-orange",
    "title-ink",
    "full-bleed-cover",
    "full-bleed-editorial",
    # 2. Chapter dividers
    "chapter-orange",
    "chapter-ink",
    # 3. Agenda / overview
    "agenda",
    "agenda-photo",
    # 4. Headline + body
    "executive-summary",
    "action-title",
    "key-takeaways",
    "recommendation",
    "next-steps",
    # 5. Text columns
    "two-column-cards",
    "three-column",
    "four-column-cards",
    "vertical-bullets",
    "horizontal-bullets",
    "text-picture",
    "quote",
    # 6. KPIs / scorecards
    "kpi-grid",
    "kpi-photo",
    "scorecard",
    # 7. Photo / visual
    "graphical",
    "photo-grid",
    "photo-strip-four",
    "chart-photo",
    # 8. Charts
    "bar-chart",
    "stacked-bar",
    "line-chart",
    "waterfall",
    # 9. Timelines / roadmaps
    "timeline",
    "gantt",
    "roadmap",
    "process-flow",
    # 10. Structured diagrams
    "pyramid",
    "funnel",
    "venn",
    "v-model",
    "2x2-matrix",
    "risk-matrix",
    "risk-register",
    "table",
    # 11. Bespoke / showcase
    "components-showcase",
    "excalidraw-diagram",
    "excalidraw-diagram-full",
    "svg-infographic",
    "svg-infographic-full",
    # 12. Closing
    "end-image",
    "end",
]


def _layout_ids(layouts_dir: Path) -> list[str]:
    available = {p.stem.replace(".slide", "") for p in layouts_dir.glob("*.slide.dsl")}
    ordered = [lid for lid in _ORDER if lid in available]
    leftover = sorted(available - set(ordered))
    return ordered + leftover


def _build_one(layout_id: str, *, brand: str, layouts_dir: Path, fixtures: Path) -> Path | None:
    layout_path = layouts_dir / f"{layout_id}.slide.dsl"
    content_path = fixtures / f"{layout_id}.yaml"
    if not content_path.exists():
        print(f"  SKIP {layout_id}: no fixture at {content_path}", file=sys.stderr)
        return None
    tmp = Path(tempfile.mkdtemp(prefix=f"master-{layout_id}-"))
    pptx = tmp / f"{layout_id}.pptx"
    cmd = [
        "uv", "run", "feinschliff", "build",
        str(layout_path),
        "--brand", brand,
        "--content", str(content_path),
        "-o", str(pptx),
        "--skip-content-lint",
        "--allow-missing-assets",
        "--allow-diagram-warnings",
    ]
    res = subprocess.run(cmd, cwd=REPO_ROOT, capture_output=True, text=True)
    if res.returncode != 0:
        shutil.rmtree(tmp, ignore_errors=True)
        print(f"  FAIL {layout_id}\n{res.stderr[-400:]}", file=sys.stderr)
        return None
    return pptx


_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _merge_into_master(per_layout_pptx: list[tuple[str, Path]], *, brand: str, out_pptx: Path) -> list[dict]:
    """Open the first .pptx as the master, clone every other deck's slide(s)
    into it as snippets. Each cloned slide's images are re-added to the
    master via the proper python-pptx image API so part names don't collide.
    Returns the snippets index for snippets.yaml."""
    first_id, first_pptx = per_layout_pptx[0]
    prs = Presentation(str(first_pptx))
    snippets = []
    for i, _ in enumerate(prs.slides):
        snippets.append({"id": first_id, "source_idx": i, "intent": f"{brand}/{first_id}"})

    dest_layout = prs.slide_layouts[0]
    for layout_id, pptx_path in per_layout_pptx[1:]:
        src = Presentation(str(pptx_path))
        for src_slide in src.slides:
            new_slide = prs.slides.add_slide(dest_layout)
            for shape in list(new_slide.shapes):
                shape._element.getparent().remove(shape._element)

            # Map src image rels → new rIds in the master before cloning XML,
            # so we can rewrite r:embed / r:link attributes on each cloned shape.
            rid_map: dict[str, str] = {}
            for src_rid, rel in src_slide.part.rels.items():
                if rel.reltype.endswith("/image"):
                    from io import BytesIO
                    _, new_rid = new_slide.part.get_or_add_image_part(BytesIO(rel.target_part.blob))
                    rid_map[src_rid] = new_rid

            for child in list(src_slide.shapes._spTree):
                if etree.QName(child).localname in ("nvGrpSpPr", "grpSpPr"):
                    continue
                cloned = copy.deepcopy(child)
                for el in cloned.iter():
                    embed = el.get(f"{_R}embed")
                    if embed and embed in rid_map:
                        el.set(f"{_R}embed", rid_map[embed])
                    link = el.get(f"{_R}link")
                    if link and link in rid_map:
                        el.set(f"{_R}link", rid_map[link])
                new_slide.shapes._spTree.append(cloned)

            snippets.append({
                "id": layout_id,
                "source_idx": len(snippets),
                "intent": f"{brand}/{layout_id}",
            })
    prs.save(str(out_pptx))
    return snippets


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--brand", default="feinschliff", choices=sorted(_BRANDS), help="Brand id to build")
    p.add_argument("--workers", type=int, default=max(1, (os.cpu_count() or 2) // 2))
    p.add_argument("--only", nargs="*", help="Restrict to specific layout ids")
    args = p.parse_args()

    cfg = _BRANDS[args.brand]
    brand_dir = cfg["brand_dir"]
    layouts_dir = brand_dir / "layouts"
    fixtures = cfg["fixtures"]
    master_pptx = brand_dir / "master.pptx"
    snippets_yaml = brand_dir / "snippets.yaml"
    layouts_yaml = brand_dir / "layouts.yaml"
    master_theme = "default"

    ids = _layout_ids(layouts_dir)
    if args.only:
        ids = [i for i in ids if i in args.only]
    print(f"building {len(ids)} layouts for {args.brand} with --workers {args.workers}…")

    results: dict[str, Path] = {}
    with futures.ThreadPoolExecutor(max_workers=args.workers) as ex:
        futs = {
            ex.submit(_build_one, lid, brand=args.brand, layouts_dir=layouts_dir, fixtures=fixtures): lid
            for lid in ids
        }
        for fut in futures.as_completed(futs):
            lid = futs[fut]
            pptx = fut.result()
            if pptx is not None:
                results[lid] = pptx
                print(f"  built {lid}")

    if not results:
        print("no layouts built; aborting", file=sys.stderr)
        return 1

    per_layout = [(lid, results[lid]) for lid in ids if lid in results]
    snippets = _merge_into_master(per_layout, brand=args.brand, out_pptx=master_pptx)
    snippets_yaml.write_text(yaml.safe_dump({"snippets": snippets}, sort_keys=False))
    layouts_yaml.write_text(
        f"# {args.brand} brand is snippet-only — every layout is a ClonePlan\n"
        "# against snippets.yaml. No FillPlan entries.\n"
        f"master_theme: {master_theme}\nlayouts: []\n"
    )

    print(f"\nwrote {master_pptx} ({master_pptx.stat().st_size / 1024:.1f} KB, {len(snippets)} slides)")
    print(f"wrote {snippets_yaml} ({len(snippets)} snippets)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
