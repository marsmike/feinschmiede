"""End-to-end render against the real BSH master template.

Skipped unless the BSH master pptx is available at the canonical local path.
This is the verify-fast smoke for PR1 — proves the module reproduces the
proof-script behavior through the public API.
"""

from __future__ import annotations

import os
from pathlib import Path

import pytest

from feinschmiede.master_template import (
    ChartSpec,
    ClonePlan,
    FillPlan,
    PictureRef,
    load_catalog,
    render,
)

HOME = Path(os.path.expanduser("~"))
BSH_MASTER = HOME / "work/pptx-templates/BSH-Master-Templates-2026-06-12.pptx"
BSH_SOURCE = HOME / "work/pptx-templates/BSH-Templates-2026-06-12.pptx"
SAMPLE_IMAGE = HOME / "work/feinschliff-bsh/.work/vfinal-035.png"

pytestmark = pytest.mark.skipif(
    not (BSH_MASTER.exists() and BSH_SOURCE.exists()),
    reason="BSH master + source decks not present locally",
)


@pytest.fixture
def bsh_pack(tmp_path: Path) -> Path:
    pack = tmp_path / "bsh-pack-v5"
    pack.mkdir()
    (pack / "master.pptx.ref").write_text(str(BSH_MASTER))
    (pack / "source_deck.pptx.ref").write_text(str(BSH_SOURCE))
    (pack / "layouts.yaml").write_text(
        """
layouts:
  - name: "Title + Graphical Content + Text"
    role: content
    placeholders:
      - idx: 0
        type: TITLE
        role: headline
      - idx: 1
        type: OBJECT
        role: body_or_chart
  - name: "Text and Picture"
    role: content
    placeholders:
      - idx: 0
        type: TITLE
      - idx: 1
        type: BODY
      - idx: 2
        type: PICTURE
"""
    )
    (pack / "snippets.yaml").write_text(
        """
snippets:
  - id: timeline-4
    source_idx: 17
    intent: bespoke timeline ribbon
"""
    )
    return pack


def test_fill_plan_text_only(bsh_pack, tmp_path):
    out = tmp_path / "out.pptx"
    plan = FillPlan(
        layout="Title + Graphical Content + Text",
        fills={
            0: "Headline",
            1: ["First bullet", "Second bullet", "Third bullet"],
        },
    )
    render(bsh_pack, [plan], out)
    assert out.exists() and out.stat().st_size > 0

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    title = next(p for p in slide.placeholders if p.placeholder_format.idx == 0)
    assert title.text_frame.text == "Headline"


def test_fill_plan_with_chart(bsh_pack, tmp_path):
    out = tmp_path / "out.pptx"
    plan = FillPlan(
        layout="Title + Graphical Content + Text",
        fills={
            0: "Chart slide",
            1: ChartSpec(
                kind="column",
                categories=["Q1", "Q2", "Q3", "Q4"],
                series=[("Revenue", [14.2, 17.8, 19.4, 22.6])],
            ),
        },
    )
    render(bsh_pack, [plan], out)

    from pptx import Presentation

    prs = Presentation(str(out))
    slide = prs.slides[0]
    assert any(shape.has_chart for shape in slide.shapes)


@pytest.mark.skipif(not SAMPLE_IMAGE.exists(), reason="no local sample image")
def test_fill_plan_with_picture(bsh_pack, tmp_path):
    out = tmp_path / "out.pptx"
    plan = FillPlan(
        layout="Text and Picture",
        fills={
            0: "Picture slide",
            1: "Caption body.",
            2: PictureRef(path=SAMPLE_IMAGE),
        },
    )
    render(bsh_pack, [plan], out)

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 1


def test_clone_plan_preserves_shape_count(bsh_pack, tmp_path):
    out = tmp_path / "out.pptx"
    plan = ClonePlan(
        snippet_id="timeline-4",
        text_replacements=[("Timeline 4", "Roadmap H1 2026")],
    )
    render(bsh_pack, [plan], out)

    from pptx import Presentation

    src = Presentation(str(BSH_SOURCE))
    dst = Presentation(str(out))
    src_count = len(list(src.slides[17].shapes._spTree))
    dst_count = len(list(dst.slides[0].shapes._spTree))
    assert dst_count >= src_count - 2

    a_ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    runs = dst.slides[0].shapes._spTree.findall(".//a:t", a_ns)
    assert any(r.text == "Roadmap H1 2026" for r in runs)


def test_strip_clears_existing_slides(bsh_pack, tmp_path):
    out = tmp_path / "empty.pptx"
    render(bsh_pack, [], out)
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 0
