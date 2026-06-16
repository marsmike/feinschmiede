"""Open the brand master, optionally overlay a theme, play the plan list."""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation

from feinschliff.master_template._brand import index_layouts, master_path, norm
from feinschliff.master_template.clone_plan import ClonePlan, apply_clone
from feinschliff.master_template.fill_plan import FillPlan, apply_fill
from feinschliff.master_template.theme_overlay import apply_theme, base_palette

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _strip_existing_slides(prs) -> None:
    """Remove every sample slide from the master's presentation.

    Drops the `sldId` entries and the presentation->slide rels. Slide Parts
    themselves stay in the package — `ClonePlan` reuses their image rels
    through the package's rel graph to dedup media by partname (load-bearing
    for same-master split clones).
    """
    rels = prs.part.rels
    for sldId in list(prs.slides._sldIdLst):
        rid = sldId.get(f"{{{_R_NS}}}id")
        prs.slides._sldIdLst.remove(sldId)
        if rid in rels:
            rels.pop(rid)


def render(
    brand_pack: Path,
    plans: Iterable[FillPlan | ClonePlan],
    out: Path,
    *,
    theme: Path | dict | None = None,
    source_deck: Path | None = None,
) -> Path:
    """Render `plans` against `brand_pack`'s master, write to `out`.

    Optional `theme` overlays scheme colors onto the master before the
    plans dispatch — same master, N visual variations. `source_deck` is
    the deck cloned slides are read from; defaults to the brand's own
    master.
    """
    brand_pack = Path(brand_pack)
    prs = Presentation(str(master_path(brand_pack)))
    if theme is not None:
        apply_theme(prs, theme, recolor_from=base_palette(brand_pack))
    _strip_existing_slides(prs)

    source_prs = None
    layout_by_name = index_layouts(prs)

    for plan in plans:
        if isinstance(plan, FillPlan):
            layout = layout_by_name.get(norm(plan.layout))
            if layout is None:
                raise KeyError(f"layout not in master: {plan.layout!r}")
            apply_fill(prs.slides.add_slide(layout), plan)
        elif isinstance(plan, ClonePlan):
            if source_prs is None:
                source_prs = Presentation(str(source_deck or master_path(brand_pack)))
            apply_clone(prs, source_prs, plan)
        else:
            raise TypeError(f"unknown plan type: {type(plan).__name__}")

    out = Path(out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
