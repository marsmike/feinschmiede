"""Layout-fill plans — pick a layout, fill placeholders by index.

Values may be `str`, `list[str]` (one paragraph per item), `PictureRef`, or
`ChartSpec`. Charts and pictures replace the OBJECT placeholder at its bbox.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Sequence


@dataclass
class PictureRef:
    path: Path
    crop: bool = True  # center-crop to bbox aspect before insertion


@dataclass
class ChartSpec:
    kind: str  # XL_CHART_TYPE member name: "column_clustered", "bar_clustered", "line"...
    categories: Sequence[str]
    series: list[tuple[str, Sequence[float]]]


@dataclass
class FillPlan:
    layout: str
    fills: dict[int, object] = field(default_factory=dict)


def apply_fill(slide, plan: FillPlan) -> None:
    by_idx = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    for idx, value in plan.fills.items():
        ph = by_idx.get(idx)
        if ph is None:
            continue
        if isinstance(value, ChartSpec):
            _replace_with_chart(slide, ph, value)
        elif isinstance(value, PictureRef):
            _fill_picture(slide, ph, value)
        elif isinstance(value, (list, tuple)):
            tf = ph.text_frame
            tf.text = str(value[0])
            for line in value[1:]:
                tf.add_paragraph().text = str(line)
        else:
            ph.text_frame.text = str(value)


def _replace_with_chart(slide, ph, spec: ChartSpec) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    L, T, W, H = ph.left, ph.top, ph.width, ph.height
    ph._element.getparent().remove(ph._element)

    data = CategoryChartData()
    data.categories = list(spec.categories)
    for name, values in spec.series:
        data.add_series(name, tuple(values))

    slide.shapes.add_chart(
        getattr(XL_CHART_TYPE, spec.kind.upper()), L, T, W, H, data
    )


def _fill_picture(slide, ph, ref: PictureRef) -> None:
    # PICTURE placeholders honor the master's authored crop / frame.
    if str(ph.placeholder_format.type).startswith("PICTURE") and hasattr(ph, "insert_picture"):
        ph.insert_picture(str(ref.path))
        return
    # OBJECT placeholders get replaced — center-crop first so
    # add_picture's stretch-to-bbox doesn't distort.
    L, T, W, H = ph.left, ph.top, ph.width, ph.height
    ph._element.getparent().remove(ph._element)
    path = _crop_to_aspect(Path(ref.path), W, H) if ref.crop else Path(ref.path)
    slide.shapes.add_picture(str(path), L, T, width=W, height=H)


def _crop_to_aspect(image_path: Path, target_w: int, target_h: int) -> Path:
    """Center-crop to `target_w/target_h` ratio. Returns a cached sibling file."""
    from PIL import Image

    ratio = target_w / target_h
    out = image_path.with_name(f"{image_path.stem}__crop{image_path.suffix}")
    if out.exists() and out.stat().st_mtime >= image_path.stat().st_mtime:
        return out

    img = Image.open(image_path)
    w, h = img.size
    if abs(w / h - ratio) < 0.001:
        return image_path
    if w / h > ratio:
        new_w = int(h * ratio)
        box = ((w - new_w) // 2, 0, (w - new_w) // 2 + new_w, h)
    else:
        new_h = int(w / ratio)
        box = (0, (h - new_h) // 2, w, (h - new_h) // 2 + new_h)
    img.crop(box).save(out, quality=92)
    return out
