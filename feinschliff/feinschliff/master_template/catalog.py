"""Inspect a brand pack's master.pptx and emit `layouts` + `snippets` catalogs.

    python -m feinschliff.master_template.catalog brands/feinschliff
"""
from __future__ import annotations

import sys
from collections import Counter
from pathlib import Path

import yaml
from pptx import Presentation

from feinschliff.master_template._brand import index_layouts, master_path, norm

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _snippet_anchors(slide) -> dict[str, int]:
    """Count `<a:t>` text runs per distinct string. Emitted only for snippets
    where queueing matters — i.e., at least one string appears >= 2 times.
    Authors composing a ClonePlan use this to size the replacements list."""
    counts = Counter(
        t.text
        for t in slide.shapes._spTree.iter(f"{{{_A_NS}}}t")
        if t.text and t.text.strip()
    )
    if not any(c >= 2 for c in counts.values()):
        return {}
    return dict(counts.most_common())


def build_catalog(brand_pack: Path) -> dict:
    prs = Presentation(str(master_path(brand_pack)))

    # index_layouts applies the same normalization + collision policy the
    # renderer uses, so the catalog can't list a layout name the renderer
    # would reject (and the two no longer carry duplicate collision loops).
    layouts = [
        {
            "name": name,
            "placeholders": [
                {
                    "idx": ph.placeholder_format.idx,
                    "type": ph.placeholder_format.type.name if ph.placeholder_format.type else None,
                    "name": ph.name,
                }
                for ph in layout.placeholders
            ],
        }
        for name, layout in index_layouts(prs).items()
    ]

    snippets = []
    for i, slide in enumerate(prs.slides):
        preview = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                preview = shape.text_frame.text.strip().splitlines()[0][:80]
                break
        entry = {
            "source_idx": i,
            "layout": norm(slide.slide_layout.name),
            "preview": preview,
        }
        anchors = _snippet_anchors(slide)
        if anchors:
            entry["anchors"] = anchors
        snippets.append(entry)

    return {"layouts": layouts, "snippets": snippets}


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print("usage: python -m feinschliff.master_template.catalog <brand_pack>", file=sys.stderr)
        return 2
    yaml.safe_dump(build_catalog(Path(argv[1])), sys.stdout, sort_keys=False, allow_unicode=True)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
