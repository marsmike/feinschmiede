"""`feinschliff deck …` subcommands: multi-slide composer + layout picker.

Subcommands:

  feinschliff deck build <plan.yaml> [-o OUT.pptx]
      Build one .pptx from a plan that lists N slides, each pinning a
      layout and inline content (or a content file).

  feinschliff deck pick <signals.yaml>
      Print a recommended layout id based on structured signals.

  feinschliff deck wireframe <layout.slide.dsl> --brand <id> [-o out.svg]
      Render a DSL layout as an annotated SVG wireframe (bounding boxes for
      text slots, picture slots, rect backgrounds). No PPTX round-trip needed.
      Add --overlay-pptx <file.pptx> to embed the actual rendered slide behind
      the wireframe boxes for pixel-accurate deviation analysis.

  feinschliff deck wireframe-sheet <plan.yaml> [-o sheet.svg]
      Render every slide in a plan as a wireframe and compose them into a
      single SVG contact sheet. Useful as a fast layout-regression baseline:
      store the SVGs in git, regenerate after DSL changes, and diff.

Plan schema (build):

  brand: feinschliff                       # default brand for slides
  out:   deck.pptx                         # output path; --output overrides
  slides:
    - layout: layouts/title-orange.slide.dsl
      content:
        pgmeta: "Q1 2026"
        title:  "..."
    - layout: layouts/quote.slide.dsl
      content_file: examples/v2/quote.yaml  # alternative to inline
    - layout: brands/gs-ramspau/layouts/stundenplan.slide.dsl
      brand:  gs-ramspau                    # per-slide override
      content_file: brands/gs-ramspau/examples/v2/stundenplan.yaml
"""
from __future__ import annotations

import argparse
import os
import sys
import time
import tempfile
from pathlib import Path

import yaml

from feinschliff.deck.content_metadata import auto_bind_slots, warn_overbudget_slots
from feinschliff.deck.orchestrate import (
    patch_set_hash as _patch_set_hash_fn,
)

from feinschmiede.dsl.tokens import load_tokens, load_tokens_with_theme
from feinschliff.content_validator import (
    emit_defects_and_abort_message, validate_content,
)
from feinschliff.defects import fatal_kinds, format_defect
from feinschmiede.brand_discovery import find_brand
from feinschliff.io.image_provider import discover_providers, get_provider
from feinschliff.io.image_providers import chain_from_brand_config

# Core verify imports — available without feinschliff-builder.
from feinschliff.verify.deck.notes_budget import validate_notes
from feinschliff.verify.deck.titles import extract_titles_from_plan
from feinschliff.verify.deck.storyline import render_contact_sheet, write_storyline_report
from feinschliff.verify.deck.claim_evidence import judge_plan, write_report as write_claim_evidence_report
from feinschliff.verify.deck.ghost_deck import judge_ghost_deck, write_ghost_deck_report
from feinschliff.verify.deck.title_lint import lint_titles
from feinschliff.pipeline_log import log_event


def _require_or_delegate_builder(feature: str) -> None:
    """Ensure a builder-backed deck subcommand can run.

    The advanced deck features (storyline, wireframe, polish, book,
    strict-static/autofix, …) are implemented on top of the
    ``feinschliff_builder`` package. Under the per-plugin launcher model the
    office venv does *not* contain that package — only the builder plugin's own
    venv does (it bundles office + engine + builder). So:

      * dev / builder venv — ``import feinschliff_builder`` succeeds → run inline
        (unchanged behaviour, the path every test exercises).
      * office-only venv with the builder plugin installed — re-exec the same
        command through the ``feinschliff-builder`` CLI (a capability call, the
        family's coupling rule), which has the package importable.
      * neither — exit 2 with an accurate install hint.
    """
    try:
        import feinschliff_builder  # noqa: F401
        return
    except ImportError:
        pass

    import shutil
    builder_cli = shutil.which("feinschliff-builder")
    if builder_cli:
        # `feinschliff deck <args>` → `feinschliff-builder deck <args>`; the
        # builder CLI re-registers this exact deck parser in a venv where
        # feinschliff_builder imports, so the inline path runs there. execv
        # replaces this process and never returns (the return is for testability).
        os.execv(builder_cli, [builder_cli, *sys.argv[1:]])
        return

    sys.stderr.write(
        f"error: '{feature}' is an advanced deck feature provided by the "
        f"feinschliff-builder plugin.\n"
        f"  Install it with: /plugin install feinschliff-builder@feinschmiede\n"
        f"  (or `uv pip install feinschliff-builder` in a dev checkout).\n"
    )
    raise SystemExit(2)


def _bundled_assets() -> Path:
    """Return the assets/ directory shipped inside this plugin."""
    return Path(__file__).resolve().parents[1] / "assets"


def _bundled_compounds() -> Path:
    """Return the compounds/ directory shipped inside the feinschmiede engine."""
    import feinschmiede

    return Path(feinschmiede.__file__).resolve().parent / "compounds"


def _find_toolkit_file(rel: str) -> Path | None:
    """Resolve *rel* against each discovered layout dir's parent (i.e. the
    plugin root), returning the first match.  Used to replace hard-coded
    ``REPO_ROOT / rel`` fallback lookups in the deck CLI.
    """
    from feinschliff.layout_discovery import all_layout_dirs
    for layout_dir in all_layout_dirs():
        candidate = (layout_dir.parent / rel).resolve()
        if candidate.is_file():
            return candidate
    return None


def register(parser: argparse.ArgumentParser) -> None:
    sub = parser.add_subparsers(dest="deck_command", required=True)

    p_build = sub.add_parser("build", help="Build a multi-slide deck from a plan YAML")
    p_build.add_argument("plan", help="Path to the deck plan YAML")
    p_build.add_argument("-o", "--output", help="Output .pptx (overrides plan.out)")
    p_build.add_argument(
        "--skip-content-lint",
        action="store_true",
        help="Skip pre-render content lints (title-length, action-verb-leading). "
             "For emergency overrides only.",
    )
    p_build.add_argument(
        "--no-image-provider",
        action="store_true",
        help="Ignore the brand's $image_provider: unbound image slots render "
             "their carried default path (or the placeholder) instead of a "
             "provider search. For showcase/fidelity builds, where a stale "
             "asset_lock.json or image cache must never override the pack's "
             "own assets.",
    )
    p_build.add_argument(
        "--slot-debug-color",
        metavar="#RRGGBB",
        help="Slot-coverage debugging: render every slot-sourced text "
             "(regular text slots AND native-payload slots) in this colour. "
             "Diff against a defaults build to see which text is bindable "
             "vs baked chrome.",
    )
    p_build.add_argument(
        "--allow-diagram-warnings",
        action="store_true",
        help="Ship even when diagram-overflow or diagram-text-too-small "
             "defects surface. Otherwise these are fatal by default — same "
             "policy as single-slide `feinschliff build`.",
    )
    p_build.add_argument(
        "--allow-missing-assets",
        action="store_true",
        help="Ship even when a picture slot points at a missing file or is "
             "unset. Default: fatal. Mark intentionally-empty slots with "
             "`optional:true` to skip the abort without this flag.",
    )
    p_build.add_argument(
        "--strict-static",
        action="store_true",
        help="Promote static-verifier WARNs to build-blocking errors. The "
             "pre-render static geometry verifier (feinschliff_builder.verify.static) "
             "runs on every build when feinschliff-builder is installed; FATAL "
             "defects (slot-overflow) abort the build by default, no flag "
             "needed. With this flag, WARN-level defects (empty-placeholder) "
             "abort too (exit 1, defects printed).",
    )
    p_build.add_argument(
        "--strict-craft",
        action="store_true",
        help="Run Knaflic deterministic craft-rules checks (no-pie-chart, "
             "chart-title-claim, title-word-count, …) over every slide after "
             "compile. Writes craft_report.md next to the output deck. "
             "Exit 1 if verdict is 'fail'; exit 0 on 'warn' or 'clean'.",
    )
    p_build.add_argument(
        "--autofix",
        action="store_true",
        help="Run the static verifier before compile and automatically apply "
             "mechanical fixes (shorten_slot, delete_word, drop_bullet, "
             "swap_layout_*) for known defect classes.  Up to 3 inner fix "
             "cycles are attempted; residual WARN defects are printed but do "
             "NOT block the compile.  Residual FATAL defects (slot-overflow) "
             "abort via the default static gate.  The fixed plan is "
             "written back to disk before compile.",
    )
    p_build.add_argument(
        "--embed-fonts",
        action="store_true",
        help="Embed brand display/body font files into the .pptx so "
             "recipients without the fonts render faithfully (opt-in; "
             "enlarges the file). "
             "No font-license (fsType) check is performed — verify your "
             "brand fonts permit embedding.",
    )
    p_build.add_argument(
        "--workers",
        type=int,
        default=16,
        help="Parallel slide-compile workers (DSL expand + diagram render "
             "fan out across processes; the final PPTX assembly stays "
             "serial). Default 16; 0 = auto (min(8, cpu/2)); 1 = sequential.",
    )
    p_build.add_argument(
        "--strict-visual",
        action="store_true",
        help="After building, render slides to PNG and run AeSlides visual "
             "metrics (whitespace ratio, visual balance, element collision). "
             "Writes out/visual_metrics_report.md. Requires soffice or "
             "pdftoppm; skipped with a warning when neither is available. "
             "Opt-in only — omitting this flag skips all PNG rendering.",
    )
    p_build.add_argument(
        "--no-slide-numbers",
        action="store_true",
        help="Suppress the automatic 'NN / TOTAL' slide-number footer that is "
             "stamped in the bottom-right corner of every slide by default.",
    )
    p_build.add_argument(
        "--theme",
        metavar="NAME",
        default=None,
        help="Color theme to apply to the default brand (e.g. 'default', 'claude'). "
             "Overrides the brand's declared $default_theme. Has no effect on "
             "brands that have no themes/ directory.",
    )
    p_build.set_defaults(func=cmd_build)

    p_pick = sub.add_parser("pick", help="Recommend a layout for the given signals")
    p_pick.add_argument("signals", help="Path to a signals YAML (or '-' for stdin)")
    p_pick.add_argument("--top-k", type=int, default=3,
                        help="Print the top-K candidates with scores (default 3)")
    p_pick.set_defaults(func=cmd_pick)

    p_pick_deck = sub.add_parser(
        "pick-deck",
        help="Run the arc-aware deck-level picker over a plan.yaml. "
             "Emits picker_report.json with chosen layouts, runners-up, "
             "and arc warnings.",
    )
    p_pick_deck.add_argument("plan", help="Path to the deck plan YAML")
    p_pick_deck.add_argument(
        "--deck-brief",
        default=None,
        help="Path to a deck_brief.yaml/json. When provided and deck_type is "
             "set, arc-position warnings are computed against the matching arc schema.",
    )
    p_pick_deck.add_argument(
        "-o", "--output",
        default=None,
        help="Output path for picker_report.json. Default: picker_report.json "
             "next to plan.yaml.",
    )
    p_pick_deck.add_argument(
        "--top-k", type=int, default=5,
        help="Candidates per slide passed to pick_layout (default 5).",
    )
    p_pick_deck.set_defaults(func=cmd_pick_deck)

    p_intake_v = sub.add_parser(
        "intake-validate",
        help="Validate a deck_brief.yaml against the schema. "
             "Exit 0 = valid, 1 = invalid (errors to stderr), 2 = plumbing.",
    )
    p_intake_v.add_argument("brief", help="Path to deck_brief.yaml")
    p_intake_v.set_defaults(func=cmd_intake_validate)

    p_intake_s = sub.add_parser(
        "intake-skeleton",
        help="Emit a deck_brief.yaml skeleton seeded by heuristic "
             "inference. The orchestrating skill fills the remaining "
             "fields via AskUserQuestion before writing the final brief.",
    )
    p_intake_s.add_argument(
        "--brief", default=None,
        help="Optional brief text, or a path to a brief file. "
             "When supplied, infer_from_text seeds known fields.",
    )
    p_intake_s.add_argument(
        "-o", "--output", default=None,
        help="Output path. Default: stdout.",
    )
    p_intake_s.set_defaults(func=cmd_intake_skeleton)

    p_commit_v = sub.add_parser(
        "commitment-validate",
        help="Validate a commitment.yaml against the schema and "
             "(optionally) cross-check arc-alignment against the matching "
             "deck_type arc schema. Exit 0 clean / 1 invalid / 2 plumbing.",
    )
    p_commit_v.add_argument("commitment", help="Path to commitment.yaml")
    p_commit_v.add_argument(
        "--check-arc",
        action="store_true",
        help="Also run check_arc_alignment against the loaded arc schema "
             "for the commitment's deck_type.",
    )
    p_commit_v.set_defaults(func=cmd_commitment_validate)

    p_storyline = sub.add_parser(
        "storyline",
        help="Emit a title-only contact sheet from a deck_plan.json / "
             "content_plan.json (Phase 2 storyline gate).",
    )
    p_storyline.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_storyline.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_storyline.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_storyline.set_defaults(func=cmd_storyline)

    p_ce = sub.add_parser(
        "claim-evidence",
        help="Mid-plan claim-evidence text gate (step 2b). "
             "Judges each claim-carrying slide for title-body coherence "
             "before render. Cheap Haiku pass, no PPTX round-trip.",
    )
    p_ce.add_argument("plan", help="Path to plan.yaml or plan.json")
    p_ce.add_argument(
        "--design-brief",
        default=None,
        help="Path to design_brief.json (optional — enables per-slide claim hints).",
    )
    p_ce.add_argument(
        "-o", "--output", required=True,
        help="Output path for claim_evidence_report.md",
    )
    p_ce.add_argument(
        "--offline", action="store_true",
        help="Skip all LLM calls; return clean verdicts (for testing).",
    )
    p_ce.add_argument(
        "--model", default="claude-haiku-4-5-20251001",
        help="Model to use for judgment (default: claude-haiku-4-5-20251001).",
    )
    p_ce.set_defaults(func=cmd_claim_evidence)

    p_gd = sub.add_parser(
        "ghost-deck",
        help="Ghost-deck title-strip check: judge whether slide titles tell a "
             "coherent argument without body content. Exit 0 = pass/warn, 1 = fail, "
             "2 = plumbing.",
    )
    p_gd.add_argument(
        "plan",
        help="Path to plan.yaml (or a flat YAML list of titles).",
    )
    p_gd.add_argument(
        "-o", "--output", required=True,
        help="Output path for ghost_deck_report.md.",
    )
    p_gd.add_argument(
        "--offline", action="store_true",
        help="Skip all LLM calls; return a pass verdict (for CI without keys).",
    )
    p_gd.add_argument(
        "--model", default="claude-haiku-4-5-20251001",
        help="Model to use for judgment (default: claude-haiku-4-5-20251001).",
    )
    p_gd.set_defaults(func=cmd_ghost_deck)

    p_tl = sub.add_parser(
        "title-lint",
        help="Deterministic title rules: empty, too-long, no-verb, and-conjunction. "
             "Zero LLM. Exit 0 = clean, 1 = issues found, 2 = plumbing.",
    )
    p_tl.add_argument(
        "plan",
        help="Path to plan.yaml (or a flat YAML list of titles).",
    )
    p_tl.add_argument(
        "-o", "--output", required=True,
        help="Output path for title_lint_report.md.",
    )
    p_tl.add_argument(
        "--json", dest="emit_json", action="store_true",
        help="Emit issues as a JSON array to stdout instead of a markdown report.",
    )
    p_tl.set_defaults(func=cmd_title_lint)

    p_wf = sub.add_parser(
        "wireframe",
        help="Render a DSL layout as an annotated SVG wireframe.",
    )
    p_wf.add_argument("layout", help="Path to a .slide.dsl file")
    p_wf.add_argument("--brand", required=True, help="Brand id (dir name under brands/)")
    p_wf.add_argument("--content", help="YAML file with slot values (optional)")
    p_wf.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels even when --content is supplied. "
             "Forces skip_interpolation=True so the wireframe shows slot structure.",
    )
    p_wf.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wf.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file whose first slide is embedded as background. "
             "Requires LibreOffice on PATH.",
    )
    p_wf.add_argument(
        "--overlay-slide", type=int, default=0,
        help="0-based slide index to use from --overlay-pptx (default 0).",
    )
    p_wf.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide image (0.0–1.0, default 0.55).",
    )
    p_wf.set_defaults(func=cmd_wireframe)

    p_wfs = sub.add_parser(
        "wireframe-sheet",
        help="Render all slides in a plan as a SVG wireframe contact sheet.",
    )
    p_wfs.add_argument("plan", help="Path to the deck plan YAML")
    p_wfs.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wfs.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file; each slide is embedded behind its wireframe.",
    )
    p_wfs.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide images (0.0–1.0, default 0.55).",
    )
    p_wfs.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels in every cell — skip interpolation "
             "even when plans provide inline content or content_file. The contact "
             "sheet is intended as a layout-regression baseline, so showing slot "
             "structure is often more useful than filled content.",
    )
    p_wfs.set_defaults(func=cmd_wireframe_sheet)

    p_polish = sub.add_parser(
        "polish",
        help="Refurbish old PPTX diagrams into brand-perfect DSL artifacts.",
    )
    p_polish.add_argument("input", help=".pptx file to refurbish")
    p_polish.add_argument("-o", "--output", required=True, help="output .pptx path")
    p_polish.add_argument("--brand", default="feinschliff")
    p_polish.add_argument(
        "--refurbish-all",
        action="store_true",
        help="Auto-accept all refurbish proposals (emit DSL for every diagram slide).",
    )
    p_polish.add_argument(
        "--no-refurbish",
        action="store_true",
        help="Skip refurbish; just copy the input to the output path unchanged.",
    )
    p_polish.add_argument(
        "--refurbish-default",
        choices=("excalidraw", "svg"),
        default=None,
        help="Force a specific emitter instead of letting kind_selector choose."
        " Only applies to --mode redesign.",
    )
    p_polish.add_argument(
        "--mode",
        choices=["cosmetic", "redesign"],
        default="cosmetic",
        help=(
            "cosmetic (default): preserve slide count + content verbatim, fix brand"
            " chrome/typography only. redesign: extract diagram IR and rebuild"
            " diagram slides (legacy behaviour). --refurbish-* flags only apply to redesign."
        ),
    )
    p_polish.set_defaults(func=cmd_polish)

    p_book = sub.add_parser(
        "book",
        help="Render an annotated speaker-book PDF from a deck plan + "
             "design brief. Front matter (takeaway, audience, frame, "
             "red_line, hook) + one page per slide with the rendered "
             "thumbnail, claim, speaker notes, audience_fit, and role.",
    )
    p_book.add_argument("plan", help="Path to the deck plan YAML.")
    p_book.add_argument(
        "--design-brief", required=True,
        help="Path to design_brief.json (front matter + per-slide role / "
             "audience_fit come from here).",
    )
    p_book.add_argument(
        "--pptx", default=None,
        help="Path to the rendered .pptx (used to extract per-slide PNG "
             "thumbnails). If omitted, the book is rendered without "
             "thumbnails — useful for fast preview while authoring.",
    )
    p_book.add_argument(
        "-o", "--output", required=True,
        help="Output path for the speaker-book .pdf.",
    )
    p_book.set_defaults(func=cmd_book)

    # `deck plan` is a thin alias for `deck storyline` — same CLI surface,
    # same handler. The mode-level work (steps 0 → 1c only, no render) is
    # the skill orchestrator's job; the CLI part is just the storyline
    # report materialization helper. See
    # skills/deck/references/modes.md::plan for the mode semantics.
    p_plan = sub.add_parser(
        "plan",
        help="Emit a title-only storyline report (alias for `deck storyline`). "
             "Implements the CLI surface of the /deck plan mode — steps 0 → 1c "
             "only, no render. See skills/deck/references/modes.md::plan.",
    )
    p_plan.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_plan.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_plan.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_plan.set_defaults(func=cmd_storyline)

    # ── Pipeline timing logs + parallel plan authoring ──────────────────
    # (`log-event`, `timing`, `plan-skeleton`, `plan-merge`) — extracted
    # into deck_subcommands/plan_log.py to keep this file from growing.
    from feinschliff.cli.deck_subcommands import plan_log
    plan_log.register(sub)

    # ── Parallel-verify aspect helpers ──────────────────────────────────
    p_aspect = sub.add_parser(
        "verify-aspect",
        help="Run one focused verify aspect on a built deck. Each aspect "
             "is independently runnable; spawning multiple aspects in "
             "parallel + collating gives the parallel-verify path.",
    )
    p_aspect.add_argument(
        "aspect",
        choices=["bbox", "font", "narrative", "brand", "image", "content",
                 "notes-coherence"],
        help="bbox = bounding-box / overflow; font = legibility / role "
             "mismatches; narrative = SCQA / claim-title across titles; "
             "brand = color contrast / token discipline; image = picture "
             "slot fit / style consistency; content = title-body coherence; "
             "notes-coherence = per-slide speaker notes track the deck's "
             "red_line.",
    )
    p_aspect.add_argument("--plan", required=True,
                          help="Path to the deck plan.yaml.")
    p_aspect.add_argument("--pptx", default=None,
                          help="Path to the built .pptx (required for bbox/font/brand/image).")
    p_aspect.add_argument("--png-dir", default=None,
                          help="Directory with rendered slide-NN.png files.")
    p_aspect.add_argument("--design-brief", default=None,
                          help="Path to design_brief.json (used by narrative).")
    p_aspect.add_argument("-o", "--output", required=True,
                          help="Output path for verify-<aspect>.json.")
    p_aspect.set_defaults(func=cmd_verify_aspect)

    p_vs = sub.add_parser(
        "verify-static",
        help="Pre-render static geometry verify: detect slot-overflow and "
             "empty-placeholder defects from a plan.yaml without rendering. "
             "Exit 0 = clean, 1 = defects found, 2 = plumbing error.",
    )
    p_vs.add_argument("plan", help="Path to the deck plan YAML")
    p_vs.add_argument(
        "--json",
        action="store_true",
        help="Emit defects as a JSON array to stdout instead of the human "
             "readable format. Shape: [{slide_index, kind, severity, "
             "message, meta}, ...]",
    )
    p_vs.add_argument(
        "--brand", default=None,
        help="Override brand. Default: from plan.brand or 'feinschliff'.",
    )
    p_vs.set_defaults(func=cmd_verify_static)

    p_af = sub.add_parser(
        "apply-fixes",
        help="Apply mechanical fixes to a plan.yaml from a verify-static "
             "defects JSON.  Mutates plan in place (or writes -o out.yaml). "
             "Exit 0 = patches applied, 1 = no patches applied.",
    )
    p_af.add_argument("plan", help="Path to the deck plan YAML to fix.")
    p_af.add_argument(
        "--defects", required=True,
        help="Path to a defects JSON file: either a flat list of Defect "
             "dicts [{slide_index, kind, severity, message, meta}, ...] "
             "as emitted by `deck verify-static --json`, or a "
             '{"defects": {slide_idx: [...]}, ...} collated shape.',
    )
    p_af.add_argument(
        "-o", "--output",
        help="Output path for the fixed plan YAML.  When omitted the plan "
             "is updated in place.",
    )
    p_af.add_argument(
        "--brand", default=None,
        help="Override brand. Default: from plan.brand or 'feinschliff'.",
    )
    p_af.set_defaults(func=cmd_apply_fixes)

    p_collate = sub.add_parser(
        "verify-collate",
        help="Merge per-aspect verify outputs into a single verify_report.md.",
    )
    p_collate.add_argument(
        "--aspect", action="append", default=[], metavar="PATH",
        help="One per aspect: path to a verify-<aspect>.json file.",
    )
    p_collate.add_argument("--plan", required=True,
                           help="Path to the deck plan.yaml (for slide titles).")
    p_collate.add_argument("--iteration", type=int, default=1,
                           help="Iteration number (header field).")
    p_collate.add_argument("--budget", type=int, default=3,
                           help="Iteration budget (header field).")
    p_collate.add_argument("--png-dir", default=None,
                           help="PNG directory (header reference).")
    p_collate.add_argument("-o", "--output", required=True,
                           help="Output path for verify_report.md.")
    p_collate.set_defaults(func=cmd_verify_collate)


def _patch_set_hash(patches: list) -> str:
    """Delegate to feinschliff.deck.orchestrate.patch_set_hash."""
    return _patch_set_hash_fn(patches)


def cmd_intake_validate(args) -> int:
    """Validate a deck_brief.yaml against the intake schema."""
    from feinschliff.intake import load_brief

    path = Path(args.brief).resolve()
    if not path.is_file():
        print(f"deck: deck_brief not found: {path}", file=sys.stderr)
        return 2
    try:
        load_brief(path)
    except ValueError as e:
        print(str(e), file=sys.stderr)
        return 1
    print(f"deck: deck_brief OK ({path})")
    return 0


def cmd_commitment_validate(args) -> int:
    """Validate commitment.yaml; optionally cross-check arc alignment."""
    from feinschliff.storyline import (
        check_arc_alignment,
        load_all_arcs,
        load_commitment,
    )

    path = Path(args.commitment).resolve()
    if not path.is_file():
        print(f"deck: commitment not found: {path}", file=sys.stderr)
        return 2
    try:
        commitment = load_commitment(path)
    except ValueError as e:
        print(str(e), file=sys.stderr)
        return 1
    if args.check_arc:
        arcs = load_all_arcs()
        arc = arcs.get(commitment["deck_type"])
        if arc is None:
            print(
                f"deck: no arc schema for deck_type={commitment['deck_type']!r}",
                file=sys.stderr,
            )
            return 1
        misses = check_arc_alignment(commitment, arc)
        if misses:
            for m in misses:
                print(f"deck: arc: {m}", file=sys.stderr)
            return 1
    print(f"deck: commitment OK ({path})")
    return 0


def cmd_intake_skeleton(args) -> int:
    """Emit a deck_brief.yaml skeleton seeded by heuristic inference."""
    from feinschliff.intake import empty_brief, infer_from_text

    brief = empty_brief()
    if args.brief:
        brief_text = args.brief
        candidate = Path(args.brief)
        if candidate.is_file():
            brief_text = candidate.read_text()
        brief.update(infer_from_text(brief_text))
    out_yaml = yaml.safe_dump(brief, allow_unicode=True, sort_keys=False)
    if args.output:
        Path(args.output).write_text(out_yaml)
        print(f"deck: wrote skeleton to {args.output}", file=sys.stderr)
    else:
        sys.stdout.write(out_yaml)
    return 0


def cmd_build(args) -> int:
    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck: plan not found: {plan_path}", file=sys.stderr)
        return 2
    plan = yaml.safe_load(plan_path.read_text()) or {}

    # Surface the soft degradation: without feinschliff-builder, per-slide
    # speaker-notes validation is silently skipped. Print a one-line hint so
    # operators know what's missing instead of finding out via a downstream
    # surprise. Suppress with FEINSCHLIFF_QUIET_NOTES_BUDGET=1 for CI runs
    # that intentionally ship without builder.
    if validate_notes is None and not os.environ.get("FEINSCHLIFF_QUIET_NOTES_BUDGET"):
        print(
            "deck build: notes-budget validation skipped "
            "(feinschliff-builder not installed). "
            "Install feinschliff-builder to enable per-slide notes lint, "
            "or set FEINSCHLIFF_QUIET_NOTES_BUDGET=1 to silence this hint.",
            file=sys.stderr,
        )

    default_brand = plan.get("brand", "feinschliff")
    default_theme = getattr(args, "theme", None) or plan.get("theme")
    # Notes lint reads the deck-level verbosity (mirrors design_brief.verbosity)
    # to pick a per-slide word budget. Unset → budget check is skipped.
    plan_verbosity = plan.get("verbosity")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    # ── brand-chrome-leak warn (skeleton→plan layout swap audit) ─────────
    # When `plan.skeleton.yaml` lives next to `plan.yaml`, compare per-slide
    # layout choices. A slide whose skeleton layout was a brand-pack layout
    # (under `brands/<name>/layouts/`) but whose final plan layout is a
    # toolkit one is a "brand-chrome leak" — the content author overrode
    # the picker's brand pick with a generic. Loud one-line warning per
    # leak so build logs surface what was given up. Suppress with
    # `FEINSCHLIFF_QUIET_BRAND_LEAK=1`. No fatal effect — the operator may
    # have a real reason to swap, but they should see it.
    if not os.environ.get("FEINSCHLIFF_QUIET_BRAND_LEAK"):
        sk_path = plan_path.with_name("plan.skeleton.yaml")
        if sk_path.is_file():
            try:
                sk_plan = yaml.safe_load(sk_path.read_text()) or {}
                sk_slides = sk_plan.get("slides") or []
                leaks: list[str] = []
                for idx, (sk_s, pl_s) in enumerate(zip(sk_slides, slides_spec), start=1):
                    sk_lay = str(sk_s.get("layout", ""))
                    pl_lay = str(pl_s.get("layout", ""))
                    if not sk_lay or not pl_lay or sk_lay == pl_lay:
                        continue
                    sk_brand = "/brands/" in sk_lay
                    pl_brand = "/brands/" in pl_lay
                    if sk_brand and not pl_brand:
                        leaks.append(
                            f"  slide {idx}: brand layout "
                            f"'{Path(sk_lay).stem}' → toolkit "
                            f"'{Path(pl_lay).stem}'"
                        )
                if leaks:
                    print(
                        f"deck build: brand-chrome-leak — "
                        f"{len(leaks)} slide(s) overrode brand picks with "
                        f"toolkit layouts:\n" + "\n".join(leaks)
                        + "\n  See deck-skill 'When overriding the picked "
                        "layout' guidance; suppress with "
                        "FEINSCHLIFF_QUIET_BRAND_LEAK=1.",
                        file=sys.stderr,
                    )
            except (yaml.YAMLError, OSError):
                # Skeleton unreadable / malformed → silent skip (a build
                # without skeleton context is the normal pre-fan-out flow).
                pass

    # Resolve build-time image provider from the default brand's
    # `$image_provider` (extends-resolved). Per-slide `brand:` overrides
    # still drive tokens/compounds, but the provider used to resolve
    # `picture query:` nodes is deck-wide — `asset_lock.json` lives next
    # to the output deck. Absent → provider is None; brands using only
    # `picture path:` build as before. A typo'd kind raises KeyError with
    # a registry listing — surfaces as the normal CLI traceback.
    discover_providers()
    provider = None
    chain = None
    try:
        default_brand_obj = find_brand(default_brand)
    except ValueError as e:
        print(f"deck: {e}", file=sys.stderr)
        return 2
    if not getattr(args, "no_image_provider", False):
        # New `$image_providers` (plural list) → ProviderChain.
        # Legacy `$image_provider` (singular dict) stays on the original
        # `get_provider` path — that one resolves arbitrary kinds via the
        # registry, including test-only kinds the chain doesn't whitelist.
        _raw_tokens = getattr(default_brand_obj, "tokens", {}) or {}
        _ip_list = _raw_tokens.get("$image_providers")
        if _ip_list is not None:
            chain = chain_from_brand_config(_ip_list, brand_root=default_brand_obj.root)
        elif default_brand_obj.image_provider_config:
            cfg = default_brand_obj.image_provider_config
            provider = get_provider(cfg["kind"], cfg.get("config"))

    # ── Pre-render static geometry verify (--strict-static) ──────────────
    if getattr(args, "strict_static", False):
        _require_or_delegate_builder("deck build --strict-static")
        from feinschliff_builder.verify.static import validate as _validate_static
        _static_bag = _validate_static(
            plan, brand=default_brand_obj, plan_dir=plan_path.parent
        )
        if _static_bag:
            for _d in _static_bag:
                _loc = _d.location or "slide ?"
                print(
                    f"deck build: static: {_loc}: "
                    f"[{_d.severity.value.upper()}] {_d.kind.value} — {_d.message}",
                    file=sys.stderr,
                )
            print(
                f"deck build: --strict-static: {len(_static_bag)} static "
                f"defect(s) found. Fix them or remove --strict-static to skip "
                f"this gate.",
                file=sys.stderr,
            )
            return 1

    # ── Auto-fix loop (--autofix) ─────────────────────────────────────────
    if getattr(args, "autofix", False):
        _require_or_delegate_builder("deck build --autofix")
        from feinschliff_builder.verify.static import validate as _validate_static_af
        from feinschliff_builder.verify.autofix import plan_fixes, apply_fixes, diff_summary

        _MAX_AUTOFIX_CYCLES = 3
        _total_patches = 0
        _seen_hashes: set[str] = set()
        _oscillation_detected = False
        for _cycle in range(_MAX_AUTOFIX_CYCLES):
            _static_bag = _validate_static_af(
                plan, brand=default_brand_obj, plan_dir=plan_path.parent
            )
            if not _static_bag:
                break
            _patches = plan_fixes(_static_bag, plan, default_brand_obj.root)
            if not _patches:
                # No mechanical fix available; leave residuals for compile.
                break
            _h = _patch_set_hash(_patches)
            if _h in _seen_hashes:
                print(
                    f"deck build: autofix cycle {_cycle + 1}: identical patch set "
                    f"seen before; halting to avoid oscillation",
                    file=sys.stderr,
                )
                _oscillation_detected = True
                break
            _seen_hashes.add(_h)
            _before = plan
            plan = apply_fixes(plan, _patches)
            _total_patches += len(_patches)
            _summary = diff_summary(_before, plan)
            print(
                f"deck build: autofix cycle {_cycle + 1}: "
                f"{len(_patches)} patch(es) applied",
            )
            if _summary:
                for _line in _summary.splitlines():
                    print(f"  {_line}")
        else:
            # Exhausted cycles — check if residuals remain.
            _residuals = _validate_static_af(
                plan, brand=default_brand_obj, plan_dir=plan_path.parent
            )
            if _residuals:
                print(
                    f"deck build: --autofix: {len(_residuals)} residual static "
                    f"defect(s) after {_MAX_AUTOFIX_CYCLES} cycle(s) — proceeding "
                    f"to compile (orchestrator may revise).",
                    file=sys.stderr,
                )
        if _total_patches > 0:
            # Write the auto-fixed plan back to disk before compile.
            plan_path.write_text(
                yaml.safe_dump(plan, sort_keys=False, allow_unicode=True),
                encoding="utf-8",
            )
            print(
                f"deck build: auto-fix passes: {_total_patches} total patch(es); "
                f"plan written back to {plan_path}"
            )
        # Re-capture slides_spec from the (potentially mutated) plan so the
        # compile loop below uses the fixed content, not the pre-fix snapshot.
        slides_spec = plan.get("slides") or []

    # ── Default static gate — fatal static defects block EVERY build ─────
    # Run the same pre-render static verifier --strict-static uses, but gate
    # the abort on fatal_kinds(): FATAL static defects (slot-overflow) abort
    # the build; WARN-only results (empty-placeholder) never block a default
    # build — promoting those to errors stays --strict-static's job. Placed
    # after --autofix so the gate judges the fixed plan, and skipped under
    # --strict-static (which already aborted on ANY defect above). When
    # feinschliff-builder is not installed the import misses and the gate
    # degrades to a no-op — same optional-import pattern as pipeline.py's
    # structural validators (never crash, never delegate for a default build).
    _validate_static_gate = None
    # --slot-debug-color builds are diagnostic renders (every slot bound to
    # its own default to visualise coverage) — slot-overflow there reflects
    # the pack's showcase copy, not deck content. Skip the fatal gate.
    if getattr(args, "slot_debug_color", None):
        pass
    elif not getattr(args, "strict_static", False):
        try:
            from feinschliff_builder.verify.static import (  # type: ignore[import]
                validate as _validate_static_gate,
            )
        except ImportError:
            _validate_static_gate = None  # builder absent → gate is a no-op
    if _validate_static_gate is not None:
        _gate_bag = _validate_static_gate(
            plan, brand=default_brand_obj, plan_dir=plan_path.parent
        )
        _gate_fatal = [d for d in _gate_bag if d.kind.value in fatal_kinds()]
        if _gate_fatal:
            for _d in _gate_fatal:
                _loc = _d.location or "slide ?"
                print(
                    f"deck build: static: {_loc}: "
                    f"[{_d.severity.value.upper()}] {_d.kind.value} — {_d.message}",
                    file=sys.stderr,
                )
            _autofix_hint = (
                ""
                if getattr(args, "autofix", False)
                else " or run --autofix"
            )
            print(
                f"deck build: aborting — {len(_gate_fatal)} fatal static "
                f"defect(s) found pre-render. Shorten the flagged content, "
                f"enable autoshrink on the layout{_autofix_hint}.",
                file=sys.stderr,
            )
            return 1

    # Compute the output deck path up front so it can serve as `deck_dir`
    # for `asset_lock.json` + `.cache/` during the build.
    out_path = Path(args.output or plan.get("out", "deck.pptx")).resolve()

    plan_dir = plan_path.parent
    slides_payload: list[tuple[list, object, Path]] = []
    compile_jobs: list[tuple[int, object, Path, dict]] = []
    content_defects_by_slide: dict[int, list] | None = (
        {} if not args.skip_content_lint else None
    )
    all_diagram_defects: list = []

    # Top-level timing phase for the build. Per-slide events would push
    # indentation past sanity here; instead we emit one `build:total`
    # event manually around the whole compile loop (see end of function)
    # and emit `build:slide` lightweight events inline via log_event.
    import time as _time
    _build_t0 = _time.perf_counter()
    log_event(plan_dir, "build:total", "start",
              slides=len(slides_spec), brand=default_brand)

    with tempfile.TemporaryDirectory() as _tmp:
        diagrams_out = Path(_tmp) / "diagrams"
        diagrams_out.mkdir()

        for i, spec in enumerate(slides_spec):
            layout_path = (plan_dir / spec["layout"]).resolve()
            if not layout_path.is_file():
                # also try toolkit-relative (plugin root / rel).
                alt = _find_toolkit_file(spec["layout"])
                if alt is not None:
                    layout_path = alt
                else:
                    print(f"deck: slide {i}: layout not found: {spec['layout']}", file=sys.stderr)
                    return 2

            brand = spec.get("brand", default_brand)
            # Per-slide theme: explicit slide `theme:` key > CLI --theme / plan theme:
            slide_theme = spec.get("theme", default_theme)
            try:
                brand_dir = find_brand(brand).root
            except ValueError as e:
                print(f"deck: slide {i}: {e}", file=sys.stderr)
                return 2

            try:
                tokens = load_tokens_with_theme(brand_dir, slide_theme)
            except ValueError as e:
                print(f"deck: slide {i}: {e}", file=sys.stderr)
                return 2
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=_bundled_compounds()
            )

            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd

            ctx = spec.get("content") or {}
            if not ctx and "content_file" in spec:
                content_path = (plan_dir / spec["content_file"]).resolve()
                if not content_path.is_file():
                    alt = _find_toolkit_file(spec["content_file"])
                    if alt is not None:
                        content_path = alt
                    else:
                        print(f"deck: slide {i}: content_file not found: {spec['content_file']}", file=sys.stderr)
                        return 2
                ctx = yaml.safe_load(content_path.read_text()) or {}
            if getattr(args, "slot_debug_color", None):
                ctx = dict(ctx)
                ctx["_slot_debug_color"] = args.slot_debug_color

            # Brand-layout slot metadata: auto-bind footer / page-number
            # slots (from deck-level `vars:` / the slide index) and derive
            # provider queries for unbound `class: replace` image slots.
            # Explicit plan bindings — including an explicit "" — always win.
            ctx = auto_bind_slots(
                ctx,
                layout_path=layout_path,
                layout_nodes=layout_nodes,
                slide_index=i + 1,
                deck_vars=plan.get("vars"),
                image_provider_available=provider is not None,
            )

            if content_defects_by_slide is not None:
                slide_index = i + 1
                # Strip `.slide.dsl` (and any leading path) to get the bare
                # layout name — this is what the structural validators key on.
                layout_name = layout_path.name
                if layout_name.endswith(".slide.dsl"):
                    layout_name = layout_name[: -len(".slide.dsl")]
                warn_overbudget_slots(ctx, layout_path=layout_path, slide_index=slide_index)
                slot_budgets = compute_slot_budgets(layout_nodes, tokens, compounds=compounds)
                chrome_bboxes: list = []
                try:
                    fm_text, _ = split_frontmatter(layout_path.read_text(encoding="utf-8"))
                    if fm_text:
                        chrome_bboxes = (yaml.safe_load(fm_text) or {}).get("chrome_bboxes") or []
                except Exception as exc:
                    print(f"deck build: WARN: could not read chrome_bboxes from "
                          f"{layout_path.name}: {exc}", file=sys.stderr)
                    chrome_bboxes = []
                slide_defects = validate_content(
                    ctx, slide_index=slide_index, layout=layout_name,
                    slot_budgets=slot_budgets, chrome_bboxes=chrome_bboxes,
                )
                if validate_notes is not None:
                    slide_defects.extend(validate_notes(
                        spec.get("notes"),
                        slide_index=slide_index,
                        is_hook=(i == 0),
                        verbosity=plan_verbosity,
                    ))
                warn_defects, fatal_defects = [], []
                for d in slide_defects:
                    (warn_defects if getattr(d, "severity", "fatal") == "warn"
                     else fatal_defects).append(d)
                for d in warn_defects:
                    print(f"deck build: WARN: {d}", file=sys.stderr)
                if fatal_defects:
                    content_defects_by_slide[slide_index] = fatal_defects

            compile_jobs.append((
                i, spec.get("notes"), brand_dir,
                dict(layout_path=layout_path, ctx=ctx, brand_dir=brand_dir,
                     slide_index=i + 1, diagrams_out_dir=diagrams_out,
                     theme=slide_theme),
            ))

        if content_defects_by_slide:
            emit_defects_and_abort_message(content_defects_by_slide, cli_name="deck build")
            return 1

        # Compile every slide — DSL expand + diagram rendering. Slides are
        # independent (the diagram cache keys on slide_index, so output
        # files are disjoint); with --workers > 1 they fan out across
        # processes and only the PPTX assembly below stays serial.
        _workers = getattr(args, "workers", 1)
        if _workers == 0:
            import os as _os
            _workers = min(8, max(1, (_os.cpu_count() or 2) // 2))
        _workers = min(_workers, max(1, len(compile_jobs)))

        def _join(i: int, notes, brand_dir: Path, slide_result) -> None:
            for d in slide_result.defects:
                print(f"deck: slide {i + 1}: {format_defect(d)}", file=sys.stderr)
            all_diagram_defects.extend(slide_result.defects)
            if notes is not None:
                slides_payload.append(
                    (slide_result.primitives, slide_result.tokens,
                     brand_dir / "assets", notes)
                )
            else:
                slides_payload.append(
                    (slide_result.primitives, slide_result.tokens,
                     brand_dir / "assets")
                )

        if _workers <= 1:
            for i, notes, brand_dir, kwargs in compile_jobs:
                _slide_t0 = _time.perf_counter()
                log_event(plan_dir, "build:slide", "start", slide=i + 1,
                          layout=kwargs["layout_path"].name)
                slide_result = compile_slide(**kwargs)
                log_event(
                    plan_dir, "build:slide", "end", slide=i + 1,
                    layout=kwargs["layout_path"].name,
                    elapsed_ms=int((_time.perf_counter() - _slide_t0) * 1000),
                    defects=len(slide_result.defects),
                )
                _join(i, notes, brand_dir, slide_result)
        else:
            from concurrent.futures import ProcessPoolExecutor
            log_event(plan_dir, "build:compile-pool", "start",
                      slides=len(compile_jobs), workers=_workers)
            _pool_t0 = _time.perf_counter()
            with ProcessPoolExecutor(max_workers=_workers) as _pool:
                _futs = [(i, notes, brand_dir,
                          _pool.submit(compile_slide, **kwargs))
                         for i, notes, brand_dir, kwargs in compile_jobs]
                # Join in submission order — defect output and the slide
                # payload sequence stay identical to a sequential build.
                for i, notes, brand_dir, fut in _futs:
                    _join(i, notes, brand_dir, fut.result())
            log_event(plan_dir, "build:compile-pool", "end",
                      slides=len(compile_jobs), workers=_workers,
                      elapsed_ms=int((_time.perf_counter() - _pool_t0) * 1000))

        allowed_to_skip: set[str] = set()
        if getattr(args, "allow_diagram_warnings", False):
            allowed_to_skip |= {"diagram-overflow", "diagram-text-too-small"}
        _fatal_diagram = [
            d for d in all_diagram_defects
            if d.kind.value in fatal_kinds() and d.kind.value not in allowed_to_skip
        ]
        if _fatal_diagram:
            print(
                f"deck build: aborting — {len(_fatal_diagram)} fatal defect(s) "
                f"across {len(slides_spec)} slide(s). Pass "
                f"--allow-diagram-warnings to demote "
                f"diagram-overflow/diagram-text-too-small (if those are the only blockers).",
                file=sys.stderr,
            )
            return 1

        prs = build_multi_slide(
            slides_payload,
            asset_root_fallback=_bundled_assets(),
            image_provider=provider,
            provider_chain=chain,
            deck_dir=out_path.parent,
            slide_numbers=not getattr(args, "no_slide_numbers", False),
        )
        missing = getattr(prs, "missing_assets", []) or []
        # --allow-missing-assets is suppressed when a provider chain is
        # configured and failed — operator must fix chain, not ship blank deck.
        chain_failed = chain is not None and any(
            e.get("kind") == "chain-miss" for e in missing
        )
        if missing and (chain_failed or not getattr(args, "allow_missing_assets", False)):
            for entry in missing:
                kind = entry.get("kind", "missing")
                path = entry.get("path") or "(unset)"
                slide_n = entry.get("slide_index", "?")
                line = entry.get("line_no", "?")
                print(
                    f"deck build: slide {slide_n}: missing asset ({kind}) "
                    f"at line {line}: {path}",
                    file=sys.stderr,
                )
            print(
                f"deck build: aborting — {len(missing)} missing required "
                f"asset(s) across {len(slides_spec)} slide(s). Mark optional "
                f"slots with `optional:true` or pass --allow-missing-assets.",
                file=sys.stderr,
            )
            return 1
        if getattr(args, "embed_fonts", False) and slides_payload:
            from feinschliff.dsl.font_embed import embed_brand_fonts

            # Per-slide `brand:` overrides exist, but font embedding is
            # deck-wide — the first slide's brand tokens win.
            embedded = embed_brand_fonts(prs, slides_payload[0][1])
            if embedded:
                print(f"embedded fonts: {', '.join(embedded)}", file=sys.stderr)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        print(f"wrote {out_path} ({len(prs.slides)} slides)")
        log_event(
            plan_dir, "build:total", "end",
            elapsed_ms=int((_time.perf_counter() - _build_t0) * 1000),
            slides=len(prs.slides), output=str(out_path),
        )

        # ── Post-build visual metrics (--strict-visual) ───────────────────
        if getattr(args, "strict_visual", False):
            import tempfile as _tempfile
            from feinschliff.quality import compute_visual_metrics
            from feinschliff.quality.visual_metrics_report import (
                write_visual_metrics_report,
            )
            _png_map: dict[int, Path] = {}
            try:
                from feinschliff_builder.verify.render_pngs import (
                    render_slides_to_png,
                )
                with _tempfile.TemporaryDirectory() as _tmp:
                    _png_map = render_slides_to_png(out_path, Path(_tmp))
                    if not _png_map:
                        print(
                            "deck build: --strict-visual: PNG render returned "
                            "no slides (soffice/pdftoppm unavailable?); "
                            "skipping visual metrics.",
                            file=sys.stderr,
                        )
                    else:
                        _vm_result = compute_visual_metrics(_png_map)
                        _report_path = out_path.parent / "visual_metrics_report.md"
                        write_visual_metrics_report(_vm_result, _report_path)
                        print(
                            f"visual metrics: {_vm_result.verdict} "
                            f"({len(_vm_result.issues)} issue(s)); "
                            f"report → {_report_path}"
                        )
                        if _vm_result.issues:
                            for _vi in _vm_result.issues:
                                print(
                                    f"  slide {_vi.slide} [{_vi.metric}/"
                                    f"{_vi.severity}]: {_vi.message}",
                                    file=sys.stderr,
                                )
            except ImportError:
                print(
                    "deck build: --strict-visual: feinschliff-builder not "
                    "installed; skipping visual metrics.",
                    file=sys.stderr,
                )
            except Exception as _ve:
                print(
                    f"deck build: --strict-visual: visual metrics failed "
                    f"({_ve}); skipping.",
                    file=sys.stderr,
                )

        # ── Post-build Knaflic craft-rules check (--strict-craft) ────────
        if getattr(args, "strict_craft", False):
            from feinschliff.quality import check_craft_rules, write_craft_report
            _craft_slides = [
                {
                    "layout": str(
                        (plan_dir / spec["layout"]).resolve()
                        if not Path(spec["layout"]).is_absolute()
                        else Path(spec["layout"])
                    ),
                    "content_inline": spec.get("content") or {},
                }
                for spec in slides_spec
            ]
            _craft_report = check_craft_rules(_craft_slides)
            _craft_report_path = out_path.parent / "craft_report.md"
            write_craft_report(_craft_report, _craft_report_path)
            print(
                f"craft rules: {_craft_report.verdict} "
                f"({len(_craft_report.issues)} issue(s)); "
                f"report → {_craft_report_path}"
            )
            for _ci in _craft_report.issues:
                print(
                    f"  slide {_ci.slide} [{_ci.rule}/{_ci.severity}]: "
                    f"{_ci.message}",
                    file=sys.stderr,
                )
            if _craft_report.verdict == "fail":
                print(
                    "deck build: --strict-craft: craft-rules verdict is 'fail'. "
                    "Fix the flagged slides or remove --strict-craft to skip.",
                    file=sys.stderr,
                )
                return 1

        return 0


def cmd_pick(args) -> int:
    from feinschliff.deck.picker import LayoutPicker

    if args.signals == "-":
        raw = sys.stdin.read()
    else:
        raw = Path(args.signals).read_text()
    signals = yaml.safe_load(raw) or {}

    picker = LayoutPicker(top_k=args.top_k)
    candidates = picker.candidates(signals)
    if not candidates:
        print("deck: no candidate layouts matched the signals", file=sys.stderr)
        return 1
    for c in candidates:
        print(f"{c.score:5.2f}  {c.layout_name:<24}  {c.reason}")
    return 0


def cmd_pick_deck(args) -> int:
    """Run the arc-aware deck-level picker; emit picker_report.json."""
    from feinschliff.picker import pick_deck
    from feinschliff.picker.report import write_picker_report

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck pick-deck: plan not found: {plan_path}", file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    brand = plan.get("brand", "feinschliff")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck pick-deck: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    deck_brief: dict | None = None
    if args.deck_brief:
        brief_path = Path(args.deck_brief).resolve()
        if not brief_path.is_file():
            print(
                f"deck pick-deck: --deck-brief not found: {brief_path}",
                file=sys.stderr,
            )
            return 2
        import json as _json
        raw = brief_path.read_text()
        try:
            deck_brief = _json.loads(raw)
        except _json.JSONDecodeError:
            deck_brief = yaml.safe_load(raw) or {}

    report = pick_deck(
        slides_spec,
        brand=brand,
        deck_brief=deck_brief,
        top_k=args.top_k,
    )

    out_path = (
        Path(args.output).resolve()
        if args.output
        else plan_path.parent / "picker_report.json"
    )
    write_picker_report(report, out_path)

    for w in report.arc_warnings:
        print(f"arc warning: {w}", file=sys.stderr)

    print(
        f"picker_report: {len(report.picks)} slide(s), "
        f"{len(report.arc_warnings)} arc warning(s) → {out_path}"
    )
    return 0



def cmd_wireframe(args) -> int:
    _require_or_delegate_builder("deck wireframe")
    from feinschliff_builder.decompile.wireframe import render_wireframe

    layout_path = Path(args.layout).resolve()
    if not layout_path.is_file():
        alt = _find_toolkit_file(args.layout)
        if alt is not None:
            layout_path = alt
        else:
            print(f"deck wireframe: layout not found: {args.layout}", file=sys.stderr)
            return 2

    if args.content:
        content_path = Path(args.content).resolve()
        if not content_path.is_file():
            print(f"deck wireframe: content file not found: {args.content}", file=sys.stderr)
            return 2
    else:
        content_path = None
    # Skip interpolation when no content is given (preserves slot labels) or
    # when --show-slots forces slot-structure mode even with content provided.
    skip_interp = (content_path is None) or args.show_slots
    try:
        primitives, tokens = _build_primitives_for_layout(
            layout_path, args.brand, content_path,
            skip_interpolation=skip_interp,
        )
    except ValueError as exc:
        print(f"deck wireframe: {exc}", file=sys.stderr)
        return 2

    bg_b64: str | None = None
    if args.overlay_pptx:
        from feinschliff.io.pptx_to_png import slide_to_b64
        try:
            bg_b64 = slide_to_b64(
                Path(args.overlay_pptx).resolve(),
                slide_index=args.overlay_slide,
            )
        except (FileNotFoundError, RuntimeError, IndexError) as exc:
            print(f"deck wireframe: overlay failed — {exc}", file=sys.stderr)
            return 2

    title = layout_path.name.replace(".slide.dsl", "")
    svg = render_wireframe(
        primitives, tokens,
        title=title,
        background_png_b64=bg_b64,
        background_opacity=args.overlay_opacity,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay" if bg_b64 else "wireframe"
    print(f"wrote {out} ({mode}, {len(primitives)} primitives)")
    return 0


def cmd_wireframe_sheet(args) -> int:
    _require_or_delegate_builder("deck wireframe-sheet")
    from feinschliff_builder.decompile.wireframe import render_wireframe_sheet

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck wireframe-sheet: plan not found: {plan_path}", file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    default_brand = plan.get("brand", "feinschliff")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck wireframe-sheet: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    plan_dir = plan_path.parent
    slides_data: list[tuple[list, object, str]] = []

    for i, spec in enumerate(slides_spec):
        layout_rel = spec["layout"]
        layout_path = (plan_dir / layout_rel).resolve()
        if not layout_path.is_file():
            alt = _find_toolkit_file(layout_rel)
            if alt is not None:
                layout_path = alt
            else:
                print(f"deck wireframe-sheet: slide {i}: layout not found: {layout_rel}",
                      file=sys.stderr)
                return 2

        brand = spec.get("brand", default_brand)
        ctx_inline = spec.get("content") or {}
        content_path: Path | None = None
        if not ctx_inline and "content_file" in spec:
            cp = (plan_dir / spec["content_file"]).resolve()
            if not cp.is_file():
                cp = _find_toolkit_file(spec["content_file"]) or cp
            content_path = cp if cp.is_file() else None

        try:
            brand_dir = find_brand(brand).root
            tokens = load_tokens(brand_dir)
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=_bundled_compounds()
            )
            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd
            ctx: dict = ctx_inline.copy()
            if not ctx and content_path:
                ctx = yaml.safe_load(content_path.read_text()) or {}
            if args.show_slots:
                # Skip interpolation so {{ slot_name }} labels survive into the cell.
                primitives, _ = expand_compounds(layout_nodes, compounds)
            else:
                interp = interpolate_nodes(layout_nodes, ctx)
                primitives, _ = expand_compounds(interp, compounds)
        except (ValueError, OSError, yaml.YAMLError, KeyError) as exc:
            print(f"deck wireframe-sheet: slide {i}: {exc}", file=sys.stderr)
            return 1

        title = layout_path.name.replace(".slide.dsl", "")
        slides_data.append((primitives, tokens, title))

    bg_list: list[str | None] | None = None
    if args.overlay_pptx:
        from feinschliff.io.pptx_to_png import pptx_to_pngs_b64
        try:
            pngs = pptx_to_pngs_b64(Path(args.overlay_pptx).resolve())
        except (FileNotFoundError, RuntimeError) as exc:
            print(f"deck wireframe-sheet: overlay failed — {exc}", file=sys.stderr)
            return 2
        if len(pngs) != len(slides_data):
            if len(pngs) < len(slides_data):
                detail = (
                    f"only {len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "unmatched slides will have no overlay. "
                    "(Some LibreOffice versions only export the first slide.)"
                )
            else:
                detail = (
                    f"{len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "extra overlays will be ignored. Check that --overlay-pptx "
                    "matches the deck plan."
                )
            print(f"deck wireframe-sheet: warning — {detail}", file=sys.stderr)
        # Pad to match slide count; truncate is implicit via index guard in render.
        bg_list = list(pngs) + [None] * max(0, len(slides_data) - len(pngs))

    svg = render_wireframe_sheet(
        slides_data,
        background_pngs_b64=bg_list,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay sheet" if bg_list else "wireframe sheet"
    print(f"wrote {out} ({mode}, {len(slides_data)} slide(s))")
    return 0


def cmd_polish(args) -> int:
    """Walk a PPTX, extract diagram IR from each slide, emit DSL artifacts,
    and rebuild a brand-perfect polished deck.

    WIRED:
    - Reads input .pptx with python-pptx.
    - Calls extract_vector.extract_from_slide() for every slide.
    - Calls kind_selector.select_kind() (or uses --refurbish-default) to pick
      the emitter target.
    - Emits .exc.dsl (excalidraw) or .svg.dsl (svg) into <output-dir>/refurbished/.
    - Writes refurbish_report.md next to the output file.
    - Rebuilds a polished PPTX from the refurbished diagram slides using the
      excalidraw-diagram or svg-infographic layout with the brand token pack.

    NOTE: Non-diagram slides (no detectable nodes) are skipped in the rebuilt
    deck. The output PPTX contains only refurbished diagram slides. If no
    diagram slides are found, the input is copied to the output as a fallback.

    --refurbish-all is the only interactive-mode variant implemented here;
    interactive per-slide confirmation is not yet wired.
    """
    _require_or_delegate_builder("deck polish")
    import shutil

    from pathlib import Path as _Path
    from pptx import Presentation
    from feinschliff_builder.diagrams.refurbish.extract_vector import extract_from_slide
    from feinschliff_builder.diagrams.refurbish.kind_selector import select_kind
    from feinschliff_builder.diagrams.refurbish.emit_excalidraw import emit as emit_excalidraw_dsl
    from feinschliff_builder.diagrams.refurbish.emit_svg import emit as emit_svg_dsl

    def _extract_slide_title(slide) -> str:
        """Pick the best title candidate from a refurbished slide.

        Order: (1) PPTX title placeholder if it has text; (2) the largest
        free-floating textbox (i.e., one whose shape isn't an auto-shape with
        its own label). Never returns a shape-label as the slide title.
        """
        # 1. Slide title placeholder
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame and title_shape.text_frame.text.strip():
                return title_shape.text_frame.text.strip().split("\n")[0]
        except (AttributeError, ValueError):
            pass
        # 2. Largest free-floating textbox (skip shapes with their own label)
        best = ""
        best_area = 0
        for sh in slide.shapes:
            try:
                _ = sh.auto_shape_type  # raises if not an auto-shape
                continue  # skip boxes/ellipses (they're nodes, not titles)
            except (ValueError, AttributeError):
                pass
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                area = (sh.width or 0) * (sh.height or 0)
                if area > best_area:
                    best = sh.text_frame.text.strip().split("\n")[0]
                    best_area = area
        return best

    src_path = _Path(args.input).resolve()
    if not src_path.is_file():
        print(f"deck polish: input not found: {src_path}", file=sys.stderr)
        return 2

    out_path = _Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if args.no_refurbish:
        shutil.copy(src_path, out_path)
        print(f"deck polish: --no-refurbish — copied {src_path.name} → {out_path}")
        return 0

    if getattr(args, "mode", "redesign") == "cosmetic":
        import subprocess
        from feinschliff.polish import cosmetic_polish

        report = cosmetic_polish(src_path, args.brand, out_path)
        for w in report.warnings:
            print(f"deck polish: warning: {w}", file=sys.stderr)
        print(
            f"deck polish: cosmetic plan written → {report.plan_path}"
            f" ({report.slides_preserved} slide(s))"
        )
        result = subprocess.run(
            ["feinschliff", "deck", "build", str(report.plan_path)],
            check=False,
        )
        return result.returncode

    refurbish_dir = out_path.parent / "refurbished"
    refurbish_dir.mkdir(exist_ok=True)
    report_lines = ["# Refurbish Report\n"]

    refurbished_slides: list[dict] = []  # each: {layout, content}

    in_pres = Presentation(str(src_path))
    for idx, slide in enumerate(in_pres.slides, start=1):
        ir = extract_from_slide(slide)
        if not ir.nodes:
            report_lines.append(f"- slide {idx}: no diagram nodes detected — skipped")
            continue

        kind = args.refurbish_default or select_kind(ir)

        if args.refurbish_all or args.refurbish_default:
            if kind == "excalidraw":
                dsl = emit_excalidraw_dsl(ir, 1720, 480)
                ext = ".exc.dsl"
                layout_name = "excalidraw-diagram.slide.dsl"
            else:
                dsl = emit_svg_dsl(ir, 1720, 480)
                ext = ".svg.dsl"
                layout_name = "svg-infographic.slide.dsl"
            artifact = refurbish_dir / f"slide-{idx}{ext}"
            artifact.write_text(dsl, encoding="utf-8")
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) → {artifact.name}"
            )

            # Best-effort: pull a title from the source slide.
            # Prefer (1) the slide's title placeholder, (2) the largest free-
            # floating textbox that isn't a shape label. Never use a box's own
            # label as the slide title — that's noise.
            title = _extract_slide_title(slide)

            # Strip the leading "canvas WxH\n" line — the layout template
            # embeds diagram_dsl inside its own canvas block.
            dsl_body = dsl.partition("\n")[2]

            refurbished_slides.append({
                "layout": layout_name,
                "content": {
                    "pgmeta": f"Slide {idx}",
                    "tracker": "Refurbished",
                    "action_title": title or f"Diagram {idx}",
                    "so_what": "",
                    "source": f"Refurbished from input slide {idx}",
                    "diagram_dsl": dsl_body,
                    "footer_left": "Feinschliff",
                    "footer_right": f"Slide {idx}",
                },
            })
        else:
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) — no flag set, skipped"
            )

    (out_path.parent / "refurbish_report.md").write_text("\n".join(report_lines), encoding="utf-8")

    if not refurbished_slides:
        # Nothing extracted — fall back to copying input
        shutil.copy(src_path, out_path)
        slide_count = len(in_pres.slides)
        print(
            f"deck polish: processed {slide_count} slide(s), "
            f"0 diagram slides found — input copied to {out_path.name}"
        )
        return 0

    _build_refurbished_deck(
        refurbished_slides,
        brand=args.brand,
        out_path=out_path,
    )

    slide_count = len(in_pres.slides)
    artifact_count = len(refurbished_slides)
    print(
        f"deck polish: processed {slide_count} slide(s), "
        f"{artifact_count} refurbished diagram slide(s) → {out_path.name}"
    )
    return 0


def _build_refurbished_deck(slides_plan: list[dict], brand: str, out_path: Path) -> None:
    """Build refurbished deck — DSL path removed."""
    raise NotImplementedError(
        "deck polish build is not available without the DSL pipeline. "
        "Use 'feinschliff render-master' for the master-template path."
    )


def cmd_book(args) -> int:
    """`feinschliff deck book` — annotated speaker-book PDF.

    Loads the deck plan + design brief, optionally renders per-slide
    PNG thumbnails from the built .pptx (via the existing
    `feinschliff_builder.verify.render_pngs.render_slides_to_png` helper), and writes a
    multi-page PDF: front matter page + one page per slide.

    The brief is the source of truth for deck-level fields (takeaway,
    audience, frame, …) and per-slide role / audience_fit. Speaker
    notes come from the plan's per-slide `notes:` field (preferred);
    when absent, the brief's per-slide `notes` is used as fallback.
    """
    _require_or_delegate_builder("deck book")
    import json as _json
    import tempfile as _tempfile

    from feinschliff.book import (
        BookSlide, DeckFrontMatter, compose_book_pdf,
    )
    from feinschliff_builder.verify.render_pngs import render_slides_to_png

    plan_path = Path(args.plan).resolve()
    brief_path = Path(args.design_brief).resolve()
    out_path = Path(args.output).resolve()

    if not plan_path.is_file():
        print(f"deck book: plan not found: {plan_path}", file=sys.stderr)
        return 2
    if not brief_path.is_file():
        print(f"deck book: design_brief not found: {brief_path}",
              file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    brief = _json.loads(brief_path.read_text(encoding="utf-8"))

    front = DeckFrontMatter(
        takeaway=brief.get("takeaway", ""),
        audience=brief.get("audience", ""),
        audience_notes=brief.get("audience_notes", ""),
        frame=brief.get("frame", ""),
        frame_rationale=brief.get("frame_rationale", ""),
        red_line=brief.get("red_line", ""),
        hook_technique=(brief.get("hook") or {}).get("technique", ""),
        hook_opener=(brief.get("hook") or {}).get("opener", ""),
        deck_title=brief.get("takeaway") or plan.get("title"),
    )

    # Optional thumbnail render. The PDF is still useful without —
    # the speaker reads notes + claims even when the .pptx isn't
    # built yet (early-iteration mode).
    thumbnails: dict[int, Path] = {}
    tmp_ctx = None
    if args.pptx:
        pptx_path = Path(args.pptx).resolve()
        if pptx_path.is_file():
            tmp_ctx = _tempfile.TemporaryDirectory()
            thumbnails = render_slides_to_png(pptx_path, Path(tmp_ctx.name))
        else:
            print(f"deck book: --pptx not found, skipping thumbnails: "
                  f"{pptx_path}", file=sys.stderr)

    plan_slides = plan.get("slides") or []
    brief_slides = {int(s.get("index", 0)): s
                    for s in (brief.get("slides") or [])}

    book_slides: list[BookSlide] = []
    for i, spec in enumerate(plan_slides):
        brief_s = brief_slides.get(i, {})
        content = spec.get("content") or {}
        claim = (brief_s.get("claim")
                 or content.get("title")
                 or content.get("action_title")
                 or "")
        notes = spec.get("notes") or brief_s.get("notes") or ""
        thumb = thumbnails.get(i + 1)  # render_slides_to_png is 1-based
        book_slides.append(BookSlide(
            index=i,
            role=brief_s.get("role", ""),
            claim=claim,
            notes=notes,
            audience_fit=brief_s.get("audience_fit", ""),
            thumbnail_path=thumb,
            section_label=spec.get("section"),
        ))

    try:
        compose_book_pdf(front, book_slides, out_path)
    finally:
        if tmp_ctx is not None:
            tmp_ctx.cleanup()

    print(f"wrote {out_path} ({len(book_slides)} slide page(s) + 1 "
          f"front-matter page; thumbnails: "
          f"{'yes' if thumbnails else 'no'})")
    return 0


def cmd_storyline(args) -> int:
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()
    try:
        titles = extract_titles_from_plan(plan_path)
    except FileNotFoundError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2

    contact_sheet = render_contact_sheet(titles, brief_summary=args.brief_summary)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    write_storyline_report(out_path, contact_sheet=contact_sheet)
    print(f"wrote {out_path} ({len([t for t in titles if t])} non-empty title(s) "
          f"across {len(titles)} slide(s))")
    return 0


def cmd_claim_evidence(args) -> int:
    """``feinschliff deck claim-evidence`` — mid-plan claim-evidence gate.

    Exit codes:
    - 0: clean (all judged slides pass)
    - 1: dirty (at least one slide has a claim-evidence defect)
    - 2: plumbing error (plan not found, parse failure, etc.)
    """
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()

    # Load plan
    try:
        plan_text = plan_path.read_text()
    except FileNotFoundError as exc:
        print(f"deck claim-evidence: {exc}", file=sys.stderr)
        return 2

    try:
        plan: dict = yaml.safe_load(plan_text) or {}
    except Exception as exc:  # noqa: BLE001
        print(f"deck claim-evidence: failed to parse plan: {exc}", file=sys.stderr)
        return 2

    # Load optional design brief
    design_brief: dict | None = None
    if args.design_brief:
        brief_path = Path(args.design_brief).resolve()
        try:
            import json as _json
            brief_text = brief_path.read_text()
            # Accept both JSON and YAML
            try:
                design_brief = _json.loads(brief_text)
            except _json.JSONDecodeError:
                design_brief = yaml.safe_load(brief_text) or {}
        except FileNotFoundError as exc:
            print(f"deck claim-evidence: {exc}", file=sys.stderr)
            return 2
        except Exception as exc:  # noqa: BLE001
            print(f"deck claim-evidence: failed to parse design brief: {exc}", file=sys.stderr)
            return 2

    slide_count = len(plan.get("slides") or [])

    try:
        results = judge_plan(
            plan,
            design_brief=design_brief,
            offline=args.offline,
            model=args.model,
        )
    except SystemExit:
        raise  # propagate ANTHROPIC_API_KEY error
    except Exception as exc:  # noqa: BLE001
        print(f"deck claim-evidence: judgment failed: {exc}", file=sys.stderr)
        return 2

    # Token-cost estimate (AC6)
    judged_count = len(results)
    if not args.offline and results:
        from feinschliff.verify.llm.prompts import claim_evidence_prompt
        # Rough estimate: average prompt length × slides judged / 4 chars/token
        sample_prompt = claim_evidence_prompt("Sample title", "Sample body text.")
        avg_tokens_per_slide = len(sample_prompt) // 4
        total_k = (avg_tokens_per_slide * judged_count) // 1000
        print(
            f"claim-evidence: {judged_count} slides judged, "
            f"~{max(total_k, 1)}k input tokens via {args.model}",
            file=sys.stderr,
        )
    else:
        print(
            f"claim-evidence: {judged_count} slides judged (--offline, 0 tokens)",
            file=sys.stderr,
        )

    overall = write_claim_evidence_report(out_path, results, slide_count=slide_count)
    print(f"wrote {out_path} (verdict: {overall}, {judged_count} judged / {slide_count} total)")

    return 0 if overall == "clean" else 1


def _load_titles_from_path(path: Path) -> list[str]:
    """Load titles from *path*.

    Accepts two shapes:
    - A plan YAML with a top-level ``slides`` list (delegated to
      ``extract_titles_from_plan``).
    - A flat YAML list of strings (titles directly).

    Raises ``FileNotFoundError`` or ``ValueError`` on bad input.
    """
    if not path.is_file():
        raise FileNotFoundError(f"file not found: {path}")
    text = path.read_text(encoding="utf-8")
    data = yaml.safe_load(text)
    if isinstance(data, list):
        # Flat list of titles
        return [str(t).strip() if t else "" for t in data]
    if isinstance(data, dict) and "slides" in data:
        # Plan-YAML shape — delegate to extract_titles_from_plan
        return extract_titles_from_plan(path)
    raise ValueError(
        f"unrecognised shape: expected a plan YAML with `slides:` or a flat list of titles ({path})"
    )


def cmd_ghost_deck(args) -> int:
    """``feinschliff deck ghost-deck`` — narrative coherence check on titles.

    Exit codes:
    - 0: pass or warn (titles form a coherent story, or minor issues only)
    - 1: fail (significant narrative gaps detected)
    - 2: plumbing error (file not found, parse failure, etc.)
    """
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()

    try:
        titles = _load_titles_from_path(plan_path)
    except FileNotFoundError as exc:
        print(f"deck ghost-deck: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"deck ghost-deck: {exc}", file=sys.stderr)
        return 2

    try:
        result = judge_ghost_deck(titles, offline=args.offline, model=args.model)
    except SystemExit:
        raise
    except Exception as exc:  # noqa: BLE001
        print(f"deck ghost-deck: judgment failed: {exc}", file=sys.stderr)
        return 2

    write_ghost_deck_report(result, out_path)
    issue_count = len(result.issues)
    print(
        f"wrote {out_path} "
        f"(verdict: {result.verdict}, {issue_count} issue(s), {len(titles)} title(s))"
    )
    return 1 if result.verdict == "fail" else 0


def cmd_title_lint(args) -> int:
    """``feinschliff deck title-lint`` — deterministic title rules, zero LLM.

    Exit codes:
    - 0: clean (no issues)
    - 1: issues found (at least one warn or fail)
    - 2: plumbing error (file not found, parse failure, etc.)
    """
    import json as _json

    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()

    try:
        titles = _load_titles_from_path(plan_path)
    except FileNotFoundError as exc:
        print(f"deck title-lint: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"deck title-lint: {exc}", file=sys.stderr)
        return 2

    issues = lint_titles(titles)

    if getattr(args, "emit_json", False):
        import dataclasses
        sys.stdout.write(_json.dumps([dataclasses.asdict(i) for i in issues], indent=2) + "\n")
    else:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        fail_count = sum(1 for i in issues if i.severity == "fail")
        warn_count = sum(1 for i in issues if i.severity == "warn")
        parts: list[str] = [
            f"# Title Lint — {len(titles)} titles",
            "",
            f"**Issues:** {len(issues)} ({fail_count} fail, {warn_count} warn)",
            "",
        ]
        if not issues:
            parts.append("_No issues._")
        else:
            for iss in issues:
                parts.append(f"- [{iss.severity.upper()}] `{iss.rule}` — {iss.message}")
        out_path.write_text("\n".join(parts) + "\n", encoding="utf-8")
        print(
            f"wrote {out_path} "
            f"({len(issues)} issue(s) across {len(titles)} title(s))"
        )

    return 1 if issues else 0


# ──────────────────────────────────────────────────────────────────────────────
# Parallel verify — focused aspect checks
# ──────────────────────────────────────────────────────────────────────────────

def cmd_verify_aspect(args) -> int:
    """`feinschliff deck verify-aspect <aspect> --plan plan.yaml -o out.json`

    Each aspect runs an independent narrow check. Designed to be spawned
    as one subagent per aspect — N aspects run in parallel, then
    `verify-collate` merges them into a single verify_report.md.

    Implementation note: aspects that need LLM judgment (narrative,
    image-style, content-cohesion) currently emit a stub finding entry
    pointing the orchestrator at the relevant PNGs/files. The deterministic
    aspects (bbox, font, brand) emit actual findings. The orchestrator
    fills LLM-judged aspects via subagent dispatch.
    """
    import json as _json
    aspect = args.aspect
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    out: dict = {
        "aspect": aspect,
        "plan": str(plan_path),
        "iteration_check_ms": None,
        "findings": [],  # list of {slide, kind, message, severity, hint}
        "summary": "",
        "needs_llm": False,
    }

    _t0 = time.perf_counter()

    if aspect == "bbox":
        # Deterministic: re-run feinschliff verify on the plan, surface
        # text-overflow / out-of-bounds / diagram-overflow defects.
        from feinschliff.pipeline import compile_slide
        import tempfile

        with tempfile.TemporaryDirectory() as _tmp:
            diagrams_out = Path(_tmp) / "diagrams"
            diagrams_out.mkdir()
            for i, spec in enumerate(slides):
                layout_path = (plan_path.parent / spec["layout"]).resolve()
                if not layout_path.is_file():
                    layout_path = _find_toolkit_file(spec["layout"]) or layout_path
                try:
                    brand_dir = find_brand(spec.get("brand")
                                           or plan.get("brand")
                                           or "feinschliff").root
                except ValueError:
                    continue
                try:
                    r = compile_slide(
                        layout_path=layout_path,
                        ctx=spec.get("content") or {},
                        brand_dir=brand_dir,
                        slide_index=i + 1,
                        diagrams_out_dir=diagrams_out,
                    )
                except Exception as e:
                    out["findings"].append({
                        "slide": i + 1, "kind": "compile-error",
                        "severity": "fatal",
                        "message": str(e)[:200], "hint": "see plan.yaml",
                    })
                    continue
                for d in r.defects:
                    out["findings"].append({
                        "slide": i + 1, "kind": d.kind.value,
                        "severity": d.severity.value,
                        "message": d.message[:240], "hint": "",
                    })

    elif aspect == "font":
        # Deterministic-ish: surface diagram-text-too-small from build,
        # plus a hint pointing the orchestrator at any slide where a text
        # primitive uses role=detail (most fragile per past runs).
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            if "size:detail" in dsl or " detail " in dsl:
                out["findings"].append({
                    "slide": i + 1, "kind": "detail-role-fragile",
                    "severity": "warn",
                    "message": "Diagram uses role=detail — close to the 10pt floor "
                               "after region/slide scaling; consider promoting to body.",
                    "hint": "search the diagram_dsl for `size:detail` or `detail `",
                })

    elif aspect == "narrative":
        # Reads titles + design_brief + storyline_report (if present).
        # The actual SCQA / claim-title judgment is LLM work; emit a
        # contact-sheet snapshot for the orchestrator subagent.
        titles = [
            (s.get("content") or {}).get("title")
            or (s.get("content") or {}).get("action_title")
            or s.get("_meta", {}).get("title", "?")
            for s in slides
        ]
        out["needs_llm"] = True
        out["contact_sheet"] = list(enumerate(titles, start=1))
        if args.design_brief:
            db = Path(args.design_brief).resolve()
            if db.is_file():
                try:
                    out["design_brief"] = _json.loads(db.read_text(encoding="utf-8"))
                except _json.JSONDecodeError:
                    pass

    elif aspect == "brand":
        # Deterministic-ish: scan diagram_dsl for fill tokens brand_bridge
        # can't resolve (semantic names + the upstream Excalidraw aliases).
        # Brand pack discipline check.
        from feinschmiede.diagrams.brand_bridge import SEMANTIC_NAMES
        from feinschmiede.diagrams.excalidraw_expand import _COLOR_ALIASES
        canonical = SEMANTIC_NAMES | frozenset(_COLOR_ALIASES)
        import re
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            for m in re.finditer(r"fill:([a-z0-9_-]+)", dsl):
                tok = m.group(1)
                if tok not in canonical:
                    out["findings"].append({
                        "slide": i + 1, "kind": "non-canonical-token",
                        "severity": "warn",
                        "message": f"Diagram uses fill:{tok}, not in the "
                                   f"{len(canonical)}-name semantic vocabulary.",
                        "hint": "use one of: " + ", ".join(sorted(canonical)[:10]) + ", …",
                    })

    elif aspect == "image":
        # Picture-slot presence + style consistency check. Surfaces slides
        # that declare an image_style but have no picture slot, or
        # vice-versa. LLM verification of actual style match is a separate
        # subagent step (needs_llm=True).
        if args.design_brief:
            try:
                db = _json.loads(Path(args.design_brief).read_text(encoding="utf-8"))
                out["image_style"] = db.get("image_style")
            except (OSError, _json.JSONDecodeError):
                pass
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            has_pic = any(k in c for k in ("hero_image", "picture", "image", "photo"))
            if has_pic and not out.get("image_style"):
                out["findings"].append({
                    "slide": i + 1, "kind": "image-style-undeclared",
                    "severity": "warn",
                    "message": "Slide uses a picture slot but design_brief.image_style "
                               "is unset.",
                    "hint": "set image_style in design_brief.json",
                })

    elif aspect == "content":
        # Title-body coherence + filler-word lint. Deterministic part:
        # scan body slots for filler words. LLM part: claim/proof match.
        FILLER = {"basically", "actually", "really", "very", "just", "simply", "in order to"}
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            for key in ("body", "supporting_body", "so_what"):
                val = c.get(key) or ""
                if not isinstance(val, str):
                    continue
                low = val.lower()
                hits = [w for w in FILLER if w in low]
                if hits:
                    out["findings"].append({
                        "slide": i + 1, "kind": "filler-word",
                        "severity": "warn",
                        "message": f"{key} contains filler: {', '.join(hits)}",
                        "hint": "Cut the filler — strong sentences don't need it.",
                    })

    elif aspect == "notes-coherence":
        # Pair the deck's red_line against each slide's (claim, notes).
        # The orchestrator LLM judges whether the spoken delivery tracks
        # the arc: drift / contradiction / off-arc tangents → dirty.
        from feinschliff.verify.deck.notes_coherence import (
            SlideForCoherence,
            render_contact_sheet as _render_notes_sheet,
        )
        out["needs_llm"] = True
        red_line = ""
        if args.design_brief:
            db_path = Path(args.design_brief).resolve()
            if db_path.is_file():
                try:
                    db = _json.loads(db_path.read_text(encoding="utf-8"))
                    red_line = db.get("red_line", "") or ""
                    out["red_line"] = red_line
                    out["design_brief"] = db
                except _json.JSONDecodeError:
                    pass
        coherence_slides: list[SlideForCoherence] = []
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            coherence_slides.append(SlideForCoherence(
                index=i,
                role=spec.get("_meta", {}).get("role", ""),
                claim=c.get("title") or c.get("action_title") or "",
                notes=spec.get("notes"),
            ))
        # Cheap deterministic pre-flag: hook slide missing notes.
        # The LLM judge handles drift / off-arc semantics.
        if coherence_slides and not (coherence_slides[0].notes or "").strip():
            out["findings"].append({
                "slide": 1, "kind": "hook-notes-missing",
                "severity": "warn",
                "message": "Hook slide has no speaker notes; expected the "
                           "deck-level storyline articulating the red_line.",
                "hint": "Author the full red_line arc into slide 1's notes.",
            })
        out["contact_sheet"] = _render_notes_sheet(red_line, coherence_slides)

    out["iteration_check_ms"] = int(
        (time.perf_counter() - _t0) * 1000
    )
    out["summary"] = (
        f"{len(out['findings'])} finding(s) "
        f"({'needs LLM' if out['needs_llm'] else 'deterministic only'})"
    )

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(_json.dumps(out, indent=2, ensure_ascii=False),
                        encoding="utf-8")
    log_event(plan_path.parent, f"verify-aspect:{aspect}", "tick",
              elapsed_ms=out["iteration_check_ms"],
              findings=len(out["findings"]))
    print(f"wrote {out_path} — {out['summary']}")
    return 0


def cmd_verify_collate(args) -> int:
    """`feinschliff deck verify-collate --aspect a.json --aspect b.json -o report.md`
    Merge per-aspect verify outputs into a unified verify_report.md.
    """
    import json as _json
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    # Aggregate findings grouped by slide index. Each finding carries its
    # source aspect so the report makes the parallel checks visible.
    by_slide: dict[int, list[dict]] = {}
    aspects_seen: list[str] = []
    needs_llm: list[str] = []
    total_findings = 0

    for path_str in args.aspect:
        p = Path(path_str).resolve()
        if not p.is_file():
            print(f"deck verify-collate: aspect file not found: {p}",
                  file=sys.stderr)
            continue
        try:
            data = _json.loads(p.read_text(encoding="utf-8"))
        except _json.JSONDecodeError as e:
            print(f"deck verify-collate: {p}: {e}", file=sys.stderr)
            continue
        aspect = data.get("aspect", p.stem)
        aspects_seen.append(aspect)
        if data.get("needs_llm"):
            needs_llm.append(aspect)
        for f in data.get("findings", []) or []:
            f = dict(f)
            f["aspect"] = aspect
            by_slide.setdefault(int(f.get("slide", 0)), []).append(f)
            total_findings += 1

    # Decide verdict: clean iff no findings + no LLM aspects waiting.
    fatal = any(
        f.get("severity") == "fatal"
        for fs in by_slide.values() for f in fs
    )
    dirty = total_findings > 0 or fatal
    verdict = "dirty" if dirty else ("pending-llm" if needs_llm else "clean")

    lines: list[str] = []
    lines.append(f"# Verify Report — {plan_path.name}")
    lines.append("")
    lines.append(f"- **Iteration:** {args.iteration} of {args.budget}")
    lines.append(f"- **Verdict:** {verdict}"
                 + (f" — {total_findings} finding(s) across "
                    f"{len(by_slide)} slide(s)" if dirty else ""))
    lines.append(f"- **Aspects checked:** {', '.join(aspects_seen) or '(none)'}")
    if needs_llm:
        lines.append(f"- **Pending LLM verdict from:** {', '.join(needs_llm)}")
    if args.png_dir:
        lines.append(f"- **Rendered PNGs:** `{args.png_dir}`")
    lines.append("")
    lines.append("---")
    lines.append("")

    for i, spec in enumerate(slides, start=1):
        title = (spec.get("content") or {}).get("title") \
                or (spec.get("content") or {}).get("action_title") \
                or spec.get("_meta", {}).get("title", "(untitled)")
        layout_name = Path(spec.get("layout", "?")).name.replace(".slide.dsl", "")
        findings = by_slide.get(i, [])
        if not findings:
            lines.append(f"## Slide {i} — {title!r} ({layout_name}) ✅")
            lines.append("")
            lines.append("_No defects._")
        else:
            lines.append(f"## Slide {i} — {title!r} ({layout_name})")
            lines.append("")
            for f in findings:
                lines.append(
                    f"- **[{f['aspect']}/{f.get('kind', '?')}]** "
                    f"{f.get('message', '')}"
                )
                if f.get("hint"):
                    lines.append(f"  → {f['hint']}")
        lines.append("")

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    log_event(plan_path.parent, "verify-collate", "tick",
              aspects=len(aspects_seen), findings=total_findings,
              verdict=verdict, iteration=args.iteration)
    print(f"wrote {out_path} — verdict: {verdict}, "
          f"{total_findings} finding(s) across {len(by_slide)} slide(s)")
    return 0


def cmd_verify_static(args) -> int:
    """`feinschliff deck verify-static <plan.yaml>` — pre-render static check.

    Inspects a plan.yaml for geometry defects that can be detected from the
    DSL + populated content without rendering (slot-overflow,
    empty-placeholder). Cheaper than a full build: catches class of defect
    in ~10-50 ms vs ~3 s/slide for a render-based check.

    Exit codes:
      0 — clean (no defects)
      1 — one or more defects found
      2 — plumbing error (plan not found, brand resolution failure, etc.)
    """
    _require_or_delegate_builder("deck verify-static")
    import json as _json
    from feinschliff_builder.verify.static import validate as _validate_static

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck verify-static: plan not found: {plan_path}", file=sys.stderr)
        return 2

    try:
        plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    except (OSError, yaml.YAMLError) as exc:
        print(f"deck verify-static: could not load plan: {exc}", file=sys.stderr)
        return 2

    brand_name = getattr(args, "brand", None) or plan.get("brand") or "feinschliff"
    try:
        brand_obj = find_brand(brand_name)
    except ValueError as exc:
        print(f"deck verify-static: {exc}", file=sys.stderr)
        return 2

    try:
        bag = _validate_static(plan, brand=brand_obj, plan_dir=plan_path.parent)
    except Exception as exc:  # noqa: BLE001
        print(f"deck verify-static: unexpected error: {exc}", file=sys.stderr)
        return 2

    if getattr(args, "json", False):
        # Produce a backward-compatible schema so `deck apply-fixes --defects`
        # can consume this output: {slide_index, kind, severity, message, meta}.
        out = []
        for d in bag:
            extra = d.extra or {}
            entry = {
                "slide_index": extra.get("slide_index", 0),
                "kind": d.kind.value,
                "severity": d.severity.value,
                "message": d.message,
                "meta": {k: v for k, v in extra.items() if k != "slide_index"},
            }
            if d.location is not None:
                entry["location"] = d.location
            out.append(entry)
        print(_json.dumps(out, indent=2, ensure_ascii=False))
    else:
        if bag:
            for d in bag:
                _loc = d.location or "slide ?"
                print(
                    f"{_loc}: [{d.severity.value.upper()}] {d.kind.value} — {d.message}"
                )
        else:
            print("verify-static: clean — no defects found")

    return 1 if bag else 0


def _parse_severity(value: str):
    """Parse a severity string that may use either engine or legacy vocabulary.

    ``deck verify-static --json`` emits ENGINE severity values because it goes
    through ``feinschliff_builder.verify.static.validate()``, which maps the
    legacy Severity enum to the ``feinschmiede.diagnostics`` engine enum before
    serialising:

        legacy FATAL  → engine "error"
        legacy WARN   → engine "warning"
        legacy INFO   → engine "info"

    Legacy values (``"fatal"``, ``"warn"``, ``"info"``) are also accepted so
    that hand-crafted defect files and older tooling continue to work.

    The inverse mapping applied here exactly undoes validate()'s forward mapping:

        engine "error"   → Severity.FATAL
        engine "warning" → Severity.WARN
        engine "info"    → Severity.INFO
        legacy "fatal"   → Severity.FATAL  (pass-through)
        legacy "warn"    → Severity.WARN   (pass-through)

    Raises ``ValueError`` for unrecognised values (caller logs + skips).
    """
    from feinschliff.defects import Severity
    _ENGINE_TO_LEGACY = {
        "error": Severity.FATAL,
        "warning": Severity.WARN,
        "info": Severity.INFO,
    }
    if value in _ENGINE_TO_LEGACY:
        return _ENGINE_TO_LEGACY[value]
    # Fall through to legacy enum constructor (handles "fatal", "warn", "info").
    return Severity(value)


def cmd_apply_fixes(args) -> int:
    """`feinschliff deck apply-fixes <plan.yaml> --defects <defects.json> [-o out.yaml]`

    Read the defects JSON (flat list OR collated shape), translate to
    deterministic FixPatch objects, apply them to the plan, and write the
    result.  Prints a markdown diff summary to stdout.

    Defects JSON shapes supported:
      - Flat list:  [{slide_index, kind, severity, message, meta}, ...]
        (output of `deck verify-static --json`)
      - Collated:   {"defects": {"1": [...], "2": [...]}, ...}
        (output shape used by cli/verify.py)

    Severity values accepted: both ENGINE vocabulary ("error", "warning",
    "info" — emitted by ``deck verify-static --json``) and LEGACY vocabulary
    ("fatal", "warn", "info" — used by hand-crafted files and older tooling).
    See ``_parse_severity()`` for the exact mapping.

    Exit codes:
      0 — at least one patch was applied
      1 — no patches applied (defects present but none mechanically fixable,
          OR no defects at all)
      2 — plumbing error
    """
    _require_or_delegate_builder("deck apply-fixes")
    import json as _json
    from feinschliff_builder.verify.autofix import plan_fixes, apply_fixes, diff_summary
    from feinschliff.defects import Defect, DefectKind

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck apply-fixes: plan not found: {plan_path}", file=sys.stderr)
        return 2

    try:
        plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    except (OSError, yaml.YAMLError) as exc:
        print(f"deck apply-fixes: could not load plan: {exc}", file=sys.stderr)
        return 2

    defects_path = Path(args.defects).resolve()
    if not defects_path.is_file():
        print(f"deck apply-fixes: defects file not found: {defects_path}", file=sys.stderr)
        return 2

    try:
        raw = _json.loads(defects_path.read_text(encoding="utf-8"))
    except (_json.JSONDecodeError, OSError) as exc:
        print(f"deck apply-fixes: could not load defects: {exc}", file=sys.stderr)
        return 2

    # Normalise both JSON shapes to a flat list of Defect objects.
    defect_dicts: list[dict] = []
    if isinstance(raw, list):
        # Flat list shape: [{slide_index, kind, severity, message, meta}, ...]
        defect_dicts = raw
    elif isinstance(raw, dict):
        # Collated shape: {"defects": {slide_idx: [...]}, ...}
        nested = raw.get("defects") or {}
        for _slide_defects in nested.values():
            if isinstance(_slide_defects, list):
                defect_dicts.extend(_slide_defects)
    else:
        print("deck apply-fixes: unrecognised defects JSON shape", file=sys.stderr)
        return 2

    defects: list[Defect] = []
    for dd in defect_dicts:
        try:
            defects.append(Defect(
                slide_index=int(dd["slide_index"]),
                kind=DefectKind(dd["kind"]),
                severity=_parse_severity(dd["severity"]),
                message=str(dd.get("message", "")),
                meta=dict(dd.get("meta") or {}),
            ))
        except (KeyError, ValueError) as exc:
            print(
                f"deck apply-fixes: skipping malformed defect entry {dd!r}: {exc}",
                file=sys.stderr,
            )

    if not defects:
        print("deck apply-fixes: no defects to process — nothing to do")
        return 1

    brand_name = getattr(args, "brand", None) or plan.get("brand") or "feinschliff"
    try:
        brand_obj = find_brand(brand_name)
    except ValueError as exc:
        print(f"deck apply-fixes: {exc}", file=sys.stderr)
        return 2

    patches = plan_fixes(defects, plan, brand_obj.root)
    if not patches:
        print("deck apply-fixes: no mechanical fixes available for these defects")
        print("Auto-fix passes: 0")
        return 1

    fixed_plan = apply_fixes(plan, patches)
    summary = diff_summary(plan, fixed_plan)

    out_path = Path(args.output).resolve() if args.output else plan_path
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        yaml.safe_dump(fixed_plan, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )

    # Print diff summary to stdout.
    if summary:
        print(summary)
    print(f"Auto-fix passes: {len(patches)}")
    print(f"wrote {out_path} ({len(patches)} patch(es) applied)")
    return 0
